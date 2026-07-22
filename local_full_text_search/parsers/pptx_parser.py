from __future__ import annotations

from pathlib import Path
from typing import Iterable

from local_full_text_search.core.errors import ParserDependencyError
from local_full_text_search.core.task_manager import CancelToken
from local_full_text_search.models.content_block import ContentBlock
from local_full_text_search.parsers.base_parser import BaseParser


class PptxParser(BaseParser):
    name = "pptx"

    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == ".pptx"

    def parse(self, file_path: Path, cancel_token: CancelToken) -> Iterable[ContentBlock]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ParserDependencyError("未安装 python-pptx，无法解析 PPTX") from exc
        presentation = Presentation(str(file_path))
        block_index = 0
        for slide_number, slide in enumerate(presentation.slides, start=1):
            cancel_token.wait_if_paused()
            cancel_token.throw_if_cancelled()
            texts = _extract_shape_texts(slide.shapes)
            if texts:
                yield self.make_block(
                    file_path,
                    block_index,
                    "pptx_slide",
                    f"第 {slide_number} 张幻灯片",
                    "\n".join(texts),
                    slide_number=slide_number,
                )
                block_index += 1
            notes = []
            try:
                notes = _extract_shape_texts(slide.notes_slide.shapes)
            except Exception:
                notes = []
            if notes:
                yield self.make_block(
                    file_path,
                    block_index,
                    "pptx_notes",
                    f"第 {slide_number} 张幻灯片；备注",
                    "\n".join(notes),
                    slide_number=slide_number,
                )
                block_index += 1


def _extract_shape_texts(shapes: object) -> list[str]:
    texts: list[str] = []
    for shape in shapes:
        if hasattr(shape, "shapes"):
            texts.extend(_extract_shape_texts(shape.shapes))
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                texts.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    texts.append(" | ".join(cells))
    return texts
