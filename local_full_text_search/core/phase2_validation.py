from __future__ import annotations

import hashlib
import json
import os
import tempfile
import threading
import time
from collections.abc import Callable
from concurrent.futures import Future
from pathlib import Path

from local_full_text_search.config.defaults import AppSettings
from local_full_text_search.core.database import DatabaseManager
from local_full_text_search.core.index_manager import (
    IndexManager,
    ParseJob,
    ParseLane,
    ParseOutcome,
)
from local_full_text_search.core.task_manager import CancelToken
from local_full_text_search.core.ocr_backend_benchmark import (
    BackendTrial,
    candidate_backends,
    choose_default_backend,
)
from local_full_text_search.ocr.ocr_engine import OcrEngine


def validate_pdf_page_pipeline(base: Path) -> dict[str, object]:
    import fitz

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    path = root / "five-hundred-pages.pdf"
    document = fitz.open()
    try:
        for page_number in range(1, 501):
            page = document.new_page()
            page.insert_text(
                (72, 72),
                f"FROZEN_PDF_PAGE_{page_number:03d}",
            )
        document.save(path)
    finally:
        document.close()
    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    started = time.perf_counter()
    summary = IndexManager(
        database,
        AppSettings(
            enable_ocr=False,
            ocr_scanned_pdf=False,
            pdf_parser_workers=2,
        ),
    ).index_root(root_id)
    elapsed_seconds = time.perf_counter() - started
    with database.connect() as connection:
        page_tasks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'pdf_native_page'
                  AND status = 'complete'
                  AND confirmed_at IS NOT NULL
                """
            ).fetchone()[0]
        )
        merge_tasks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'document_merge'
                  AND status = 'complete'
                  AND confirmed_at IS NOT NULL
                """
            ).fetchone()[0]
        )
        pages = [
            int(row["page_number"])
            for row in connection.execute(
                """
                SELECT page_number FROM content_blocks
                WHERE block_type = 'pdf_page'
                ORDER BY page_number, block_index
                """
            )
        ]
        digest = hashlib.sha256(
            "\n".join(
                str(row["raw_text"] or "")
                for row in connection.execute(
                    """
                    SELECT raw_text FROM content_blocks
                    ORDER BY page_number, block_index
                    """
                )
            ).encode("utf-8")
        ).hexdigest()
        normal_search_hits = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM content_blocks
                WHERE raw_text LIKE '%FROZEN_PDF_PAGE_250%'
                """
            ).fetchone()[0]
        )
    cancel_gate = _validate_pdf_cancel_gate(
        base / "cancel-gate",
        root,
    )
    mode_consistency = _validate_pdf_mode_consistency(
        base / "performance-gate",
        root,
        expected_digest=digest,
        expected_pages=pages,
        expected_search_hits=normal_search_hits,
    )
    passed = bool(
        summary.failed == 0
        and page_tasks == 500
        and merge_tasks == 1
        and pages == list(range(1, 501))
        and bool(cancel_gate["passed"])
        and bool(mode_consistency["passed"])
    )
    return {
        "passed": passed,
        "elapsed_seconds": round(elapsed_seconds, 3),
        "page_tasks": page_tasks,
        "merge_tasks": merge_tasks,
        "ordered_pages": len(pages),
        "content_digest": digest,
        "cancel_gate": cancel_gate,
        "mode_consistency": mode_consistency,
        "summary": summary.to_dict(),
    }


def _validate_pdf_cancel_gate(
    base: Path,
    root: Path,
) -> dict[str, object]:
    base.mkdir(parents=True, exist_ok=True)
    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    manager = IndexManager(
        database,
        AppSettings(
            enable_ocr=False,
            ocr_scanned_pdf=False,
            pdf_parser_workers=1,
        ),
        run_context={"execution_mode": "normal"},
    )
    token = CancelToken()
    results: list[object] = []
    errors: list[str] = []

    def run() -> None:
        try:
            results.append(manager.index_root(root_id, token))
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
    )
    os.environ["LFTS_VALIDATION_SAFE_POINT_DELAY_MS"] = "20"
    worker = threading.Thread(
        target=run,
        name="lfts-pdf-cancel-validation",
    )
    worker.start()
    confirmed_before_cancel = 0
    deadline = time.monotonic() + 45.0
    try:
        while worker.is_alive() and time.monotonic() < deadline:
            with database.connect() as connection:
                confirmed_before_cancel = int(
                    connection.execute(
                        """
                        SELECT COUNT(*) FROM parse_tasks
                        WHERE task_type = 'pdf_native_page'
                          AND status = 'complete'
                          AND confirmed_at IS NOT NULL
                        """
                    ).fetchone()[0]
                )
            if confirmed_before_cancel > 0:
                token.cancel(force=True)
                break
            time.sleep(0.02)
        worker.join(timeout=30.0)
    finally:
        if worker.is_alive():
            token.cancel(force=True)
            manager.force_terminate_processes()
            worker.join(timeout=5.0)
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
            ] = previous_delay
    with database.connect() as connection:
        confirmed_pages = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'pdf_native_page'
                  AND status = 'complete'
                  AND confirmed_at IS NOT NULL
                """
            ).fetchone()[0]
        )
        pending_pages = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'pdf_native_page'
                  AND status != 'complete'
                """
            ).fetchone()[0]
        )
        published_blocks = int(
            connection.execute(
                "SELECT COUNT(*) FROM content_blocks"
            ).fetchone()[0]
        )
        merge_completed = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'document_merge'
                  AND status = 'complete'
                """
            ).fetchone()[0]
        )
    summary = results[0] if results else None
    cancelled = bool(
        summary is not None and getattr(summary, "cancelled", False)
    )
    search_ready = bool(database.index_readiness()["ready"])
    passed = bool(
        not errors
        and not worker.is_alive()
        and cancelled
        and confirmed_pages > 0
        and pending_pages > 0
        and published_blocks == 0
        and merge_completed == 0
        and not search_ready
    )
    return {
        "passed": passed,
        "cancelled": cancelled,
        "confirmed_before_cancel": confirmed_before_cancel,
        "confirmed_pages": confirmed_pages,
        "pending_pages": pending_pages,
        "published_blocks": published_blocks,
        "merge_completed": merge_completed,
        "search_ready": search_ready,
        "worker_exited": not worker.is_alive(),
        "errors": errors,
    }


def _validate_pdf_mode_consistency(
    base: Path,
    root: Path,
    *,
    expected_digest: str,
    expected_pages: list[int],
    expected_search_hits: int,
) -> dict[str, object]:
    base.mkdir(parents=True, exist_ok=True)
    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    summary = IndexManager(
        database,
        AppSettings(
            enable_ocr=False,
            ocr_scanned_pdf=False,
            pdf_parser_workers=4,
            process_parser_workers=4,
            process_pending_tasks=8,
        ),
        run_context={"execution_mode": "performance"},
    ).index_root(root_id)
    with database.connect() as connection:
        pages = [
            int(row["page_number"])
            for row in connection.execute(
                """
                SELECT page_number FROM content_blocks
                WHERE block_type = 'pdf_page'
                ORDER BY page_number, block_index
                """
            )
        ]
        digest = hashlib.sha256(
            "\n".join(
                str(row["raw_text"] or "")
                for row in connection.execute(
                    """
                    SELECT raw_text FROM content_blocks
                    ORDER BY page_number, block_index
                    """
                )
            ).encode("utf-8")
        ).hexdigest()
        search_hits = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM content_blocks
                WHERE raw_text LIKE '%FROZEN_PDF_PAGE_250%'
                """
            ).fetchone()[0]
        )
    digest_matches = digest == expected_digest
    page_order_matches = pages == expected_pages
    search_hits_match = search_hits == expected_search_hits
    passed = bool(
        summary.failed == 0
        and digest_matches
        and page_order_matches
        and search_hits_match
    )
    return {
        "passed": passed,
        "digest_matches": digest_matches,
        "page_order_matches": page_order_matches,
        "search_hits_match": search_hits_match,
        "content_digest": digest,
        "page_count": len(pages),
        "search_hits": search_hits,
        "summary": summary.to_dict(),
    }


def validate_ocr_adaptive_v2(base: Path) -> dict[str, object]:
    import numpy as np
    from PIL import Image

    base.mkdir(parents=True, exist_ok=True)
    source = base / "adaptive.png"
    checkpoint = base / "adaptive-checkpoint.json"
    pixels = np.random.default_rng(20260730).integers(
        0,
        256,
        size=(1_200, 2_200, 3),
        dtype=np.uint8,
    )
    Image.fromarray(pixels).save(source)

    class Detector:
        def predict(self, image: object) -> list[dict[str, object]]:
            height, width = image.shape[:2]
            return [
                {
                    "dt_polys": np.asarray(
                        [
                            [
                                [20, 20],
                                [min(width - 1, 360), 20],
                                [min(width - 1, 360), 100],
                                [20, 100],
                            ]
                        ],
                        dtype=np.float32,
                    )
                }
            ]

    first_calls: list[int] = []

    class CrashingRecognizer:
        def predict(
            self,
            crops: list[object],
        ) -> list[dict[str, object]]:
            first_calls.append(len(crops))
            if len(first_calls) == 3:
                raise RuntimeError("injected recognition batch crash")
            return [
                {"rec_text": "冻结版自适应识别", "rec_score": 0.42}
                for _ in crops
            ]

    first = OcrEngine(
        det_limit_side_len=960,
        microbatch_max_requests=2,
        microbatch_wait_ms=0,
    )
    first._detector = Detector()
    first._recognizer = CrashingRecognizer()
    crash_observed = False
    try:
        first.recognize_adaptive(
            source,
            checkpoint_path=checkpoint,
        )
    except RuntimeError as exc:
        crash_observed = "injected recognition batch crash" in str(
            exc
        )

    resumed_calls: list[int] = []

    class ResumedRecognizer:
        def predict(
            self,
            crops: list[object],
        ) -> list[dict[str, object]]:
            resumed_calls.append(len(crops))
            return [
                {"rec_text": "冻结版自适应识别", "rec_score": 0.42}
                for _ in crops
            ]

    resumed = OcrEngine(
        det_limit_side_len=960,
        microbatch_max_requests=2,
        microbatch_wait_ms=0,
    )
    resumed._detector = Detector()
    resumed._recognizer = ResumedRecognizer()
    result = resumed.recognize_adaptive(
        source,
        checkpoint_path=checkpoint,
    )
    passed = bool(
        crash_observed
        and not checkpoint.exists()
        and int(
            result.extra.get(
                "checkpoint_recognition_batches_reused",
                0,
            )
        )
        >= 1
        and int(result.extra.get("adaptive_regions_remaining", -1))
        == 0
        and float(result.extra.get("coverage_ratio", 0.0)) == 1.0
    )
    return {
        "passed": passed,
        "crash_observed": crash_observed,
        "first_inference_batches": first_calls,
        "resumed_inference_batches": resumed_calls,
        "checkpoint_reused": result.extra.get(
            "checkpoint_recognition_batches_reused",
            0,
        ),
        "adaptive_regions_created": result.extra.get(
            "adaptive_regions_created",
            0,
        ),
        "adaptive_regions_split": result.extra.get(
            "adaptive_regions_split",
            0,
        ),
        "coverage_ratio": result.extra.get("coverage_ratio", 0.0),
    }


def validate_ocr_backend_gate(base: Path) -> dict[str, object]:
    def trial(
        backend: str,
        elapsed: int,
    ) -> BackendTrial:
        return BackendTrial(
            backend=backend,
            elapsed_ms=elapsed,
            model_load_ms=10,
            first_batch_ms=5,
            peak_rss_bytes=128 * 1024**2,
            accuracy_digest="same",
            text_digest="same",
            box_digest="same",
            confidence_digest="same",
            offline=True,
            error="",
        )

    baseline = [
        trial("paddle_cpu_mkldnn_off", elapsed)
        for elapsed in (1000, 980, 1020)
    ]
    candidate = [
        trial("paddle_cpu_mkldnn_on", elapsed)
        for elapsed in (750, 760, 740)
    ]
    decision = choose_default_backend(
        baseline,
        [candidate],
        rss_budget_bytes=4 * 1024**3,
    )
    candidates = candidate_backends()
    mechanism_passed = bool(
        decision.selected_backend == "paddle_cpu_mkldnn_on"
        and decision.speedup_ratio >= 0.20
        and len(candidates) == 4
        and all(item.same_model_required for item in candidates)
    )
    real_scheduler = _validate_real_unified_ocr_lane(
        base / "unified-lane"
    )
    worker_crash_recovery = _validate_real_ocr_worker_rebuild(
        base / "worker-rebuild"
    )
    passed = bool(
        mechanism_passed
        and real_scheduler["passed"]
        and worker_crash_recovery["passed"]
    )
    return {
        "passed": passed,
        "mechanism_only": False,
        "real_three_run_ab_required": True,
        "real_scheduler": real_scheduler,
        "worker_crash_recovery": worker_crash_recovery,
        "candidates": [
            {
                "backend": item.backend,
                "available": item.available,
                "same_model_required": item.same_model_required,
            }
            for item in candidates
        ],
        "decision": {
            "passed": mechanism_passed,
            "selected_backend": decision.selected_backend,
            "changed": decision.changed,
            "speedup_ratio": decision.speedup_ratio,
            "rejections": decision.rejections,
        },
    }


def _draw_ocr_validation_image(
    path: Path,
    label: str,
    *,
    width: int = 1200,
    height: int = 800,
) -> None:
    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    font_path = (
        Path(os.environ.get("WINDIR", "C:/Windows"))
        / "Fonts"
        / "arial.ttf"
    )
    font = (
        ImageFont.truetype(str(font_path), 36)
        if font_path.is_file()
        else ImageFont.load_default()
    )
    for line_number in range(1, 12):
        draw.text(
            (45, 30 + (line_number - 1) * 62),
            (
                f"{label} LINE {line_number:02d} "
                f"VALUE {line_number * 173}"
            ),
            fill="black",
            font=font,
        )
    path.parent.mkdir(parents=True, exist_ok=True)
    image.save(path)


def _unified_ocr_settings() -> AppSettings:
    return AppSettings(
        enable_ocr=True,
        ocr_images=True,
        ocr_scanned_pdf=True,
        min_ocr_image_pixels=0,
        enable_parse_cache=False,
        parser_workers=1,
        process_parser_workers=1,
        process_pending_tasks=1,
        pdf_parser_workers=1,
        ocr_workers=1,
        ocr_pending_tasks=1,
        ocr_microbatch_parent_jobs=4,
        slow_file_workers=1,
        slow_pending_tasks=1,
        index_write_batch_size=8,
    )


def _latest_index_metrics(
    database: DatabaseManager,
) -> dict[str, object]:
    runs = database.recent_index_runs_since(
        "1970-01-01T00:00:00+00:00"
    )
    if not runs:
        return {}
    summary = runs[-1].get("summary")
    if not isinstance(summary, dict):
        return {}
    metrics = summary.get("metrics")
    return dict(metrics) if isinstance(metrics, dict) else {}


def _wait_for_processes_to_exit(
    worker_pids: set[int],
    *,
    timeout: float = 10.0,
) -> list[int]:
    import psutil

    pending = {pid for pid in worker_pids if pid > 0}
    deadline = time.monotonic() + max(0.0, timeout)
    while pending and time.monotonic() < deadline:
        pending = {pid for pid in pending if psutil.pid_exists(pid)}
        if pending:
            time.sleep(0.05)
    return sorted(pid for pid in pending if psutil.pid_exists(pid))


def _validate_real_unified_ocr_lane(
    base: Path,
) -> dict[str, object]:
    import fitz
    import zipfile

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    token = hashlib.sha256(str(base).encode("utf-8")).hexdigest()[:8]

    pdf_source = root / "00-large-scanned.pdf"
    document = fitz.open()
    try:
        for page_number in range(1, 5):
            page_image = base / f"pdf-page-{page_number}.png"
            _draw_ocr_validation_image(
                page_image,
                f"UNIFIED PDF {token} PAGE {page_number:02d}",
            )
            page = document.new_page(width=720, height=480)
            page.insert_image(page.rect, filename=str(page_image))
        document.save(pdf_source)
    finally:
        document.close()

    standalone = root / "01-small-image.png"
    _draw_ocr_validation_image(
        standalone,
        f"UNIFIED IMAGE {token}",
        width=900,
        height=620,
    )
    zip_member = base / "zip-member.png"
    _draw_ocr_validation_image(
        zip_member,
        f"UNIFIED ZIP IMAGE {token}",
        width=920,
        height=640,
    )
    archive_path = root / "02-images.zip"
    with zipfile.ZipFile(
        archive_path,
        "w",
        compression=zipfile.ZIP_DEFLATED,
    ) as archive:
        archive.write(zip_member, arcname="nested/zip-image.png")

    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    summary = IndexManager(
        database,
        _unified_ocr_settings(),
        run_context={
            "execution_mode": "normal",
            "validation_kind": "unified_ocr_real_lane",
        },
    ).index_root(root_id)
    with database.connect() as connection:
        rows = [
            dict(row)
            for row in connection.execute(
                """
                SELECT request.source_kind, request.source_unit,
                       request.confirmed_at, request.id,
                       task.worker_pid
                FROM ocr_requests request
                JOIN parse_tasks task
                  ON task.id = request.parent_task_id
                WHERE request.status = 'confirmed'
                ORDER BY request.confirmed_at, request.id
                """
            )
        ]
        extras = [
            json.loads(str(row[0] or "{}"))
            for row in connection.execute(
                """
                SELECT extra_json FROM content_blocks
                WHERE source_type = 'ocr'
                ORDER BY id
                """
            )
        ]
    source_kinds = sorted(
        {str(row.get("source_kind") or "") for row in rows}
    )
    worker_pids = {
        int(row.get("worker_pid") or 0)
        for row in rows
        if int(row.get("worker_pid") or 0) > 0
    }
    model_load_counts = [
        int(extra.get("ocr_model_load_count") or 0)
        for extra in extras
    ]
    model_states = {
        str(extra.get("ocr_model_state") or "")
        for extra in extras
        if extra.get("ocr_model_state")
    }
    pdf_positions = [
        index
        for index, row in enumerate(rows)
        if row.get("source_kind") == "pdf_page"
    ]
    small_positions = [
        index
        for index, row in enumerate(rows)
        if row.get("source_kind") in {"image", "zip_image"}
    ]
    small_before_pdf_drained = bool(
        pdf_positions
        and small_positions
        and min(small_positions) < max(pdf_positions)
    )
    metrics = _latest_index_metrics(database)
    ocr_metrics = dict(metrics.get("ocr_metrics") or {})
    alive_after = _wait_for_processes_to_exit(worker_pids)
    model_load_count = max(model_load_counts, default=0)
    passed = bool(
        summary.failed == 0
        and source_kinds == ["image", "pdf_page", "zip_image"]
        and len(rows) == 6
        and len(worker_pids) == 1
        and model_load_count == 2
        and model_states == {"ready"}
        and int(ocr_metrics.get("ocr_worker_count_peak") or 0) == 1
        and int(ocr_metrics.get("ocr_model_load_count") or 0) == 2
        and small_before_pdf_drained
        and not alive_after
    )
    return {
        "passed": passed,
        "confirmed_source_kinds": source_kinds,
        "confirmed_request_count": len(rows),
        "confirmed_order": [
            {
                "source_kind": str(row.get("source_kind") or ""),
                "source_unit": str(row.get("source_unit") or ""),
            }
            for row in rows
        ],
        "shared_worker_count": len(worker_pids),
        "worker_pids": sorted(worker_pids),
        "model_state": (
            "ready" if model_states == {"ready"} else "unknown"
        ),
        "model_load_count_per_worker": model_load_count,
        "reported_worker_peak": int(
            ocr_metrics.get("ocr_worker_count_peak") or 0
        ),
        "reported_model_load_count": int(
            ocr_metrics.get("ocr_model_load_count") or 0
        ),
        "small_source_completed_before_pdf_drained": (
            small_before_pdf_drained
        ),
        "worker_pids_after_completion": alive_after,
        "summary": summary.to_dict(),
    }


def _validate_real_ocr_worker_rebuild(
    base: Path,
) -> dict[str, object]:
    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    token = hashlib.sha256(str(base).encode("utf-8")).hexdigest()[:8]
    source = root / "crash-and-recover.png"
    _draw_ocr_validation_image(
        source,
        f"OCR WORKER REBUILD {token}",
    )
    marker_path = base / "ocr-worker-crash.json"
    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    progress_events: list[dict[str, object]] = []
    previous_marker = os.environ.get(
        "LFTS_VALIDATION_OCR_CRASH_MARKER"
    )
    os.environ["LFTS_VALIDATION_OCR_CRASH_MARKER"] = str(
        marker_path
    )
    try:
        summary = IndexManager(
            database,
            _unified_ocr_settings(),
            run_context={
                "execution_mode": "normal",
                "validation_kind": "ocr_worker_rebuild",
            },
        ).index_root(
            root_id,
            progress_callback=lambda payload: progress_events.append(
                dict(payload)
            ),
        )
    finally:
        if previous_marker is None:
            os.environ.pop(
                "LFTS_VALIDATION_OCR_CRASH_MARKER",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_OCR_CRASH_MARKER"
            ] = previous_marker
    marker = (
        json.loads(marker_path.read_text(encoding="utf-8"))
        if marker_path.is_file()
        else {}
    )
    with database.connect() as connection:
        request = connection.execute(
            """
            SELECT request.status, task.worker_pid
            FROM ocr_requests request
            JOIN parse_tasks task
              ON task.id = request.parent_task_id
            ORDER BY request.id DESC LIMIT 1
            """
        ).fetchone()
        source_success = bool(
            connection.execute(
                """
                SELECT COUNT(*) FROM files
                WHERE path = ? AND parse_status = 'success'
                """,
                (str(source),),
            ).fetchone()[0]
        )
        extras = [
            json.loads(str(row[0] or "{}"))
            for row in connection.execute(
                """
                SELECT extra_json FROM content_blocks
                WHERE source_type = 'ocr'
                """
            )
        ]
        attempts = [
            dict(row)
            for row in connection.execute(
                """
                SELECT attempt_no, status, worker_pid,
                       error_code, error_message
                FROM parse_task_attempts
                ORDER BY attempt_no
                """
            )
        ]
    crashed_pid = int(marker.get("worker_pid") or 0)
    replacement_pid = int(request["worker_pid"] or 0) if request else 0
    replacement_load_count = max(
        [
            int(extra.get("ocr_model_load_count") or 0)
            for extra in extras
        ],
        default=0,
    )
    metrics = _latest_index_metrics(database)
    hang_metrics = dict(metrics.get("hang_metrics") or {})
    rebuild_events = [
        event
        for event in progress_events
        if event.get("diagnostic_state") == "rebuilding_pool"
    ]
    interrupted_attempts = [
        attempt
        for attempt in attempts
        if attempt.get("status") == "interrupted"
        and attempt.get("error_code") == "PROCESS_WORKER_CRASH"
    ]
    rebuild_reason_recorded = bool(
        len(interrupted_attempts) == 1
        and str(
            interrupted_attempts[0].get("error_message") or ""
        ).strip()
    )
    worker_pids = {
        pid for pid in (crashed_pid, replacement_pid) if pid > 0
    }
    alive_after = _wait_for_processes_to_exit(worker_pids)
    pool_rebuild_count = int(
        hang_metrics.get("pool_rebuild_count") or 0
    )
    passed = bool(
        summary.failed == 0
        and marker
        and crashed_pid > 0
        and marker.get("model_state") == "ready"
        and int(marker.get("model_load_count") or 0) == 2
        and pool_rebuild_count == 1
        and rebuild_reason_recorded
        and replacement_pid > 0
        and replacement_pid != crashed_pid
        and request is not None
        and request["status"] == "confirmed"
        and replacement_load_count == 2
        and source_success
        and not alive_after
    )
    return {
        "passed": passed,
        "crash_observed": bool(marker),
        "crashed_worker_pid": crashed_pid,
        "crashed_worker_model_state": str(
            marker.get("model_state") or ""
        ),
        "crashed_worker_model_load_count": int(
            marker.get("model_load_count") or 0
        ),
        "crash_reason": str(marker.get("reason") or ""),
        "pool_rebuild_count": pool_rebuild_count,
        "rebuild_reason_recorded": rebuild_reason_recorded,
        "rebuild_reasons": [
            str(event.get("diagnostic_reason") or "")
            for event in rebuild_events
        ],
        "parse_attempts": attempts,
        "replacement_worker_pid": replacement_pid,
        "replacement_worker_changed": bool(
            replacement_pid > 0 and replacement_pid != crashed_pid
        ),
        "replacement_worker_model_load_count": (
            replacement_load_count
        ),
        "request_status": (
            str(request["status"]) if request is not None else ""
        ),
        "source_success": source_success,
        "worker_pids_after_completion": alive_after,
        "summary": summary.to_dict(),
    }


def validate_safe_pause(base: Path) -> dict[str, object]:
    from openpyxl import Workbook

    base.mkdir(parents=True, exist_ok=True)
    root = base / "safe-pause-xlsx"
    root.mkdir()
    source = root / "safe-pause-rows.xlsx"
    workbook = Workbook(write_only=True)
    sheet = workbook.create_sheet("SafePause")
    for row_number in range(1, 20_001):
        sheet.append(
            [
                row_number,
                f"SAFE_PAUSE_ROW_{row_number:05d}",
                "用于验证真实 XLSX 行批安全暂停与继续",
            ]
        )
    workbook.save(source)

    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    settings = AppSettings(
        enable_ocr=False,
        parser_workers=1,
        process_parser_workers=1,
        pdf_parser_workers=1,
        ocr_workers=1,
        slow_file_workers=1,
        xlsx_sheet_workers=1,
        index_write_batch_size=8,
    )
    manager = IndexManager(
        database,
        settings,
        run_context={
            "execution_mode": "normal",
            "validation_kind": "safe_pause_real_xlsx",
        },
    )
    token = CancelToken()
    progress_events: list[dict[str, object]] = []
    active_progress = threading.Event()
    finished = threading.Event()
    outcomes: list[object] = []
    errors: list[str] = []

    def progress(payload: dict[str, object]) -> None:
        progress_events.append(dict(payload))
        if (
            str(payload.get("current_file") or "").lower().endswith(
                ".xlsx"
            )
            and int(payload.get("active_completed_units") or 0) >= 250
        ):
            active_progress.set()

    def run_index() -> None:
        try:
            outcomes.append(
                manager.index_root(
                    root_id,
                    token,
                    progress,
                )
            )
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")
        finally:
            finished.set()

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
    )
    os.environ["LFTS_VALIDATION_SAFE_POINT_DELAY_MS"] = "80"
    worker = threading.Thread(
        target=run_index,
        name="lfts-safe-pause-real-xlsx",
    )
    worker.start()
    safe = False
    resumed = False
    observation_seconds = 0.0
    observation: dict[str, object] = {}
    status: dict[str, object] = {}
    requested_at = 0.0
    try:
        started = active_progress.wait(timeout=45.0)
        if not started:
            raise RuntimeError(
                "真实 XLSX 未在 45 秒内进入可观测的行批进度"
            )
        token.pause()
        requested_at = time.perf_counter()
        manager.request_pause()
        deadline = requested_at + 8.0
        while time.perf_counter() < deadline:
            if manager.is_safely_paused():
                safe = True
                break
            time.sleep(0.02)
        status = manager.pause_status()
        if not safe:
            raise RuntimeError(
                f"真实 XLSX 未在安全单位加 2 秒调度误差内暂停：{status}"
            )
        observation = _observe_real_paused_index(
            database,
            source,
            progress_events,
            seconds=5.0,
        )
        observation_seconds = float(
            observation["observation_seconds"]
        )
        manager.request_resume()
        token.resume()
        resumed = True
        if not finished.wait(timeout=90.0):
            raise RuntimeError("真实 XLSX 继续后未在 90 秒内完成")
    finally:
        if not resumed:
            manager.request_resume()
            token.resume()
        if not finished.wait(timeout=2.0):
            token.cancel(force=True)
            manager.force_terminate_processes()
        worker.join(timeout=5.0)
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
            ] = previous_delay

    final_summary = (
        outcomes[0]
        if outcomes
        else None
    )
    acknowledgement = next(
        (
            item
            for item in list(status.get("acknowledgements") or [])
            if str(item.get("safe_unit_type") or "") == "sheet_row"
        ),
        {},
    )
    acknowledged_cursor = int(
        acknowledgement.get("cursor") or 0
    )
    with database.connect() as connection:
        final_cursor = int(
            connection.execute(
                """
                SELECT COALESCE(
                    MAX(CAST(progress_cursor AS INTEGER)),
                    0
                )
                FROM parse_tasks
                WHERE parent_task_id IS NULL
                """
            ).fetchone()[0]
        )
        duplicate_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM (
                    SELECT file_id, block_index, COUNT(*) AS copies
                    FROM content_blocks
                    GROUP BY file_id, block_index
                    HAVING copies > 1
                )
                """
            ).fetchone()[0]
        )
    failed_delta = int(
        len(database.failed_files(limit=10_000))
        - int(observation.get("failed_files_before") or 0)
    )
    resume_cursor_advanced = bool(
        acknowledged_cursor > 0
        and final_cursor > acknowledged_cursor
    )
    xlsx_result = {
        "started": active_progress.is_set(),
        "safe_pause_confirmed": safe,
        "pause_latency_seconds": round(
            max(0.0, time.perf_counter() - requested_at)
            if not safe
            else float(status.get("pause_latency_seconds") or 0),
            6,
        ),
        "acknowledgement": acknowledgement,
        "observation_seconds": observation_seconds,
        "progress_delta": int(
            observation.get("progress_delta") or 0
        ),
        "database_write_delta": int(
            observation.get("database_write_delta") or 0
        ),
        "source_read_bytes_delta": int(
            observation.get("source_read_bytes_delta") or 0
        ),
        "paused_cpu_average": float(
            observation.get("paused_cpu_average") or 0.0
        ),
        "resume_cursor_advanced": resume_cursor_advanced,
        "acknowledged_cursor": acknowledged_cursor,
        "final_cursor": final_cursor,
        "duplicate_blocks": duplicate_blocks,
        "failed_delta": failed_delta,
        "source_unchanged": bool(
            observation.get("source_unchanged")
        ),
        "summary": (
            final_summary.to_dict()
            if final_summary is not None
            else {}
        ),
        "errors": errors,
    }
    passed = bool(
        not errors
        and final_summary is not None
        and final_summary.failed == 0
        and xlsx_result["started"]
        and safe
        and observation_seconds >= 5.0
        and xlsx_result["progress_delta"] == 0
        and xlsx_result["database_write_delta"] == 0
        and xlsx_result["source_read_bytes_delta"] == 0
        and xlsx_result["paused_cpu_average"] <= 5.0
        and resume_cursor_advanced
        and duplicate_blocks == 0
        and failed_delta == 0
        and xlsx_result["source_unchanged"]
    )
    pdf_result = _validate_pdf_native_page_pause(
        base / "pdf-native-page"
    )
    zip_result = _validate_zip_member_xlsx_pause(
        base / "zip-member-xlsx"
    )
    pptx_result = _validate_pptx_slide_pause(
        base / "pptx-slide"
    )
    image_ocr_result = _validate_image_ocr_pause(
        base / "image-ocr"
    )
    pdf_ocr_result = _validate_pdf_ocr_pause(
        base / "pdf-ocr"
    )
    legacy_result = _validate_legacy_office_pause(
        base / "legacy-office"
    )
    planning_results = _validate_planning_pause_gates(
        base / "planning"
    )
    return {
        "passed": bool(
            passed
            and pdf_result["passed"]
            and zip_result["passed"]
            and pptx_result["passed"]
            and image_ocr_result["passed"]
            and pdf_ocr_result["passed"]
            and legacy_result["passed"]
            and all(
                bool(item["passed"])
                for item in planning_results.values()
            )
        ),
        "mechanism_only": False,
        "status": status,
        "formats": {
            "xlsx": xlsx_result,
            "pdf_native_page": pdf_result,
            "zip_member_xlsx": zip_result,
            "pptx_slide": pptx_result,
            "image_ocr_batch": image_ocr_result,
            "pdf_ocr_region": pdf_ocr_result,
            "legacy_office_conversion": legacy_result,
        },
        "planning": planning_results,
        "resource_idle_real_format_gate_required": False,
        "remaining_format_gates": [],
    }


def _create_legacy_word_doc(source: Path) -> None:
    import gc

    import pythoncom
    import win32com.client

    word = None
    document = None
    pythoncom.CoInitialize()
    try:
        automation_errors: list[str] = []
        for prog_id in ("Word.Application", "KWPS.Application"):
            try:
                word = win32com.client.DispatchEx(prog_id)
                break
            except Exception as exc:
                automation_errors.append(f"{prog_id}: {exc}")
        if word is None:
            raise RuntimeError(
                "No Word-compatible automation server is available: "
                + "; ".join(automation_errors)
            )
        word.Visible = False
        word.DisplayAlerts = 0
        document = word.Documents.Add()
        validation_nonce = time.time_ns()
        paragraphs = [
            (
                f"LEGACY SAFE PAUSE PARAGRAPH {number:04d} "
                f"VALUE {number * 211} RUN {validation_nonce}"
            )
            for number in range(1, 601)
        ]
        paragraphs[399] += " LEGACY_SAFE_PAUSE_BLOCK_400"
        document.Content.Text = "\r".join(paragraphs)
        document.SaveAs2(str(source), FileFormat=0)
    finally:
        if document is not None:
            try:
                document.Close(False)
            except Exception:
                pass
        if word is not None:
            try:
                word.Quit()
            except Exception:
                pass
        document = None
        word = None
        gc.collect()


def _create_legacy_word_doc_isolated(source: Path) -> None:
    import subprocess
    import sys

    if getattr(sys, "frozen", False):
        command = [
            sys.executable,
            "--generate-legacy-word-doc",
            str(source),
        ]
    else:
        command = [
            sys.executable,
            "-c",
            (
                "import sys; from pathlib import Path; "
                "from local_full_text_search.core.phase2_validation "
                "import _create_legacy_word_doc; "
                "_create_legacy_word_doc(Path(sys.argv[1]))"
            ),
            str(source),
        ]
    generator = subprocess.run(
        command,
        cwd=str(Path.cwd()),
        capture_output=True,
        text=True,
        timeout=60.0,
        creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
    )
    if generator.returncode != 0:
        raise RuntimeError(
            "Unable to create the legacy Office validation file: "
            + (generator.stderr.strip() or generator.stdout.strip())
        )
    if source.read_bytes()[:8] != b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
        raise RuntimeError("Office did not create a valid OLE .doc validation file")


def _validate_legacy_office_pause(base: Path) -> dict[str, object]:
    import psutil

    from local_full_text_search.parsers.legacy_office_parser import (
        registered_office_processes_alive,
    )

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    source = root / "safe-pause-legacy.doc"
    _create_legacy_word_doc_isolated(source)

    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    manager = IndexManager(
        database,
        AppSettings(
            enable_ocr=False,
            enable_parse_cache=False,
            legacy_conversion_cache=True,
            fast_ooxml_enabled=True,
            parser_workers=1,
            process_parser_workers=1,
            slow_file_workers=1,
            index_write_batch_size=8,
        ),
        run_context={
            "execution_mode": "normal",
            "validation_kind": "safe_pause_legacy_office_conversion",
        },
    )
    token = CancelToken()
    progress_events: list[dict[str, object]] = []
    active_progress = threading.Event()
    finished = threading.Event()
    outcomes: list[object] = []
    errors: list[str] = []

    def progress(payload: dict[str, object]) -> None:
        event = dict(payload)
        progress_events.append(event)
        if (
            str(event.get("current_file") or "").lower().endswith(".doc")
            and str(event.get("active_phase") or "")
            == "legacy_office_open"
        ):
            active_progress.set()

    def run_index() -> None:
        try:
            outcomes.append(manager.index_root(root_id, token, progress))
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")
        finally:
            finished.set()

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
    )
    os.environ["LFTS_VALIDATION_SAFE_POINT_DELAY_MS"] = "80"
    worker = threading.Thread(
        target=run_index,
        name="lfts-safe-pause-legacy-office",
    )
    worker.start()
    safe = False
    resumed = False
    external_process_seen = False
    observed_office_pids: list[int] = []
    office_pids_while_paused: list[int] = []
    office_pids_after_completion: list[int] = []
    observation: dict[str, object] = {}
    status: dict[str, object] = {}
    spool_dir: Path | None = None
    try:
        if not active_progress.wait(timeout=45.0):
            raise RuntimeError(
                "Legacy Office did not enter the real conversion phase"
            )
        process_deadline = time.monotonic() + 15.0
        while time.monotonic() < process_deadline:
            spool_dir = manager._pause_spool_dir
            if spool_dir is not None:
                observed_office_pids = (
                    registered_office_processes_alive(spool_dir)
                )
                if observed_office_pids:
                    external_process_seen = True
                    break
            if finished.is_set():
                break
            time.sleep(0.02)
        if not external_process_seen:
            raise RuntimeError(
                "Legacy Office conversion did not register a live external process"
            )
        token.pause()
        manager.request_pause()
        deadline = time.monotonic() + 30.0
        while time.monotonic() < deadline:
            if manager.is_safely_paused():
                safe = True
                break
            time.sleep(0.05)
        status = manager.pause_status()
        if not safe:
            raise RuntimeError(
                f"Legacy Office did not reach a safe pause: {status}"
            )
        office_pids_while_paused = (
            registered_office_processes_alive(spool_dir)
            if spool_dir is not None
            else []
        )
        parse_pids = [
            int(item.get("worker_pid") or 0)
            for item in list(status.get("acknowledgements") or [])
            if int(item.get("worker_pid") or 0) > 0
        ]
        observation = _observe_real_paused_index(
            database,
            source,
            progress_events,
            seconds=5.0,
            read_process_ids=parse_pids,
        )
        manager.request_resume()
        token.resume()
        resumed = True
        if not finished.wait(timeout=120.0):
            raise RuntimeError("Legacy Office did not finish after resume")
    finally:
        if not resumed:
            manager.request_resume()
            token.resume()
        if not finished.wait(timeout=2.0):
            token.cancel(force=True)
            manager.force_terminate_processes()
        worker.join(timeout=5.0)
        if spool_dir is not None:
            office_pids_after_completion = (
                registered_office_processes_alive(spool_dir)
            )
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
            ] = previous_delay

    acknowledgement = next(
        (
            dict(item)
            for item in list(status.get("acknowledgements") or [])
            if str(item.get("safe_unit_type") or "").startswith(
                "legacy_"
            )
        ),
        {},
    )
    acknowledged_cursor = int(acknowledgement.get("cursor") or 0)
    worker_pid = int(acknowledgement.get("worker_pid") or 0)
    with database.connect() as connection:
        final_cursor = int(
            connection.execute(
                """
                SELECT COALESCE(MAX(CAST(progress_cursor AS INTEGER)), 0)
                FROM parse_tasks
                """
            ).fetchone()[0]
        )
        file_row = connection.execute(
            """
            SELECT parse_status FROM files
            WHERE source_kind = 'file' AND is_deleted = 0
            """
        ).fetchone()
        search_token_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM content_blocks
                WHERE raw_text LIKE '%LEGACY_SAFE_PAUSE_BLOCK_400%'
                """
            ).fetchone()[0]
        )
        converter_rows = [
            json.loads(str(row[0] or "{}"))
            for row in connection.execute(
                """
                SELECT extra_json FROM content_blocks
                WHERE raw_text <> ''
                """
            )
        ]
        duplicate_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM (
                    SELECT file_id, block_index, COUNT(*) AS copies
                    FROM content_blocks
                    GROUP BY file_id, block_index
                    HAVING copies > 1
                )
                """
            ).fetchone()[0]
        )
    final_summary = outcomes[0] if outcomes else None
    result = {
        "started": active_progress.is_set(),
        "external_process_seen": external_process_seen,
        "observed_office_pids": observed_office_pids,
        "safe_pause_confirmed": safe,
        "office_pids_while_paused": office_pids_while_paused,
        "status": status,
        "acknowledgement": acknowledgement,
        "observation_seconds": float(
            observation.get("observation_seconds") or 0.0
        ),
        "progress_delta": int(observation.get("progress_delta") or 0),
        "database_write_delta": int(
            observation.get("database_write_delta") or 0
        ),
        "source_read_bytes_delta": int(
            observation.get("source_read_bytes_delta") or 0
        ),
        "paused_cpu_average": float(
            observation.get("paused_cpu_average") or 0.0
        ),
        "source_unchanged": bool(observation.get("source_unchanged")),
        "acknowledged_cursor": acknowledged_cursor,
        "final_cursor": final_cursor,
        "resume_cursor_advanced": bool(
            final_cursor > acknowledged_cursor
        ),
        "source_success": bool(
            file_row is not None and str(file_row[0]) == "success"
        ),
        "search_token_blocks": search_token_blocks,
        "conversion_cache_reused_after_resume": bool(
            converter_rows
            and all(
                str(item.get("legacy_converter") or "") == "cache"
                for item in converter_rows
            )
        ),
        "duplicate_blocks": duplicate_blocks,
        "worker_pid": worker_pid,
        "worker_exited": bool(
            worker_pid > 0 and not psutil.pid_exists(worker_pid)
        ),
        "office_pids_after_completion": office_pids_after_completion,
        "summary": (
            final_summary.to_dict()
            if final_summary is not None
            else {}
        ),
        "errors": errors,
    }
    result["passed"] = bool(
        not errors
        and final_summary is not None
        and final_summary.failed == 0
        and result["started"]
        and external_process_seen
        and safe
        and not office_pids_while_paused
        and acknowledgement
        and result["observation_seconds"] >= 5.0
        and result["progress_delta"] == 0
        and result["database_write_delta"] == 0
        and result["source_read_bytes_delta"] == 0
        and result["paused_cpu_average"] <= 5.0
        and result["source_unchanged"]
        and result["resume_cursor_advanced"]
        and result["source_success"]
        and search_token_blocks == 1
        and result["conversion_cache_reused_after_resume"]
        and duplicate_blocks == 0
        and result["worker_exited"]
        and not office_pids_after_completion
    )
    return result


def _validate_pdf_ocr_pause(base: Path) -> dict[str, object]:
    import fitz
    import psutil
    from PIL import Image, ImageDraw, ImageFont

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    scan_image = base / "scan-page.png"
    image = Image.new("RGB", (1500, 1100), "white")
    draw = ImageDraw.Draw(image)
    font_path = (
        Path(os.environ.get("WINDIR", "C:/Windows"))
        / "Fonts"
        / "arial.ttf"
    )
    font = (
        ImageFont.truetype(str(font_path), 42)
        if font_path.is_file()
        else ImageFont.load_default()
    )
    for line_number in range(1, 19):
        draw.text(
            (65, 30 + (line_number - 1) * 56),
            f"PDF OCR SAFE PAUSE LINE {line_number:02d} CODE {line_number * 193}",
            fill="black",
            font=font,
        )
    validation_marker = hashlib.sha256(
        str(base.resolve()).encode("utf-8")
    ).hexdigest()[:12]
    draw.text(
        (65, 1035),
        f"RUN {validation_marker}",
        fill="black",
        font=font,
    )
    image.save(scan_image)
    source = root / "safe-pause-scanned.pdf"
    page_count = 4
    document = fitz.open()
    try:
        for _ in range(page_count):
            page = document.new_page(width=750, height=550)
            page.insert_image(page.rect, filename=str(scan_image))
        document.save(source)
    finally:
        document.close()

    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    manager = IndexManager(
        database,
        AppSettings(
            enable_ocr=True,
            ocr_images=False,
            ocr_scanned_pdf=True,
            enable_parse_cache=False,
            parser_workers=1,
            process_parser_workers=1,
            pdf_parser_workers=1,
            ocr_workers=1,
            ocr_pending_tasks=1,
            index_write_batch_size=8,
        ),
        run_context={
            "execution_mode": "normal",
            "validation_kind": "safe_pause_pdf_ocr_region",
        },
    )
    token = CancelToken()
    progress_events: list[dict[str, object]] = []
    active_progress = threading.Event()
    finished = threading.Event()
    outcomes: list[object] = []
    errors: list[str] = []

    def progress(payload: dict[str, object]) -> None:
        event = dict(payload)
        progress_events.append(event)
        phase = str(event.get("active_phase") or "")
        if (
            str(event.get("current_file") or "")
            .lower()
            .endswith(".pdf")
            and (
                phase.startswith("pdf_embedded_")
                or phase.startswith("pdf_ocr_")
                or phase.startswith("pdf_region_")
            )
        ):
            active_progress.set()

    def run_index() -> None:
        try:
            outcomes.append(manager.index_root(root_id, token, progress))
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")
        finally:
            finished.set()

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
    )
    os.environ["LFTS_VALIDATION_SAFE_POINT_DELAY_MS"] = "250"
    worker = threading.Thread(
        target=run_index,
        name="lfts-safe-pause-pdf-ocr",
    )
    worker.start()
    safe = False
    resumed = False
    observation: dict[str, object] = {}
    status: dict[str, object] = {}
    pending_pages_while_paused = 0
    try:
        if not active_progress.wait(timeout=90.0):
            raise RuntimeError(
                "PDF OCR 未在 90 秒内进入页内区域进展"
            )
        token.pause()
        manager.request_pause()
        deadline = time.monotonic() + 60.0
        while time.monotonic() < deadline:
            if manager.is_safely_paused():
                safe = True
                break
            time.sleep(0.05)
        status = manager.pause_status()
        if not safe:
            raise RuntimeError(
                f"PDF OCR 未在当前区域批完成后暂停：{status}"
            )
        with database.connect() as connection:
            pending_pages_while_paused = int(
                connection.execute(
                    """
                    SELECT COUNT(*) FROM parse_tasks
                    WHERE task_type = 'pdf_ocr_page'
                      AND status IN ('queued', 'running', 'paused')
                    """
                ).fetchone()[0]
            )
        parse_pids = [
            int(item.get("worker_pid") or 0)
            for item in list(status.get("acknowledgements") or [])
            if int(item.get("worker_pid") or 0) > 0
        ]
        observation = _observe_real_paused_index(
            database,
            source,
            progress_events,
            seconds=5.0,
            read_process_ids=parse_pids,
        )
        manager.request_resume()
        token.resume()
        resumed = True
        if not finished.wait(timeout=300.0):
            raise RuntimeError("PDF OCR 继续后未完成")
    finally:
        if not resumed:
            manager.request_resume()
            token.resume()
        if not finished.wait(timeout=2.0):
            token.cancel(force=True)
            manager.force_terminate_processes()
        worker.join(timeout=5.0)
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
            ] = previous_delay

    acknowledgement = next(
        (
            dict(item)
            for item in list(status.get("acknowledgements") or [])
            if str(item.get("safe_unit_type") or "")
            == "pdf_ocr_page"
        ),
        {},
    )
    acknowledged_cursor = int(acknowledgement.get("cursor") or 0)
    worker_pid = int(acknowledgement.get("worker_pid") or 0)
    with database.connect() as connection:
        final_cursor = int(
            connection.execute(
                """
                SELECT COALESCE(MAX(CAST(progress_cursor AS INTEGER)), 0)
                FROM parse_tasks
                WHERE task_type = 'pdf_ocr_page'
                """
            ).fetchone()[0]
        )
        completed_ocr_pages = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'pdf_ocr_page'
                  AND status = 'complete'
                  AND confirmed_at IS NOT NULL
                """
            ).fetchone()[0]
        )
        merge_tasks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'document_merge'
                  AND status = 'complete'
                  AND confirmed_at IS NOT NULL
                """
            ).fetchone()[0]
        )
        ordered_page_numbers = [
            int(row[0])
            for row in connection.execute(
                """
                SELECT page_number FROM content_blocks
                WHERE block_type = 'pdf_page_ocr'
                ORDER BY page_number
                """
            )
        ]
        ocr_text_chars = int(
            connection.execute(
                """
                SELECT COALESCE(SUM(LENGTH(raw_text)), 0)
                FROM content_blocks
                WHERE block_type = 'pdf_page_ocr'
                """
            ).fetchone()[0]
        )
        duplicate_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM (
                    SELECT file_id, block_index, COUNT(*) AS copies
                    FROM content_blocks
                    GROUP BY file_id, block_index
                    HAVING copies > 1
                )
                """
            ).fetchone()[0]
        )
        source_success = bool(
            connection.execute(
                """
                SELECT COUNT(*) FROM files
                WHERE source_kind = 'file'
                  AND is_deleted = 0
                  AND parse_status = 'success'
                """
            ).fetchone()[0]
        )
    final_summary = outcomes[0] if outcomes else None
    worker_exited = bool(
        worker_pid > 0 and not psutil.pid_exists(worker_pid)
    )
    result = {
        "started": active_progress.is_set(),
        "safe_pause_confirmed": safe,
        "status": status,
        "acknowledgement": acknowledgement,
        "observation_seconds": float(
            observation.get("observation_seconds") or 0.0
        ),
        "progress_delta": int(observation.get("progress_delta") or 0),
        "database_write_delta": int(
            observation.get("database_write_delta") or 0
        ),
        "source_read_bytes_delta": int(
            observation.get("source_read_bytes_delta") or 0
        ),
        "paused_cpu_average": float(
            observation.get("paused_cpu_average") or 0.0
        ),
        "source_unchanged": bool(observation.get("source_unchanged")),
        "pending_pages_while_paused": pending_pages_while_paused,
        "acknowledged_cursor": acknowledged_cursor,
        "final_cursor": final_cursor,
        "resume_cursor_advanced": bool(
            acknowledged_cursor > 0 and final_cursor > acknowledged_cursor
        ),
        "page_count": page_count,
        "completed_ocr_pages": completed_ocr_pages,
        "merge_tasks": merge_tasks,
        "ordered_pages": len(ordered_page_numbers),
        "ocr_text_chars": ocr_text_chars,
        "source_success": source_success,
        "duplicate_blocks": duplicate_blocks,
        "worker_pid": worker_pid,
        "worker_exited": worker_exited,
        "summary": (
            final_summary.to_dict()
            if final_summary is not None
            else {}
        ),
        "errors": errors,
    }
    result["passed"] = bool(
        not errors
        and final_summary is not None
        and final_summary.failed == 0
        and result["started"]
        and safe
        and acknowledgement
        and result["observation_seconds"] >= 5.0
        and result["progress_delta"] == 0
        and result["database_write_delta"] == 0
        and result["source_read_bytes_delta"] == 0
        and result["paused_cpu_average"] <= 5.0
        and result["source_unchanged"]
        and pending_pages_while_paused > 0
        and result["resume_cursor_advanced"]
        and completed_ocr_pages == page_count
        and merge_tasks == 1
        and ordered_page_numbers == list(range(1, page_count + 1))
        and ocr_text_chars > 0
        and source_success
        and duplicate_blocks == 0
        and worker_exited
    )
    return result


def _validate_image_ocr_pause(base: Path) -> dict[str, object]:
    import psutil
    from PIL import Image, ImageDraw, ImageFont

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    source = root / "safe-pause-image.png"
    image = Image.new("RGB", (1800, 1400), "white")
    draw = ImageDraw.Draw(image)
    font_path = (
        Path(os.environ.get("WINDIR", "C:/Windows"))
        / "Fonts"
        / "arial.ttf"
    )
    font = (
        ImageFont.truetype(str(font_path), 44)
        if font_path.is_file()
        else ImageFont.load_default()
    )
    for line_number in range(1, 25):
        draw.text(
            (70, 35 + (line_number - 1) * 54),
            f"IMAGE SAFE PAUSE LINE {line_number:02d} VALUE {line_number * 137}",
            fill="black",
            font=font,
        )
    validation_marker = hashlib.sha256(
        str(base.resolve()).encode("utf-8")
    ).hexdigest()[:12]
    draw.text(
        (70, 1330),
        f"RUN {validation_marker}",
        fill="black",
        font=font,
    )
    image.save(source)

    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    manager = IndexManager(
        database,
        AppSettings(
            enable_ocr=True,
            ocr_images=True,
            ocr_scanned_pdf=False,
            min_ocr_image_pixels=0,
            enable_parse_cache=False,
            parser_workers=1,
            process_parser_workers=1,
            ocr_workers=1,
            ocr_pending_tasks=1,
            slow_file_workers=1,
            index_write_batch_size=8,
        ),
        run_context={
            "execution_mode": "normal",
            "validation_kind": "safe_pause_image_ocr_batch",
        },
    )
    token = CancelToken()
    progress_events: list[dict[str, object]] = []
    active_progress = threading.Event()
    finished = threading.Event()
    outcomes: list[object] = []
    errors: list[str] = []

    def progress(payload: dict[str, object]) -> None:
        event = dict(payload)
        progress_events.append(event)
        if (
            str(event.get("current_file") or "")
            .lower()
            .endswith(".png")
            and str(event.get("active_phase") or "").startswith(
                "ocr_"
            )
        ):
            active_progress.set()

    def run_index() -> None:
        try:
            outcomes.append(manager.index_root(root_id, token, progress))
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")
        finally:
            finished.set()

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
    )
    os.environ["LFTS_VALIDATION_SAFE_POINT_DELAY_MS"] = "250"
    worker = threading.Thread(
        target=run_index,
        name="lfts-safe-pause-image-ocr",
    )
    worker.start()
    safe = False
    resumed = False
    observation: dict[str, object] = {}
    status: dict[str, object] = {}
    try:
        if not active_progress.wait(timeout=90.0):
            raise RuntimeError(
                "图片 OCR 未在 90 秒内进入检测或识别批进展"
            )
        token.pause()
        manager.request_pause()
        deadline = time.monotonic() + 60.0
        while time.monotonic() < deadline:
            if manager.is_safely_paused():
                safe = True
                break
            time.sleep(0.05)
        status = manager.pause_status()
        if not safe:
            raise RuntimeError(
                f"图片 OCR 未在当前推理批完成后暂停：{status}"
            )
        parse_pids = [
            int(item.get("worker_pid") or 0)
            for item in list(status.get("acknowledgements") or [])
            if int(item.get("worker_pid") or 0) > 0
        ]
        observation = _observe_real_paused_index(
            database,
            source,
            progress_events,
            seconds=5.0,
            read_process_ids=parse_pids,
        )
        manager.request_resume()
        token.resume()
        resumed = True
        if not finished.wait(timeout=240.0):
            raise RuntimeError("图片 OCR 继续后未完成")
    finally:
        if not resumed:
            manager.request_resume()
            token.resume()
        if not finished.wait(timeout=2.0):
            token.cancel(force=True)
            manager.force_terminate_processes()
        worker.join(timeout=5.0)
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
            ] = previous_delay

    acknowledgement = next(
        (
            dict(item)
            for item in list(status.get("acknowledgements") or [])
            if str(item.get("safe_unit_type") or "").startswith(
                "ocr_"
            )
        ),
        {},
    )
    acknowledged_cursor = int(acknowledgement.get("cursor") or 0)
    worker_pid = int(acknowledgement.get("worker_pid") or 0)
    with database.connect() as connection:
        final_cursor = int(
            connection.execute(
                """
                SELECT COALESCE(MAX(CAST(progress_cursor AS INTEGER)), 0)
                FROM parse_tasks
                """
            ).fetchone()[0]
        )
        file_row = connection.execute(
            """
            SELECT parse_status FROM files
            WHERE source_kind = 'file' AND is_deleted = 0
            """
        ).fetchone()
        block_row = connection.execute(
            """
            SELECT COALESCE(SUM(LENGTH(raw_text)), 0), extra_json
            FROM content_blocks
            WHERE block_type = 'image_ocr'
            """
        ).fetchone()
        duplicate_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM (
                    SELECT file_id, block_index, COUNT(*) AS copies
                    FROM content_blocks
                    GROUP BY file_id, block_index
                    HAVING copies > 1
                )
                """
            ).fetchone()[0]
        )
    extra = json.loads(str(block_row[1] or "{}")) if block_row else {}
    final_summary = outcomes[0] if outcomes else None
    worker_exited = bool(
        worker_pid > 0 and not psutil.pid_exists(worker_pid)
    )
    result = {
        "started": active_progress.is_set(),
        "safe_pause_confirmed": safe,
        "status": status,
        "acknowledgement": acknowledgement,
        "observation_seconds": float(
            observation.get("observation_seconds") or 0.0
        ),
        "progress_delta": int(observation.get("progress_delta") or 0),
        "database_write_delta": int(
            observation.get("database_write_delta") or 0
        ),
        "source_read_bytes_delta": int(
            observation.get("source_read_bytes_delta") or 0
        ),
        "paused_cpu_average": float(
            observation.get("paused_cpu_average") or 0.0
        ),
        "source_unchanged": bool(observation.get("source_unchanged")),
        "acknowledged_cursor": acknowledged_cursor,
        "final_cursor": final_cursor,
        "resume_cursor_advanced": bool(
            acknowledged_cursor > 0 and final_cursor > acknowledged_cursor
        ),
        "source_success": bool(
            file_row is not None and str(file_row[0]) == "success"
        ),
        "ocr_text_chars": int(block_row[0] or 0) if block_row else 0,
        "engine": str(extra.get("engine") or ""),
        "duplicate_blocks": duplicate_blocks,
        "worker_pid": worker_pid,
        "worker_exited": worker_exited,
        "summary": (
            final_summary.to_dict()
            if final_summary is not None
            else {}
        ),
        "errors": errors,
    }
    result["passed"] = bool(
        not errors
        and final_summary is not None
        and final_summary.failed == 0
        and result["started"]
        and safe
        and acknowledgement
        and result["observation_seconds"] >= 5.0
        and result["progress_delta"] == 0
        and result["database_write_delta"] == 0
        and result["source_read_bytes_delta"] == 0
        and result["paused_cpu_average"] <= 5.0
        and result["source_unchanged"]
        and result["resume_cursor_advanced"]
        and result["source_success"]
        and result["ocr_text_chars"] > 0
        and result["engine"] == "PaddleOCR"
        and duplicate_blocks == 0
        and worker_exited
    )
    return result


def _validate_pptx_slide_pause(base: Path) -> dict[str, object]:
    from pptx import Presentation
    from pptx.util import Inches

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    source = root / "safe-pause-slides.pptx"
    slide_count = 120
    presentation = Presentation()
    for slide_number in range(1, slide_count + 1):
        slide = presentation.slides.add_slide(
            presentation.slide_layouts[6]
        )
        box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(1),
        )
        box.text = f"PPTX_SAFE_PAUSE_SLIDE_{slide_number:03d}"
        slide.notes_slide.notes_text_frame.text = (
            f"PPTX_SAFE_PAUSE_NOTE_{slide_number:03d}"
        )
    presentation.save(source)

    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    manager = IndexManager(
        database,
        AppSettings(
            enable_ocr=False,
            parser_workers=1,
            process_parser_workers=1,
            slow_file_workers=1,
            fast_ooxml_enabled=True,
            index_write_batch_size=8,
        ),
        run_context={
            "execution_mode": "normal",
            "validation_kind": "safe_pause_pptx_slide",
        },
    )
    token = CancelToken()
    progress_events: list[dict[str, object]] = []
    active_progress = threading.Event()
    finished = threading.Event()
    outcomes: list[object] = []
    errors: list[str] = []

    def progress(payload: dict[str, object]) -> None:
        event = dict(payload)
        progress_events.append(event)
        if (
            str(event.get("current_file") or "")
            .lower()
            .endswith(".pptx")
            and str(event.get("active_phase") or "")
            == "pptx_slide"
            and int(event.get("active_completed_units") or 0) >= 5
        ):
            active_progress.set()

    def run_index() -> None:
        try:
            outcomes.append(manager.index_root(root_id, token, progress))
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")
        finally:
            finished.set()

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
    )
    os.environ["LFTS_VALIDATION_SAFE_POINT_DELAY_MS"] = "80"
    worker = threading.Thread(
        target=run_index,
        name="lfts-safe-pause-pptx-slide",
    )
    worker.start()
    safe = False
    resumed = False
    observation: dict[str, object] = {}
    status: dict[str, object] = {}
    try:
        if not active_progress.wait(timeout=45.0):
            raise RuntimeError(
                "PPTX 未在 45 秒内进入幻灯片语义进展"
            )
        token.pause()
        manager.request_pause()
        deadline = time.monotonic() + 8.0
        while time.monotonic() < deadline:
            if manager.is_safely_paused():
                safe = True
                break
            time.sleep(0.02)
        status = manager.pause_status()
        if not safe:
            raise RuntimeError(
                f"PPTX 未在整张幻灯片后暂停：{status}"
            )
        observation = _observe_real_paused_index(
            database,
            source,
            progress_events,
            seconds=5.0,
        )
        manager.request_resume()
        token.resume()
        resumed = True
        if not finished.wait(timeout=120.0):
            raise RuntimeError("PPTX 继续后未完成")
    finally:
        if not resumed:
            manager.request_resume()
            token.resume()
        if not finished.wait(timeout=2.0):
            token.cancel(force=True)
            manager.force_terminate_processes()
        worker.join(timeout=5.0)
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
            ] = previous_delay

    acknowledgement = next(
        (
            dict(item)
            for item in list(status.get("acknowledgements") or [])
            if str(item.get("safe_unit_type") or "")
            == "pptx_slide"
        ),
        {},
    )
    acknowledged_cursor = int(acknowledgement.get("cursor") or 0)
    with database.connect() as connection:
        final_cursor = int(
            connection.execute(
                """
                SELECT COALESCE(MAX(CAST(progress_cursor AS INTEGER)), 0)
                FROM parse_tasks
                """
            ).fetchone()[0]
        )
        ordered_slide_numbers = [
            int(row[0])
            for row in connection.execute(
                """
                SELECT DISTINCT slide_number FROM content_blocks
                WHERE slide_number IS NOT NULL
                ORDER BY slide_number
                """
            )
        ]
        search_token_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM content_blocks
                WHERE raw_text LIKE '%PPTX_SAFE_PAUSE_SLIDE_100%'
                """
            ).fetchone()[0]
        )
        duplicate_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM (
                    SELECT file_id, block_index, COUNT(*) AS copies
                    FROM content_blocks
                    GROUP BY file_id, block_index
                    HAVING copies > 1
                )
                """
            ).fetchone()[0]
        )
    final_summary = outcomes[0] if outcomes else None
    resume_cursor_advanced = bool(
        acknowledged_cursor > 0
        and final_cursor > acknowledged_cursor
    )
    result = {
        "started": active_progress.is_set(),
        "safe_pause_confirmed": safe,
        "status": status,
        "acknowledgement": acknowledgement,
        "observation_seconds": float(
            observation.get("observation_seconds") or 0.0
        ),
        "progress_delta": int(observation.get("progress_delta") or 0),
        "database_write_delta": int(
            observation.get("database_write_delta") or 0
        ),
        "source_read_bytes_delta": int(
            observation.get("source_read_bytes_delta") or 0
        ),
        "paused_cpu_average": float(
            observation.get("paused_cpu_average") or 0.0
        ),
        "source_unchanged": bool(observation.get("source_unchanged")),
        "acknowledged_cursor": acknowledged_cursor,
        "final_cursor": final_cursor,
        "resume_cursor_advanced": resume_cursor_advanced,
        "slide_count": slide_count,
        "ordered_slides": len(ordered_slide_numbers),
        "search_token_blocks": search_token_blocks,
        "duplicate_blocks": duplicate_blocks,
        "summary": (
            final_summary.to_dict()
            if final_summary is not None
            else {}
        ),
        "errors": errors,
    }
    result["passed"] = bool(
        not errors
        and final_summary is not None
        and final_summary.failed == 0
        and result["started"]
        and safe
        and acknowledgement
        and result["observation_seconds"] >= 5.0
        and result["progress_delta"] == 0
        and result["database_write_delta"] == 0
        and result["source_read_bytes_delta"] == 0
        and result["paused_cpu_average"] <= 5.0
        and result["source_unchanged"]
        and resume_cursor_advanced
        and ordered_slide_numbers == list(range(1, slide_count + 1))
        and search_token_blocks == 1
        and duplicate_blocks == 0
    )
    return result


def _validate_planning_pause_gates(
    base: Path,
) -> dict[str, dict[str, object]]:
    base.mkdir(parents=True, exist_ok=True)

    hash_root = base / "content-hash" / "files"
    hash_root.mkdir(parents=True)
    hash_source = hash_root / "safe-pause-hash.txt"
    chunk = (
        "HASH_SAFE_PAUSE_CONTENT line with deterministic bytes\n"
    ).encode("utf-8")
    target_size = 8 * 1024 * 1024
    with hash_source.open("wb") as stream:
        remaining = target_size
        while remaining > 0:
            payload = chunk[:remaining]
            stream.write(payload)
            remaining -= len(payload)
    hash_result = _validate_one_planning_pause(
        base / "content-hash",
        hash_root,
        hash_source,
        phase="content_hash",
        minimum_progress=2 * 1024 * 1024,
        expected_files=1,
    )

    directory_root = base / "directory" / "files"
    directory_root.mkdir(parents=True)
    source_files = 80
    for directory_number in range(1, source_files + 1):
        directory = directory_root / f"batch-{directory_number:03d}"
        directory.mkdir()
        (directory / f"document-{directory_number:03d}.txt").write_text(
            f"DIRECTORY_SAFE_PAUSE_FILE_{directory_number:03d}",
            encoding="utf-8",
        )
    directory_result = _validate_one_planning_pause(
        base / "directory",
        directory_root,
        directory_root,
        phase="directory_enumeration",
        minimum_progress=5,
        expected_files=source_files,
    )
    return {
        "content_hash": hash_result,
        "directory_enumeration": directory_result,
    }


def _validate_one_planning_pause(
    base: Path,
    root: Path,
    observed_source: Path,
    *,
    phase: str,
    minimum_progress: int,
    expected_files: int,
) -> dict[str, object]:
    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    manager = IndexManager(
        database,
        AppSettings(
            enable_ocr=False,
            parser_workers=1,
            process_parser_workers=1,
            slow_file_workers=1,
            planning_discovery_batch_size=128,
            index_write_batch_size=8,
        ),
        run_context={
            "execution_mode": "normal",
            "validation_kind": f"safe_pause_{phase}",
        },
    )
    token = CancelToken()
    progress_events: list[dict[str, object]] = []
    active_progress = threading.Event()
    finished = threading.Event()
    outcomes: list[object] = []
    errors: list[str] = []

    def progress(payload: dict[str, object]) -> None:
        event = dict(payload)
        progress_events.append(event)
        if str(event.get("planning_phase") or "") != phase:
            return
        value = (
            int(event.get("planning_bytes_read") or 0)
            if phase == "content_hash"
            else int(event.get("planning_completed") or 0)
        )
        if value >= minimum_progress:
            active_progress.set()

    def run_index() -> None:
        try:
            outcomes.append(manager.index_root(root_id, token, progress))
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")
        finally:
            finished.set()

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_PLANNING_SAFE_POINT_DELAY_MS"
    )
    os.environ[
        "LFTS_VALIDATION_PLANNING_SAFE_POINT_DELAY_MS"
    ] = "80"
    worker = threading.Thread(
        target=run_index,
        name=f"lfts-safe-pause-{phase}",
    )
    worker.start()
    safe = False
    resumed = False
    observation: dict[str, object] = {}
    status: dict[str, object] = {}
    try:
        if not active_progress.wait(timeout=45.0):
            raise RuntimeError(
                f"{phase} 未在 45 秒内进入可观测语义进展"
            )
        token.pause()
        manager.request_pause()
        deadline = time.monotonic() + 8.0
        while time.monotonic() < deadline:
            if manager.is_safely_paused():
                safe = True
                break
            time.sleep(0.02)
        status = manager.pause_status()
        if not safe:
            raise RuntimeError(
                f"{phase} 未在安全单位后暂停：{status}"
            )
        observation = _observe_real_paused_index(
            database,
            observed_source,
            progress_events,
            seconds=5.0,
            read_process_ids=[
                int(item.get("worker_pid") or 0)
                for item in list(
                    status.get("planning_acknowledgements") or []
                )
                if int(item.get("worker_pid") or 0) > 0
            ],
        )
        manager.request_resume()
        token.resume()
        resumed = True
        if not finished.wait(timeout=120.0):
            raise RuntimeError(f"{phase} 继续后未完成")
    finally:
        if not resumed:
            manager.request_resume()
            token.resume()
        if not finished.wait(timeout=2.0):
            token.cancel(force=True)
            manager.force_terminate_processes()
        worker.join(timeout=5.0)
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_PLANNING_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_PLANNING_SAFE_POINT_DELAY_MS"
            ] = previous_delay

    acknowledgements = [
        dict(item)
        for item in list(
            status.get("planning_acknowledgements") or []
        )
        if str(item.get("safe_unit_type") or "") == phase
    ]
    acknowledgement = acknowledgements[0] if acknowledgements else {}
    acknowledged_cursor = _planning_cursor_value(
        str(acknowledgement.get("cursor") or "")
    )
    relevant_events = [
        event
        for event in progress_events
        if str(event.get("planning_phase") or "") == phase
    ]
    final_cursor = max(
        [
            int(event.get("planning_bytes_read") or 0)
            if phase == "content_hash"
            else int(event.get("planning_completed") or 0)
            for event in relevant_events
        ],
        default=0,
    )
    with database.connect() as connection:
        indexed_files = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM files
                WHERE source_kind = 'file'
                  AND is_deleted = 0
                  AND parse_status = 'success'
                """
            ).fetchone()[0]
        )
        duplicate_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM (
                    SELECT file_id, block_index, COUNT(*) AS copies
                    FROM content_blocks
                    GROUP BY file_id, block_index
                    HAVING copies > 1
                )
                """
            ).fetchone()[0]
        )
    final_summary = outcomes[0] if outcomes else None
    resume_cursor_advanced = bool(
        acknowledged_cursor >= 0
        and final_cursor > acknowledged_cursor
    )
    result = {
        "started": active_progress.is_set(),
        "safe_pause_confirmed": safe,
        "status": status,
        "acknowledgements": acknowledgements,
        "observation_seconds": float(
            observation.get("observation_seconds") or 0.0
        ),
        "progress_delta": int(observation.get("progress_delta") or 0),
        "database_write_delta": int(
            observation.get("database_write_delta") or 0
        ),
        "source_read_bytes_delta": int(
            observation.get("source_read_bytes_delta") or 0
        ),
        "paused_cpu_average": float(
            observation.get("paused_cpu_average") or 0.0
        ),
        "source_unchanged": bool(observation.get("source_unchanged")),
        "acknowledged_cursor": acknowledged_cursor,
        "final_cursor": final_cursor,
        "resume_cursor_advanced": resume_cursor_advanced,
        "source_files": expected_files,
        "indexed_files": indexed_files,
        "source_success": indexed_files == expected_files,
        "duplicate_blocks": duplicate_blocks,
        "summary": (
            final_summary.to_dict()
            if final_summary is not None
            else {}
        ),
        "errors": errors,
    }
    result["passed"] = bool(
        not errors
        and final_summary is not None
        and final_summary.failed == 0
        and result["started"]
        and safe
        and acknowledgements
        and result["observation_seconds"] >= 5.0
        and result["progress_delta"] == 0
        and result["database_write_delta"] == 0
        and result["source_read_bytes_delta"] == 0
        and result["paused_cpu_average"] <= 5.0
        and result["source_unchanged"]
        and resume_cursor_advanced
        and indexed_files == expected_files
        and duplicate_blocks == 0
    )
    return result


def _planning_cursor_value(cursor: str) -> int:
    try:
        return int(str(cursor).rsplit(":", 1)[-1])
    except ValueError:
        return -1


def _validate_zip_member_xlsx_pause(
    base: Path,
) -> dict[str, object]:
    import zipfile

    from openpyxl import Workbook

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    member_source = base / "safe-member.xlsx"
    workbook = Workbook(write_only=True)
    sheet = workbook.create_sheet("ZipPause")
    for row_number in range(1, 12_001):
        sheet.append(
            [
                row_number,
                f"ZIP_SAFE_PAUSE_ROW_{row_number:05d}",
                "ZIP 成员行批安全暂停",
            ]
        )
    workbook.save(member_source)
    source = root / "safe-pause-member.zip"
    with zipfile.ZipFile(
        source,
        "w",
        compression=zipfile.ZIP_DEFLATED,
    ) as archive:
        archive.write(
            member_source,
            arcname="nested/safe-member.xlsx",
        )

    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    manager = IndexManager(
        database,
        AppSettings(
            enable_ocr=False,
            parser_workers=1,
            process_parser_workers=1,
            slow_file_workers=1,
            xlsx_sheet_workers=1,
            index_write_batch_size=8,
        ),
        run_context={
            "execution_mode": "normal",
            "validation_kind": "safe_pause_zip_member_xlsx",
        },
    )
    token = CancelToken()
    progress_events: list[dict[str, object]] = []
    active_progress = threading.Event()
    finished = threading.Event()
    outcomes: list[object] = []
    errors: list[str] = []

    def progress(payload: dict[str, object]) -> None:
        progress_events.append(dict(payload))
        if (
            str(payload.get("current_file") or "")
            .lower()
            .endswith(".xlsx")
            and int(payload.get("active_completed_units") or 0)
            >= 250
        ):
            active_progress.set()

    def run_index() -> None:
        try:
            outcomes.append(manager.index_root(root_id, token, progress))
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")
        finally:
            finished.set()

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
    )
    os.environ["LFTS_VALIDATION_SAFE_POINT_DELAY_MS"] = "80"
    worker = threading.Thread(
        target=run_index,
        name="lfts-safe-pause-zip-member-xlsx",
    )
    worker.start()
    safe = False
    resumed = False
    observation: dict[str, object] = {}
    status: dict[str, object] = {}
    try:
        if not active_progress.wait(timeout=45.0):
            raise RuntimeError(
                "ZIP 内 XLSX 未在 45 秒内进入行批进度"
            )
        token.pause()
        manager.request_pause()
        deadline = time.perf_counter() + 8.0
        while time.perf_counter() < deadline:
            if manager.is_safely_paused():
                safe = True
                break
            time.sleep(0.02)
        status = manager.pause_status()
        if not safe:
            raise RuntimeError(
                f"ZIP 内 XLSX 未在安全单位后暂停：{status}"
            )
        observation = _observe_real_paused_index(
            database,
            source,
            progress_events,
            seconds=5.0,
        )
        manager.request_resume()
        token.resume()
        resumed = True
        if not finished.wait(timeout=90.0):
            raise RuntimeError("ZIP 内 XLSX 继续后未完成")
    finally:
        if not resumed:
            manager.request_resume()
            token.resume()
        if not finished.wait(timeout=2.0):
            token.cancel(force=True)
            manager.force_terminate_processes()
        worker.join(timeout=5.0)
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
            ] = previous_delay

    acknowledgement = next(
        (
            item
            for item in list(status.get("acknowledgements") or [])
            if str(item.get("safe_unit_type") or "")
            == "sheet_row"
        ),
        {},
    )
    acknowledged_cursor = int(acknowledgement.get("cursor") or 0)
    with database.connect() as connection:
        final_cursor = int(
            connection.execute(
                """
                SELECT COALESCE(MAX(CAST(progress_cursor AS INTEGER)), 0)
                FROM parse_tasks
                """
            ).fetchone()[0]
        )
        member_success = bool(
            connection.execute(
                """
                SELECT COUNT(*) FROM files
                WHERE source_kind = 'zip_member'
                  AND internal_path = 'nested/safe-member.xlsx'
                  AND parse_status = 'success'
                """
            ).fetchone()[0]
        )
        search_token_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM content_blocks
                WHERE raw_text LIKE '%ZIP_SAFE_PAUSE_ROW_10000%'
                """
            ).fetchone()[0]
        )
        duplicate_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM (
                    SELECT file_id, block_index, COUNT(*) AS copies
                    FROM content_blocks
                    GROUP BY file_id, block_index
                    HAVING copies > 1
                )
                """
            ).fetchone()[0]
        )
    final_summary = outcomes[0] if outcomes else None
    failed_delta = int(
        len(database.failed_files(limit=10_000))
        - int(observation.get("failed_files_before") or 0)
    )
    resume_cursor_advanced = bool(
        acknowledged_cursor > 0
        and final_cursor > acknowledged_cursor
    )
    result = {
        "started": active_progress.is_set(),
        "safe_pause_confirmed": safe,
        "acknowledgement": acknowledgement,
        "observation_seconds": float(
            observation.get("observation_seconds") or 0.0
        ),
        "progress_delta": int(
            observation.get("progress_delta") or 0
        ),
        "database_write_delta": int(
            observation.get("database_write_delta") or 0
        ),
        "source_read_bytes_delta": int(
            observation.get("source_read_bytes_delta") or 0
        ),
        "paused_cpu_average": float(
            observation.get("paused_cpu_average") or 0.0
        ),
        "source_unchanged": bool(
            observation.get("source_unchanged")
        ),
        "resume_cursor_advanced": resume_cursor_advanced,
        "acknowledged_cursor": acknowledged_cursor,
        "final_cursor": final_cursor,
        "member_success": member_success,
        "search_token_blocks": search_token_blocks,
        "duplicate_blocks": duplicate_blocks,
        "failed_delta": failed_delta,
        "summary": (
            final_summary.to_dict()
            if final_summary is not None
            else {}
        ),
        "errors": errors,
    }
    result["passed"] = bool(
        not errors
        and final_summary is not None
        and final_summary.failed == 0
        and result["started"]
        and safe
        and result["observation_seconds"] >= 5.0
        and result["progress_delta"] == 0
        and result["database_write_delta"] == 0
        and result["source_read_bytes_delta"] == 0
        and result["paused_cpu_average"] <= 5.0
        and result["source_unchanged"]
        and resume_cursor_advanced
        and member_success
        and search_token_blocks == 1
        and duplicate_blocks == 0
        and failed_delta == 0
    )
    return result


def _validate_pdf_native_page_pause(base: Path) -> dict[str, object]:
    import fitz

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    source = root / "safe-pause-pages.pdf"
    page_count = 160
    document = fitz.open()
    try:
        for page_number in range(1, page_count + 1):
            page = document.new_page()
            page.insert_text(
                (72, 72),
                f"SAFE_PAUSE_PDF_PAGE_{page_number:03d}",
            )
        document.save(source)
    finally:
        document.close()

    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    pdf_batch_size = 32
    manager = IndexManager(
        database,
        AppSettings(
            enable_ocr=False,
            ocr_scanned_pdf=False,
            parser_workers=1,
            pdf_parser_workers=1,
            pdf_pending_tasks=1,
            pdf_page_batch_size=pdf_batch_size,
            index_write_batch_size=8,
        ),
        run_context={
            "execution_mode": "normal",
            "validation_kind": "safe_pause_real_pdf_native_page",
        },
    )
    token = CancelToken()
    progress_events: list[dict[str, object]] = []
    finished = threading.Event()
    outcomes: list[object] = []
    errors: list[str] = []

    def progress(payload: dict[str, object]) -> None:
        progress_events.append(dict(payload))

    def run_index() -> None:
        try:
            outcomes.append(manager.index_root(root_id, token, progress))
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")
        finally:
            finished.set()

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
    )
    os.environ["LFTS_VALIDATION_SAFE_POINT_DELAY_MS"] = "60"
    worker = threading.Thread(
        target=run_index,
        name="lfts-safe-pause-real-pdf-native-page",
    )
    worker.start()
    started = False
    safe = False
    resumed = False
    status: dict[str, object] = {}
    observation: dict[str, object] = {}
    completed_before_pause = 0
    pending_pages_while_paused = 0
    try:
        deadline = time.monotonic() + 45.0
        while time.monotonic() < deadline:
            with database.connect() as connection:
                completed_before_pause = int(
                    connection.execute(
                        """
                        SELECT COUNT(*) FROM parse_tasks
                        WHERE task_type = 'pdf_native_page'
                          AND status = 'complete'
                        """
                    ).fetchone()[0]
                )
                running_pages = int(
                    connection.execute(
                        """
                        SELECT COUNT(*) FROM parse_tasks
                        WHERE task_type = 'pdf_native_page'
                          AND status = 'running'
                        """
                    ).fetchone()[0]
                )
                pending_pages = int(
                    connection.execute(
                        """
                        SELECT COUNT(*) FROM parse_tasks
                        WHERE task_type = 'pdf_native_page'
                          AND status IN ('queued', 'running')
                        """
                    ).fetchone()[0]
                )
            if (
                completed_before_pause >= pdf_batch_size
                and running_pages > 0
                and pending_pages > 0
            ):
                started = True
                break
            time.sleep(0.02)
        if not started:
            raise RuntimeError(
                "真实 PDF 未在 45 秒内进入可观测的页任务进度"
            )
        token.pause()
        manager.request_pause()
        pause_deadline = time.monotonic() + 8.0
        while time.monotonic() < pause_deadline:
            if manager.is_safely_paused():
                safe = True
                break
            time.sleep(0.02)
        status = manager.pause_status()
        if not safe:
            raise RuntimeError(
                f"真实 PDF 页任务未在 8 秒内安全暂停：{status}"
            )
        with database.connect() as connection:
            pending_pages_while_paused = int(
                connection.execute(
                    """
                    SELECT COUNT(*) FROM parse_tasks
                    WHERE task_type = 'pdf_native_page'
                      AND status IN ('queued', 'running', 'paused')
                    """
                ).fetchone()[0]
            )
        observation = _observe_real_paused_index(
            database,
            source,
            progress_events,
            seconds=5.0,
        )
        manager.request_resume()
        token.resume()
        resumed = True
        if not finished.wait(timeout=90.0):
            raise RuntimeError("真实 PDF 页任务继续后未在 90 秒内完成")
    finally:
        if not resumed:
            manager.request_resume()
            token.resume()
        if not finished.wait(timeout=2.0):
            token.cancel(force=True)
            manager.force_terminate_processes()
        worker.join(timeout=5.0)
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
            ] = previous_delay

    final_summary = outcomes[0] if outcomes else None
    with database.connect() as connection:
        completed_after_resume = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'pdf_native_page'
                  AND status = 'complete'
                  AND confirmed_at IS NOT NULL
                """
            ).fetchone()[0]
        )
        merge_tasks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'document_merge'
                  AND status = 'complete'
                  AND confirmed_at IS NOT NULL
                """
            ).fetchone()[0]
        )
        ordered_page_numbers = [
            int(row[0])
            for row in connection.execute(
                """
                SELECT page_number FROM content_blocks
                WHERE block_type = 'pdf_page'
                ORDER BY page_number, block_index, id
                """
            )
        ]
        duplicate_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM (
                    SELECT file_id, block_index, COUNT(*) AS copies
                    FROM content_blocks
                    GROUP BY file_id, block_index
                    HAVING copies > 1
                )
                """
            ).fetchone()[0]
        )
    acknowledgements = [
        dict(item)
        for item in list(status.get("acknowledgements") or [])
        if str(item.get("safe_unit_type") or "")
        == "pdf_native_page"
    ]
    failed_delta = int(
        len(database.failed_files(limit=10_000))
        - int(observation.get("failed_files_before") or 0)
    )
    resume_pages_advanced = bool(
        completed_after_resume > completed_before_pause
    )
    passed = bool(
        not errors
        and final_summary is not None
        and final_summary.failed == 0
        and started
        and safe
        and acknowledgements
        and float(observation.get("observation_seconds") or 0.0) >= 5.0
        and int(observation.get("progress_delta") or 0) == 0
        and int(observation.get("database_write_delta") or 0) == 0
        and int(observation.get("source_read_bytes_delta") or 0) == 0
        and float(observation.get("paused_cpu_average") or 0.0) <= 5.0
        and bool(observation.get("source_unchanged"))
        and pending_pages_while_paused > 0
        and resume_pages_advanced
        and completed_after_resume == page_count
        and merge_tasks == 1
        and ordered_page_numbers == list(range(1, page_count + 1))
        and duplicate_blocks == 0
        and failed_delta == 0
    )
    return {
        "passed": passed,
        "started": started,
        "safe_pause_confirmed": safe,
        "status": status,
        "acknowledgements": acknowledgements,
        "observation_seconds": float(
            observation.get("observation_seconds") or 0.0
        ),
        "progress_delta": int(observation.get("progress_delta") or 0),
        "database_write_delta": int(
            observation.get("database_write_delta") or 0
        ),
        "source_read_bytes_delta": int(
            observation.get("source_read_bytes_delta") or 0
        ),
        "paused_cpu_average": float(
            observation.get("paused_cpu_average") or 0.0
        ),
        "source_unchanged": bool(observation.get("source_unchanged")),
        "completed_pages_before_pause": completed_before_pause,
        "pending_pages_while_paused": pending_pages_while_paused,
        "completed_pages_after_resume": completed_after_resume,
        "resume_pages_advanced": resume_pages_advanced,
        "page_tasks": page_count,
        "merge_tasks": merge_tasks,
        "ordered_pages": len(ordered_page_numbers),
        "duplicate_blocks": duplicate_blocks,
        "failed_delta": failed_delta,
        "summary": (
            final_summary.to_dict()
            if final_summary is not None
            else {}
        ),
        "errors": errors,
    }


def _observe_real_paused_index(
    database: DatabaseManager,
    source: Path,
    progress_events: list[dict[str, object]],
    *,
    seconds: float,
    read_process_ids: list[int] | None = None,
) -> dict[str, object]:
    import psutil

    # A final heartbeat may already be queued when the coordinator confirms
    # the safe point.  Establish the observation baseline only after that
    # callback stream has been quiet for one heartbeat interval.
    settle_deadline = time.perf_counter() + 2.0
    stable_since = time.perf_counter()
    observed_count = len(progress_events)
    while time.perf_counter() < settle_deadline:
        time.sleep(0.05)
        current_count = len(progress_events)
        if current_count != observed_count:
            observed_count = current_count
            stable_since = time.perf_counter()
        if time.perf_counter() - stable_since >= 0.8:
            break
    observation_started = time.perf_counter()
    progress_before = len(progress_events)
    database_before = _paused_database_signature(database)
    source_before = _source_identity(source)
    processes = (
        _processes_for_ids(read_process_ids)
        if read_process_ids is not None
        else _process_tree()
    )
    io_before = _process_read_bytes(processes)
    for process in processes:
        try:
            process.cpu_percent(interval=None)
        except (psutil.Error, OSError):
            continue
    time.sleep(max(5.0, float(seconds)))
    elapsed = time.perf_counter() - observation_started
    cpu_total = 0.0
    cpu_samples = 0
    for process in processes:
        try:
            cpu_total += float(process.cpu_percent(interval=None))
            cpu_samples += 1
        except (psutil.Error, OSError):
            continue
    io_after = _process_read_bytes(processes)
    database_after = _paused_database_signature(database)
    source_after = _source_identity(source)
    failed_before = int(database_before[-1])
    return {
        "observation_seconds": round(elapsed, 6),
        "progress_delta": len(progress_events) - progress_before,
        "database_write_delta": int(
            database_after != database_before
        ),
        "source_read_bytes_delta": max(0, io_after - io_before),
        "paused_cpu_average": round(
            cpu_total / max(1, cpu_samples),
            6,
        ),
        "source_unchanged": source_before == source_after,
        "failed_files_before": failed_before,
    }


def _process_tree() -> list[object]:
    import psutil

    parent = psutil.Process()
    return [parent, *parent.children(recursive=True)]


def _processes_for_ids(process_ids: list[int]) -> list[object]:
    import psutil

    processes: dict[int, object] = {}
    for process_id in process_ids:
        try:
            process = psutil.Process(int(process_id))
            processes[process.pid] = process
            for child in process.children(recursive=True):
                processes[child.pid] = child
        except (psutil.Error, OSError, ValueError):
            continue
    return list(processes.values())


def _process_read_bytes(processes: list[object]) -> int:
    import psutil

    total = 0
    for process in processes:
        try:
            total += int(process.io_counters().read_bytes)
        except (psutil.Error, OSError, AttributeError):
            continue
    return total


def _paused_database_signature(
    database: DatabaseManager,
) -> tuple[object, ...]:
    with database.connect() as connection:
        blocks = connection.execute(
            """
            SELECT COUNT(*), COALESCE(SUM(LENGTH(raw_text)), 0)
            FROM content_blocks
            """
        ).fetchone()
        task_states = tuple(
            tuple(row)
            for row in connection.execute(
                """
                SELECT status, COUNT(*)
                FROM parse_tasks
                GROUP BY status
                ORDER BY status
                """
            )
        )
        failed = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM files
                WHERE parse_status LIKE 'failed%'
                """
            ).fetchone()[0]
        )
    return (
        int(blocks[0]),
        int(blocks[1]),
        task_states,
        failed,
    )


def _source_identity(path: Path) -> tuple[int, int]:
    stat = path.stat()
    return int(stat.st_size), int(stat.st_mtime_ns)


def validate_paused_mode_switch(base: Path) -> dict[str, object]:
    from openpyxl import Workbook

    base.mkdir(parents=True, exist_ok=True)
    root = base / "mode-switch-xlsx"
    root.mkdir()
    source = root / "mode-switch-rows.xlsx"
    workbook = Workbook(write_only=True)
    sheet = workbook.create_sheet("ModeSwitch")
    for row_number in range(1, 20_001):
        sheet.append(
            [
                row_number,
                f"SAFE_MODE_SWITCH_ROW_{row_number:05d}",
                "双向模式切换不得改变正文",
            ]
        )
    workbook.save(source)

    normal = AppSettings(
        enable_ocr=False,
        enable_parse_cache=False,
        parser_workers=1,
        process_parser_workers=1,
        process_pending_tasks=1,
        pdf_parser_workers=1,
        pdf_pending_tasks=1,
        ocr_workers=1,
        slow_file_workers=1,
        xlsx_sheet_workers=1,
        index_write_batch_size=8,
    )
    performance = AppSettings.from_dict(normal.to_dict())
    performance.parser_workers = 2
    performance.process_parser_workers = 2
    performance.process_pending_tasks = 2
    performance.pdf_parser_workers = 2
    performance.ocr_workers = 2
    performance.xlsx_sheet_workers = 2
    performance.index_cpu_token_budget = 6
    performance.index_memory_budget_mb = 3072

    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    manager = IndexManager(
        database,
        normal,
        run_context={"execution_mode": "normal"},
    )
    token = CancelToken()
    progress_cursors: list[int] = []
    outcomes: list[object] = []
    errors: list[str] = []
    finished = threading.Event()

    def progress(payload: dict[str, object]) -> None:
        if str(payload.get("active_phase") or "") != "sheet_row":
            return
        cursor = int(payload.get("active_completed_units") or 0)
        if cursor > 0:
            progress_cursors.append(cursor)

    def run_index() -> None:
        try:
            outcomes.append(
                manager.index_root(root_id, token, progress)
            )
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")
        finally:
            finished.set()

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
    )
    os.environ["LFTS_VALIDATION_SAFE_POINT_DELAY_MS"] = "80"
    worker = threading.Thread(
        target=run_index,
        name="lfts-mode-switch-real-xlsx",
    )
    worker.start()
    switches: list[str] = []
    paused_after_switches: list[bool] = []
    cursors: list[tuple[int, int]] = []
    resumed = False
    try:
        first_before = _wait_for_xlsx_cursor(
            progress_cursors,
            greater_than=0,
            timeout=45.0,
        )
        first_ack = _pause_real_index(manager, token)
        applied = manager.apply_settings_while_paused(
            performance,
            execution_mode="performance",
            effective_profile={
                "mode": "performance",
                "process_parser_workers": 2,
                "ocr_workers": 2,
            },
        )
        if not applied:
            raise RuntimeError(
                "真实 XLSX 普通模式切性能模式失败"
            )
        switches.append("normal_to_performance")
        paused_after_switches.append(
            _assert_switch_did_not_resume(
                manager,
                database,
                progress_cursors,
            )
        )
        manager.request_resume()
        token.resume()
        resumed = True
        first_after = _wait_for_xlsx_cursor(
            progress_cursors,
            greater_than=max(first_before, first_ack),
            timeout=45.0,
        )
        cursors.append((max(first_before, first_ack), first_after))

        resumed = False
        second_ack = _pause_real_index(manager, token)
        applied = manager.apply_settings_while_paused(
            normal,
            execution_mode="normal",
            effective_profile={
                "mode": "normal",
                "process_parser_workers": 1,
                "ocr_workers": 1,
            },
        )
        if not applied:
            raise RuntimeError(
                "真实 XLSX 性能模式切普通模式失败"
            )
        switches.append("performance_to_normal")
        paused_after_switches.append(
            _assert_switch_did_not_resume(
                manager,
                database,
                progress_cursors,
            )
        )
        manager.request_resume()
        token.resume()
        resumed = True
        second_after = _wait_for_xlsx_cursor(
            progress_cursors,
            greater_than=second_ack,
            timeout=45.0,
        )
        cursors.append((second_ack, second_after))
        if not finished.wait(timeout=90.0):
            raise RuntimeError(
                "真实 XLSX 双向切换后未完成"
            )
    finally:
        if not resumed:
            manager.request_resume()
            token.resume()
        if not finished.wait(timeout=2.0):
            token.cancel(force=True)
            manager.force_terminate_processes()
        worker.join(timeout=5.0)
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
            ] = previous_delay

    candidate_snapshot = _content_and_search_snapshot(
        database,
        "SAFE_MODE_SWITCH_ROW_01000",
    )
    control_database = DatabaseManager(base / "control.db")
    control_database.initialize()
    control_root_id = control_database.add_root(root)
    control_summary = IndexManager(
        control_database,
        normal,
        run_context={"execution_mode": "normal"},
    ).index_root(control_root_id)
    control_snapshot = _content_and_search_snapshot(
        control_database,
        "SAFE_MODE_SWITCH_ROW_01000",
    )
    duplicate_blocks = int(
        candidate_snapshot["duplicate_blocks"]
    )
    rollback_injection_passed = _validate_mode_switch_rollback(
        base / "rollback"
    )
    format_results = {
        "image_ocr": _validate_image_ocr_mode_switch(
            base / "image-ocr"
        ),
        "pdf_ocr_page": _validate_pdf_ocr_mode_switch(
            base / "pdf-ocr"
        ),
        "zip_member": _validate_zip_member_mode_switch(
            base / "zip-member"
        ),
        "legacy_office": _validate_legacy_office_mode_switch(
            base / "legacy-office"
        ),
    }
    cursor_advanced = bool(
        len(cursors) == 2
        and all(after > before for before, after in cursors)
    )
    final_summary = outcomes[0] if outcomes else None
    passed = bool(
        not errors
        and final_summary is not None
        and final_summary.failed == 0
        and control_summary.failed == 0
        and switches
        == [
            "normal_to_performance",
            "performance_to_normal",
        ]
        and all(paused_after_switches)
        and cursor_advanced
        and duplicate_blocks == 0
        and candidate_snapshot["content_digest"]
        == control_snapshot["content_digest"]
        and candidate_snapshot["search_hits"]
        == control_snapshot["search_hits"]
        and rollback_injection_passed
        and all(
            bool(item["passed"])
            for item in format_results.values()
        )
    )
    return {
        "passed": passed,
        "mechanism_only": False,
        "format": "xlsx",
        "switches": switches,
        "remained_paused_after_each_switch": bool(
            paused_after_switches
            and all(paused_after_switches)
        ),
        "cursor_pairs": cursors,
        "cursor_advanced_after_each_resume": cursor_advanced,
        "duplicate_blocks": duplicate_blocks,
        "content_digest_matches_control": (
            candidate_snapshot["content_digest"]
            == control_snapshot["content_digest"]
        ),
        "search_hits_match_control": (
            candidate_snapshot["search_hits"]
            == control_snapshot["search_hits"]
        ),
        "rollback_injection_passed": rollback_injection_passed,
        "formats": format_results,
        "candidate": candidate_snapshot,
        "control": control_snapshot,
        "errors": errors,
        "real_long_format_gate_required": False,
        "remaining_format_gates": [],
    }


def _performance_mode_settings(normal: AppSettings) -> AppSettings:
    performance = AppSettings.from_dict(normal.to_dict())
    performance.parser_workers = max(2, int(normal.parser_workers))
    performance.process_parser_workers = max(
        2,
        int(normal.process_parser_workers),
    )
    performance.process_pending_tasks = max(
        2,
        int(normal.process_pending_tasks),
    )
    performance.pdf_parser_workers = max(
        2,
        int(normal.pdf_parser_workers),
    )
    performance.ocr_workers = max(2, int(normal.ocr_workers))
    performance.xlsx_sheet_workers = max(
        2,
        int(normal.xlsx_sheet_workers),
    )
    performance.index_cpu_token_budget = max(
        6,
        int(normal.index_cpu_token_budget),
    )
    performance.index_memory_budget_mb = max(
        3072,
        int(normal.index_memory_budget_mb),
    )
    return performance


def _mode_switch_ack_cursor(status: dict[str, object]) -> int:
    def cursor_value(item: object) -> int:
        if not isinstance(item, dict):
            return 0
        cursor = item.get("cursor")
        if cursor is not None and str(cursor) != "":
            return int(cursor)
        return int(item.get("completed_units") or 0)

    return max(
        (
            cursor_value(item)
            for item in list(status.get("acknowledgements") or [])
        ),
        default=0,
    )


def _wait_for_safe_pause(
    manager: IndexManager,
    token: CancelToken,
    *,
    timeout: float,
) -> dict[str, object]:
    token.pause()
    manager.request_pause()
    deadline = time.monotonic() + float(timeout)
    while time.monotonic() < deadline:
        if manager.is_safely_paused():
            status = manager.pause_status()
            if status.get("state") == "paused":
                return status
        time.sleep(0.02)
    raise RuntimeError(
        f"Real task did not reach a safe pause: {manager.pause_status()}"
    )


def _mode_switch_stayed_paused(
    manager: IndexManager,
    database: DatabaseManager,
    progress_events: list[dict[str, object]],
) -> bool:
    progress_before = len(progress_events)
    database_before = _paused_database_signature(database)
    time.sleep(0.5)
    return bool(
        manager.pause_status().get("state") == "paused"
        and len(progress_events) == progress_before
        and _paused_database_signature(database) == database_before
    )


def _run_real_mode_switch_task(
    base: Path,
    root: Path,
    source: Path,
    normal: AppSettings,
    *,
    progress_predicate: Callable[[dict[str, object]], bool],
    thread_name: str,
    start_timeout: float,
    resume_timeout: float,
    safe_pause_timeout: float = 60.0,
    delay_ms: int = 150,
    require_external_process: bool = False,
    checkpoint_advance_ready: (
        Callable[[DatabaseManager, int], bool] | None
    ) = None,
    reverse_switch: bool = True,
) -> tuple[dict[str, object], DatabaseManager]:
    from local_full_text_search.parsers.legacy_office_parser import (
        registered_office_processes_alive,
    )

    database = DatabaseManager(base / "index.db")
    database.initialize()
    root_id = database.add_root(root)
    performance = _performance_mode_settings(normal)
    manager = IndexManager(
        database,
        normal,
        run_context={
            "execution_mode": "normal",
            "validation_kind": thread_name,
        },
    )
    token = CancelToken()
    progress_events: list[dict[str, object]] = []
    relevant_events: list[dict[str, object]] = []
    active_progress = threading.Event()
    finished = threading.Event()
    outcomes: list[object] = []
    errors: list[str] = []

    def progress(payload: dict[str, object]) -> None:
        event = dict(payload)
        progress_events.append(event)
        if progress_predicate(event):
            relevant_events.append(event)
            active_progress.set()

    def run_index() -> None:
        try:
            outcomes.append(manager.index_root(root_id, token, progress))
        except BaseException as exc:
            errors.append(f"{type(exc).__name__}: {exc}")
        finally:
            finished.set()

    previous_delay = os.environ.get(
        "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
    )
    os.environ["LFTS_VALIDATION_SAFE_POINT_DELAY_MS"] = str(
        max(0, int(delay_ms))
    )
    worker = threading.Thread(target=run_index, name=thread_name)
    worker.start()
    switches: list[str] = []
    paused_after_switches: list[bool] = []
    statuses: list[dict[str, object]] = []
    cursors: list[int] = []
    observed_office_pids: list[int] = []
    external_process_seen = False
    resumed = False
    spool_dir: Path | None = None
    try:
        if not active_progress.wait(timeout=float(start_timeout)):
            raise RuntimeError("Real format did not enter semantic progress")
        if require_external_process:
            deadline = time.monotonic() + 15.0
            while time.monotonic() < deadline:
                spool_dir = manager._pause_spool_dir
                if spool_dir is not None:
                    observed_office_pids = (
                        registered_office_processes_alive(spool_dir)
                    )
                    if observed_office_pids:
                        external_process_seen = True
                        break
                if finished.is_set():
                    break
                time.sleep(0.02)
            if not external_process_seen:
                raise RuntimeError(
                    "Real legacy conversion did not register an Office process"
                )

        first_status = _wait_for_safe_pause(
            manager,
            token,
            timeout=safe_pause_timeout,
        )
        statuses.append(first_status)
        cursors.append(_mode_switch_ack_cursor(first_status))
        if not manager.apply_settings_while_paused(
            performance,
            execution_mode="performance",
            effective_profile={
                "mode": "performance",
                "process_parser_workers": (
                    performance.process_parser_workers
                ),
                "ocr_workers": performance.ocr_workers,
            },
        ):
            raise RuntimeError("normal to performance mode switch failed")
        switches.append("normal_to_performance")
        paused_after_switches.append(
            _mode_switch_stayed_paused(
                manager,
                database,
                progress_events,
            )
        )

        if not reverse_switch:
            manager.request_resume()
            token.resume()
            resumed = True
            if not finished.wait(timeout=float(resume_timeout)):
                raise RuntimeError(
                    "Real format did not finish after the mode switch"
                )
        else:
            relevant_before_resume = len(relevant_events)
            manager.request_resume()
            token.resume()
            resumed = True
            resume_progress_not_before = time.monotonic() + 0.5
            progress_deadline = time.monotonic() + float(resume_timeout)
            while time.monotonic() < progress_deadline:
                if (
                    time.monotonic() >= resume_progress_not_before
                    and len(relevant_events) > relevant_before_resume
                    and (
                        checkpoint_advance_ready is None
                        or checkpoint_advance_ready(database, cursors[0])
                    )
                ):
                    break
                if finished.is_set():
                    raise RuntimeError(
                        "Real format completed before the reverse mode switch: "
                        f"first_pause={first_status!r}"
                    )
                time.sleep(0.02)
            else:
                raise RuntimeError(
                    "Real format did not make progress after the first switch"
                )

            resumed = False
            second_status = _wait_for_safe_pause(
                manager,
                token,
                timeout=safe_pause_timeout,
            )
            advancement_deadline = time.monotonic() + float(resume_timeout)
            while (
                _mode_switch_ack_cursor(second_status) <= cursors[0]
                and time.monotonic() < advancement_deadline
            ):
                relevant_before_retry = len(relevant_events)
                manager.request_resume()
                token.resume()
                resumed = True
                retry_progress_not_before = time.monotonic() + 0.5
                while time.monotonic() < advancement_deadline:
                    if (
                        time.monotonic() >= retry_progress_not_before
                        and len(relevant_events) > relevant_before_retry
                        and (
                            checkpoint_advance_ready is None
                            or checkpoint_advance_ready(
                                database,
                                cursors[0],
                            )
                        )
                    ):
                        break
                    if finished.is_set():
                        raise RuntimeError(
                            "Real format completed before checkpoint advancement: "
                            f"first_pause={first_status!r}, "
                            f"second_pause={second_status!r}"
                        )
                    time.sleep(0.02)
                resumed = False
                second_status = _wait_for_safe_pause(
                    manager,
                    token,
                    timeout=safe_pause_timeout,
                )
            if _mode_switch_ack_cursor(second_status) <= cursors[0]:
                raise RuntimeError(
                    "Real format did not advance its confirmed cursor"
                )
            statuses.append(second_status)
            cursors.append(_mode_switch_ack_cursor(second_status))
            if not manager.apply_settings_while_paused(
                normal,
                execution_mode="normal",
                effective_profile={
                    "mode": "normal",
                    "process_parser_workers": normal.process_parser_workers,
                    "ocr_workers": normal.ocr_workers,
                },
            ):
                raise RuntimeError("performance to normal mode switch failed")
            switches.append("performance_to_normal")
            paused_after_switches.append(
                _mode_switch_stayed_paused(
                    manager,
                    database,
                    progress_events,
                )
            )
            manager.request_resume()
            token.resume()
            resumed = True
            if not finished.wait(timeout=float(resume_timeout)):
                raise RuntimeError(
                    "Real format did not finish after the reverse switch"
                )
    finally:
        if not resumed:
            manager.request_resume()
            token.resume()
        if not finished.wait(timeout=2.0):
            token.cancel(force=True)
            manager.force_terminate_processes()
        worker.join(timeout=5.0)
        if spool_dir is None:
            spool_dir = manager._pause_spool_dir
        if previous_delay is None:
            os.environ.pop(
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS",
                None,
            )
        else:
            os.environ[
                "LFTS_VALIDATION_SAFE_POINT_DELAY_MS"
            ] = previous_delay

    with database.connect() as connection:
        final_cursor = int(
            connection.execute(
                """
                SELECT COALESCE(MAX(CAST(progress_cursor AS INTEGER)), 0)
                FROM parse_tasks
                """
            ).fetchone()[0]
        )
        duplicate_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM (
                    SELECT file_id, block_index, COUNT(*) AS copies
                    FROM content_blocks
                    GROUP BY file_id, block_index
                    HAVING copies > 1
                )
                """
            ).fetchone()[0]
        )
    cursor_pairs = (
        [(cursors[0], cursors[1]), (cursors[1], final_cursor)]
        if len(cursors) == 2
        else ([(cursors[0], final_cursor)] if len(cursors) == 1 else [])
    )
    cursor_advanced = bool(
        len(cursor_pairs) == (2 if reverse_switch else 1)
        and all(after > before for before, after in cursor_pairs)
    )
    final_summary = outcomes[0] if outcomes else None
    office_pids_after = (
        registered_office_processes_alive(spool_dir)
        if spool_dir is not None
        else []
    )
    result = {
        "switches": switches,
        "statuses": statuses,
        "remained_paused_after_each_switch": bool(
            paused_after_switches and all(paused_after_switches)
        ),
        "cursor_pairs": cursor_pairs,
        "cursor_advanced_after_each_resume": cursor_advanced,
        "duplicate_blocks": duplicate_blocks,
        "external_process_seen": external_process_seen,
        "observed_office_pids": observed_office_pids,
        "office_pids_after": office_pids_after,
        "summary": (
            final_summary.to_dict()
            if final_summary is not None
            else {}
        ),
        "errors": errors,
    }
    result["passed"] = bool(
        not errors
        and final_summary is not None
        and final_summary.failed == 0
        and switches
        == (
            [
                "normal_to_performance",
                "performance_to_normal",
            ]
            if reverse_switch
            else ["normal_to_performance"]
        )
        and result["remained_paused_after_each_switch"]
        and cursor_advanced
        and duplicate_blocks == 0
        and (
            not require_external_process
            or (external_process_seen and not office_pids_after)
        )
    )
    return result, database


def _mode_switch_control_comparison(
    base: Path,
    root: Path,
    normal: AppSettings,
    candidate_database: DatabaseManager,
    query: str,
) -> tuple[dict[str, object], dict[str, object], object]:
    candidate = _content_and_search_snapshot(
        candidate_database,
        query,
    )
    control_database = DatabaseManager(base / "control.db")
    control_database.initialize()
    control_root_id = control_database.add_root(root)
    control_summary = IndexManager(
        control_database,
        normal,
        run_context={"execution_mode": "normal"},
    ).index_root(control_root_id)
    control = _content_and_search_snapshot(
        control_database,
        query,
    )
    return candidate, control, control_summary


def _validate_image_ocr_mode_switch(
    base: Path,
) -> dict[str, object]:
    from PIL import Image, ImageDraw, ImageFont

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    source = root / "mode-switch-image.png"
    image = Image.new("RGB", (1900, 1700), "white")
    draw = ImageDraw.Draw(image)
    validation_nonce = time.time_ns()
    font_path = (
        Path(os.environ.get("WINDIR", "C:/Windows"))
        / "Fonts"
        / "arial.ttf"
    )
    font = (
        ImageFont.truetype(str(font_path), 40)
        if font_path.is_file()
        else ImageFont.load_default()
    )
    for line_number in range(1, 31):
        draw.text(
            (55, 25 + (line_number - 1) * 54),
            (
                f"IMAGE MODE SWITCH {validation_nonce} LINE {line_number:02d} "
                f"VALUE {line_number * 173}"
            ),
            fill="black",
            font=font,
        )
    image.save(source)
    normal = AppSettings(
        enable_ocr=True,
        ocr_images=True,
        ocr_scanned_pdf=False,
        min_ocr_image_pixels=0,
        enable_parse_cache=False,
        parser_workers=1,
        process_parser_workers=1,
        process_pending_tasks=1,
        pdf_parser_workers=1,
        ocr_workers=1,
        ocr_pending_tasks=1,
        slow_file_workers=1,
        index_write_batch_size=8,
    )
    result, database = _run_real_mode_switch_task(
        base,
        root,
        source,
        normal,
        progress_predicate=lambda event: (
            str(event.get("current_file") or "")
            .lower()
            .endswith(".png")
            and str(event.get("active_phase") or "").startswith("ocr_")
        ),
        thread_name="lfts-mode-switch-image-ocr",
        start_timeout=120.0,
        resume_timeout=300.0,
        delay_ms=250,
    )
    candidate, control, control_summary = (
        _mode_switch_control_comparison(
            base,
            root,
            normal,
            database,
            "IMAGEMODESWITCH",
        )
    )
    with database.connect() as connection:
        row = connection.execute(
            """
            SELECT COALESCE(SUM(LENGTH(raw_text)), 0), extra_json
            FROM content_blocks
            WHERE block_type = 'image_ocr'
            """
        ).fetchone()
        source_success = bool(
            connection.execute(
                """
                SELECT COUNT(*) FROM files
                WHERE source_kind = 'file'
                  AND is_deleted = 0
                  AND parse_status = 'success'
                """
            ).fetchone()[0]
        )
    extra = json.loads(str(row[1] or "{}")) if row else {}
    result.update(
        {
            "source_success": source_success,
            "ocr_text_chars": int(row[0] or 0) if row else 0,
            "engine": str(extra.get("engine") or ""),
            "candidate": candidate,
            "control": control,
            "content_digest_matches_control": (
                candidate["content_digest"] == control["content_digest"]
            ),
            "search_hits_match_control": (
                candidate["search_hits"] == control["search_hits"]
            ),
            "control_summary": control_summary.to_dict(),
        }
    )
    result["passed"] = bool(
        result["passed"]
        and control_summary.failed == 0
        and source_success
        and result["ocr_text_chars"] > 0
        and result["engine"] == "PaddleOCR"
        and candidate["search_hits"] > 0
        and result["content_digest_matches_control"]
        and result["search_hits_match_control"]
    )
    return result


def _validate_pdf_ocr_mode_switch(
    base: Path,
) -> dict[str, object]:
    import fitz
    from PIL import Image, ImageDraw, ImageFont

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    image = Image.new("RGB", (1500, 1100), "white")
    draw = ImageDraw.Draw(image)
    validation_nonce = time.time_ns()
    font_path = (
        Path(os.environ.get("WINDIR", "C:/Windows"))
        / "Fonts"
        / "arial.ttf"
    )
    font = (
        ImageFont.truetype(str(font_path), 40)
        if font_path.is_file()
        else ImageFont.load_default()
    )
    for line_number in range(1, 19):
        draw.text(
            (55, 25 + (line_number - 1) * 56),
            (
                f"PDF OCR MODE SWITCH {validation_nonce} LINE {line_number:02d} "
                f"VALUE {line_number * 197}"
            ),
            fill="black",
            font=font,
        )
    source = root / "mode-switch-scanned.pdf"
    page_count = 8
    document = fitz.open()
    try:
        for page_number in range(1, page_count + 1):
            page_image = image.copy()
            page_draw = ImageDraw.Draw(page_image)
            page_draw.text(
                (1050, 1000),
                f"UNIQUE PAGE {validation_nonce} {page_number:02d}",
                fill="black",
                font=font,
            )
            page_image_path = base / f"mode-switch-page-{page_number}.png"
            page_image.save(page_image_path)
            page = document.new_page(width=750, height=550)
            page.insert_image(page.rect, filename=str(page_image_path))
        document.save(source)
    finally:
        document.close()
    normal = AppSettings(
        enable_ocr=True,
        ocr_images=False,
        ocr_scanned_pdf=True,
        enable_parse_cache=False,
        parser_workers=1,
        process_parser_workers=1,
        process_pending_tasks=1,
        pdf_parser_workers=1,
        pdf_pending_tasks=1,
        ocr_workers=1,
        ocr_pending_tasks=1,
        ocr_microbatch_parent_jobs=1,
        index_write_batch_size=8,
    )
    result, database = _run_real_mode_switch_task(
        base,
        root,
        source,
        normal,
        progress_predicate=lambda event: (
            str(event.get("current_file") or "")
            .lower()
            .endswith(".pdf")
            and (
                str(event.get("active_phase") or "").startswith(
                    "pdf_ocr_"
                )
                or str(event.get("active_phase") or "").startswith(
                    "pdf_region_"
                )
            )
        ),
        thread_name="lfts-mode-switch-pdf-ocr",
        start_timeout=120.0,
        resume_timeout=360.0,
        delay_ms=1_000,
        checkpoint_advance_ready=_pdf_ocr_page_checkpoint_advanced,
    )
    candidate, control, control_summary = (
        _mode_switch_control_comparison(
            base,
            root,
            normal,
            database,
            "PDF OCR MODE SWITCH",
        )
    )
    with database.connect() as connection:
        completed_pages = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'pdf_ocr_page'
                  AND status = 'complete'
                  AND confirmed_at IS NOT NULL
                """
            ).fetchone()[0]
        )
        ordered_page_numbers = [
            int(row[0])
            for row in connection.execute(
                """
                SELECT page_number FROM content_blocks
                WHERE block_type = 'pdf_page_ocr'
                ORDER BY page_number
                """
            )
        ]
        ocr_text_chars = int(
            connection.execute(
                """
                SELECT COALESCE(SUM(LENGTH(raw_text)), 0)
                FROM content_blocks
                WHERE block_type = 'pdf_page_ocr'
                """
            ).fetchone()[0]
        )
    result.update(
        {
            "completed_pages": completed_pages,
            "ordered_pages": len(ordered_page_numbers),
            "ocr_text_chars": ocr_text_chars,
            "candidate": candidate,
            "control": control,
            "content_digest_matches_control": (
                candidate["content_digest"] == control["content_digest"]
            ),
            "search_hits_match_control": (
                candidate["search_hits"] == control["search_hits"]
            ),
            "control_summary": control_summary.to_dict(),
        }
    )
    result["passed"] = bool(
        result["passed"]
        and control_summary.failed == 0
        and completed_pages == page_count
        and ordered_page_numbers == list(range(1, page_count + 1))
        and ocr_text_chars > 0
        and candidate["search_hits"] > 0
        and result["content_digest_matches_control"]
        and result["search_hits_match_control"]
    )
    return result


def _pdf_ocr_page_checkpoint_advanced(
    database: DatabaseManager,
    acknowledged_page: int,
) -> bool:
    with database.connect() as connection:
        completed = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM parse_tasks
                WHERE task_type = 'pdf_ocr_page'
                  AND status = 'complete'
                  AND confirmed_at IS NOT NULL
                """
            ).fetchone()[0]
        )
    return completed >= max(1, int(acknowledged_page))


def _validate_zip_member_mode_switch(
    base: Path,
) -> dict[str, object]:
    import zipfile

    from openpyxl import Workbook

    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    member_source = base / "mode-switch-member.xlsx"
    workbook = Workbook(write_only=True)
    sheet = workbook.create_sheet("ModeSwitch")
    for row_number in range(1, 12_001):
        sheet.append(
            [
                row_number,
                f"ZIP_MODE_SWITCH_ROW_{row_number:05d}",
                "ZIP member bidirectional mode switch",
            ]
        )
    workbook.save(member_source)
    source = root / "mode-switch-member.zip"
    with zipfile.ZipFile(
        source,
        "w",
        compression=zipfile.ZIP_DEFLATED,
    ) as archive:
        archive.write(
            member_source,
            arcname="nested/mode-switch-member.xlsx",
        )
    normal = AppSettings(
        enable_ocr=False,
        enable_parse_cache=False,
        parser_workers=1,
        process_parser_workers=1,
        process_pending_tasks=1,
        slow_file_workers=1,
        xlsx_sheet_workers=1,
        index_write_batch_size=8,
    )
    result, database = _run_real_mode_switch_task(
        base,
        root,
        source,
        normal,
        progress_predicate=lambda event: (
            str(event.get("current_file") or "")
            .lower()
            .endswith(".xlsx")
            and str(event.get("active_phase") or "") == "sheet_row"
            and int(event.get("active_completed_units") or 0) > 0
        ),
        thread_name="lfts-mode-switch-zip-member",
        start_timeout=60.0,
        resume_timeout=180.0,
        delay_ms=100,
    )
    candidate, control, control_summary = (
        _mode_switch_control_comparison(
            base,
            root,
            normal,
            database,
            "ZIP_MODE_SWITCH_ROW_10000",
        )
    )
    with database.connect() as connection:
        member_success = bool(
            connection.execute(
                """
                SELECT COUNT(*) FROM files
                WHERE source_kind = 'zip_member'
                  AND internal_path = 'nested/mode-switch-member.xlsx'
                  AND parse_status = 'success'
                """
            ).fetchone()[0]
        )
    result.update(
        {
            "member_success": member_success,
            "candidate": candidate,
            "control": control,
            "content_digest_matches_control": (
                candidate["content_digest"] == control["content_digest"]
            ),
            "search_hits_match_control": (
                candidate["search_hits"] == control["search_hits"]
            ),
            "control_summary": control_summary.to_dict(),
        }
    )
    result["passed"] = bool(
        result["passed"]
        and control_summary.failed == 0
        and member_success
        and candidate["search_hits"] == 1
        and result["content_digest_matches_control"]
        and result["search_hits_match_control"]
    )
    return result


def _validate_legacy_office_mode_switch(
    base: Path,
) -> dict[str, object]:
    base.mkdir(parents=True, exist_ok=True)
    root = base / "files"
    root.mkdir()
    source = root / "mode-switch-legacy.doc"
    _create_legacy_word_doc_isolated(source)
    normal = AppSettings(
        enable_ocr=False,
        enable_parse_cache=False,
        legacy_conversion_cache=True,
        fast_ooxml_enabled=True,
        parser_workers=1,
        process_parser_workers=1,
        process_pending_tasks=1,
        slow_file_workers=1,
        index_write_batch_size=8,
    )
    result, database = _run_real_mode_switch_task(
        base,
        root,
        source,
        normal,
        progress_predicate=lambda event: (
            str(event.get("current_file") or "")
            .lower()
            .endswith(".doc")
            and str(event.get("active_phase") or "")
            in {"legacy_office_open", "legacy_converted_block"}
        ),
        thread_name="lfts-mode-switch-legacy-office",
        start_timeout=90.0,
        resume_timeout=240.0,
        delay_ms=100,
        require_external_process=True,
        reverse_switch=False,
    )
    candidate, control, control_summary = (
        _mode_switch_control_comparison(
            base,
            root,
            normal,
            database,
            "LEGACY_SAFE_PAUSE_BLOCK_400",
        )
    )
    with database.connect() as connection:
        source_success = bool(
            connection.execute(
                """
                SELECT COUNT(*) FROM files
                WHERE source_kind = 'file'
                  AND is_deleted = 0
                  AND parse_status = 'success'
                """
            ).fetchone()[0]
        )
        converter_rows = [
            json.loads(str(row[0] or "{}"))
            for row in connection.execute(
                """
                SELECT extra_json FROM content_blocks
                WHERE raw_text <> ''
                """
            )
        ]
    result.update(
        {
            "source_success": source_success,
            "conversion_cache_reused_after_resume": bool(
                converter_rows
                and all(
                    str(item.get("legacy_converter") or "") == "cache"
                    for item in converter_rows
                )
            ),
            "candidate": candidate,
            "control": control,
            "content_digest_matches_control": (
                candidate["content_digest"] == control["content_digest"]
            ),
            "search_hits_match_control": (
                candidate["search_hits"] == control["search_hits"]
            ),
            "control_summary": control_summary.to_dict(),
        }
    )
    result["passed"] = bool(
        result["passed"]
        and control_summary.failed == 0
        and source_success
        and result["conversion_cache_reused_after_resume"]
        and candidate["search_hits"] == 1
        and result["content_digest_matches_control"]
        and result["search_hits_match_control"]
    )
    return result


def _wait_for_xlsx_cursor(
    progress_cursors: list[int],
    *,
    greater_than: int,
    timeout: float,
) -> int:
    deadline = time.monotonic() + float(timeout)
    while time.monotonic() < deadline:
        current = max(progress_cursors, default=0)
        if current > int(greater_than):
            return current
        time.sleep(0.02)
    raise RuntimeError(
        "真实 XLSX 行游标未在规定时间内继续前进"
    )


def _pause_real_index(
    manager: IndexManager,
    token: CancelToken,
) -> int:
    token.pause()
    manager.request_pause()
    deadline = time.monotonic() + 8.0
    while time.monotonic() < deadline:
        if manager.is_safely_paused():
            status = manager.pause_status()
            acknowledgements = list(
                status.get("acknowledgements") or []
            )
            cursor = max(
                (
                    int(
                        item.get("completed_units")
                        or item.get("cursor")
                        or 0
                    )
                    for item in acknowledgements
                    if str(item.get("safe_unit_type") or "")
                    == "sheet_row"
                ),
                default=0,
            )
            if cursor > 0:
                return cursor
        time.sleep(0.02)
    raise RuntimeError(
        f"真实 XLSX 未进入安全暂停：{manager.pause_status()}"
    )


def _assert_switch_did_not_resume(
    manager: IndexManager,
    database: DatabaseManager,
    progress_cursors: list[int],
) -> bool:
    before_cursor = max(progress_cursors, default=0)
    before_database = _paused_database_signature(database)
    time.sleep(0.5)
    after_database = _paused_database_signature(database)
    return bool(
        manager.pause_status()["state"] == "paused"
        and max(progress_cursors, default=0) == before_cursor
        and after_database == before_database
    )


def _content_and_search_snapshot(
    database: DatabaseManager,
    query: str,
) -> dict[str, object]:
    digest = hashlib.sha256()
    with database.connect() as connection:
        for row in connection.execute(
            """
            SELECT f.path, cb.block_index, cb.block_type,
                   cb.location_text, cb.raw_text
            FROM content_blocks cb
            JOIN files f ON f.id = cb.file_id
            ORDER BY f.path, cb.block_index, cb.id
            """
        ):
            for value in row:
                digest.update(
                    str(value if value is not None else "").encode(
                        "utf-8"
                    )
                )
                digest.update(b"\0")
        search_hits = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM content_blocks
                WHERE raw_text LIKE ?
                """,
                (f"%{query}%",),
            ).fetchone()[0]
        )
        duplicate_blocks = int(
            connection.execute(
                """
                SELECT COUNT(*) FROM (
                    SELECT file_id, block_index, COUNT(*) AS copies
                    FROM content_blocks
                    GROUP BY file_id, block_index
                    HAVING copies > 1
                )
                """
            ).fetchone()[0]
        )
    return {
        "content_digest": digest.hexdigest(),
        "search_hits": search_hits,
        "duplicate_blocks": duplicate_blocks,
    }


def _validate_mode_switch_rollback(base: Path) -> bool:
    from concurrent.futures import ThreadPoolExecutor

    base.mkdir(parents=True, exist_ok=True)
    original = AppSettings(parser_workers=1)
    manager = IndexManager(
        DatabaseManager(base / "rollback.db"),
        original,
        run_context={"execution_mode": "normal"},
    )
    old_executor = ThreadPoolExecutor(max_workers=1)
    lane = ParseLane(
        "normal",
        old_executor,
        1,
        1024,
        worker_count=1,
    )
    manager._pause_lanes = {"normal": lane}
    manager._pause_executors = [old_executor]
    manager._pause_process_executors = []
    manager._pause_spool_dir = base
    manager._pause_state = "paused"

    def fail_candidate(
        _jobs: list[ParseJob],
        _spool: Path,
    ) -> tuple[dict[str, ParseLane], list[object], list[object]]:
        raise RuntimeError("injected candidate pool failure")

    manager._create_lanes = fail_candidate  # type: ignore[method-assign]
    try:
        applied = manager.apply_settings_while_paused(
            AppSettings(parser_workers=2),
            execution_mode="performance",
            effective_profile={"mode": "performance"},
        )
        return bool(
            not applied
            and manager.settings is original
            and manager.pause_status()["state"] == "paused"
            and manager.run_context["execution_mode"] == "normal"
        )
    finally:
        old_executor.shutdown(
            wait=True,
            cancel_futures=True,
        )


def write_validation_result(
    name: str,
    validator: object,
) -> int:
    result_path = Path(f"{name}_validation_result.json")
    try:
        with tempfile.TemporaryDirectory(
            prefix=f"lfts_{name}_"
        ) as tmp:
            if callable(validator):
                payload = validator(Path(tmp))
            else:
                raise TypeError("validator is not callable")
    except Exception as exc:
        payload = {
            "passed": False,
            "error": f"{type(exc).__name__}: {exc}",
        }
    result_path.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(json.dumps(payload, ensure_ascii=False, indent=2))
    return 0 if payload.get("passed") else 1
