from __future__ import annotations

import argparse
import json
import shutil
import tempfile
import time
from pathlib import Path

from local_full_text_search.config.defaults import AppSettings
from local_full_text_search.core.database import DatabaseManager
from local_full_text_search.core.index_manager import IndexManager
from local_full_text_search.core.search_engine import SearchEngine
from local_full_text_search.models.search_query import SearchQuery


def create_corpus(root: Path, scale: int) -> None:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation

    root.mkdir(parents=True, exist_ok=True)
    for index in range(20 * scale):
        (root / f"text-{index:03d}.txt").write_text(
            "\n".join(f"BENCHMARK_TEXT_HIT row {row}" for row in range(300)),
            encoding="utf-8",
        )
    for index in range(8 * scale):
        document = Document()
        for paragraph in range(400):
            document.add_paragraph(f"BENCHMARK_DOCX_HIT {index} paragraph {paragraph}")
        document.save(root / f"document-{index:03d}.docx")
    for index in range(4 * scale):
        workbook = Workbook(write_only=True)
        sheet = workbook.create_sheet("Data")
        for row in range(2000):
            sheet.append([row, "BENCHMARK_XLSX_HIT", f"value-{index}-{row}"])
        workbook.save(root / f"workbook-{index:03d}.xlsx")
    for index in range(4 * scale):
        presentation = Presentation()
        for slide_index in range(40):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            box = slide.shapes.add_textbox(914400, 914400, 5000000, 914400)
            box.text = f"BENCHMARK_PPTX_HIT {index} slide {slide_index}"
        presentation.save(root / f"slides-{index:03d}.pptx")
    shutil.copy2(root / "document-000.docx", root / "document-duplicate.docx")


def run_case(corpus: Path, db_path: Path, *, fast: bool) -> dict[str, object]:
    settings = AppSettings(
        enable_ocr=False,
        fast_ooxml_enabled=fast,
        enable_parse_cache=fast,
        defer_fts_during_full_scan=fast,
        large_office_process_min_bytes=0,
        process_parser_workers=2,
        process_max_tasks_per_child=32,
    )
    db = DatabaseManager(db_path)
    db.initialize()
    root_id = db.add_root(corpus)
    started = time.perf_counter()
    summary = IndexManager(db, settings).index_root(root_id)
    elapsed_seconds = time.perf_counter() - started
    hits = {
        term: SearchEngine(db).search(SearchQuery(text=term)).total_confirmed
        for term in (
            "BENCHMARK_TEXT_HIT",
            "BENCHMARK_DOCX_HIT",
            "BENCHMARK_XLSX_HIT",
            "BENCHMARK_PPTX_HIT",
        )
    }
    with db.connect() as con:
        latest_run = con.execute(
            "SELECT summary_json FROM index_runs ORDER BY started_at DESC LIMIT 1"
        ).fetchone()
    return {
        "fast_ooxml": fast,
        "elapsed_seconds": round(elapsed_seconds, 3),
        "database_bytes": db_path.stat().st_size,
        "summary": summary.__dict__ if hasattr(summary, "__dict__") else {
            field: getattr(summary, field)
            for field in summary.__dataclass_fields__
        },
        "stats": db.stats(),
        "search_hits": hits,
        "run_metrics": json.loads(str(latest_run["summary_json"])) if latest_run else {},
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--scale", type=int, default=1)
    parser.add_argument("--output", type=Path, default=Path("build_validation/index_benchmark.json"))
    args = parser.parse_args()
    args.output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="lfts_benchmark_") as tmp:
        base = Path(tmp)
        corpus = base / "corpus"
        create_corpus(corpus, max(1, args.scale))
        legacy = run_case(corpus, base / "legacy.db", fast=False)
        optimized = run_case(corpus, base / "optimized.db", fast=True)
        baseline = float(legacy["elapsed_seconds"])
        current = float(optimized["elapsed_seconds"])
        report = {
            "legacy": legacy,
            "optimized": optimized,
            "speedup": round(baseline / current, 3) if current else None,
            "elapsed_reduction_percent": round((baseline - current) / baseline * 100, 2) if baseline else None,
        }
        args.output.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
        print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
