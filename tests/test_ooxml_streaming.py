from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from openpyxl.comments import Comment
from pptx import Presentation

from local_full_text_search.config.defaults import AppSettings
from local_full_text_search.core.task_manager import CancelToken
from local_full_text_search.parsers.ooxml.docx_stream_parser import DocxStreamParser
from local_full_text_search.parsers.ooxml.pptx_stream_parser import PptxStreamParser
from local_full_text_search.parsers.ooxml.xlsx_stream_parser import (
    XlsxStreamParser,
    encode_xlsx_cursor,
)
from local_full_text_search.parsers.parser_registry import ParserRegistry


class StreamingOoxmlTests(unittest.TestCase):
    def test_docx_streams_body_table_header_and_footer(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "sample.docx"
            document = Document()
            document.add_paragraph("DOCX_STREAM_BODY")
            table = document.add_table(rows=1, cols=2)
            table.cell(0, 0).text = "DOCX_TABLE_A"
            table.cell(0, 1).text = "DOCX_TABLE_B"
            document.sections[0].header.paragraphs[0].text = "DOCX_HEADER"
            document.sections[0].footer.paragraphs[0].text = "DOCX_FOOTER"
            document.save(path)

            blocks = list(DocxStreamParser().parse(path, CancelToken()))
            text = "\n".join(block.raw_text for block in blocks)

            self.assertIn("DOCX_STREAM_BODY", text)
            self.assertIn("DOCX_TABLE_A", text)
            self.assertIn("DOCX_TABLE_B", text)
            self.assertIn("DOCX_HEADER", text)
            self.assertIn("DOCX_FOOTER", text)

    def test_pptx_streams_slide_and_existing_notes(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "sample.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            box = slide.shapes.add_textbox(914400, 914400, 3657600, 914400)
            box.text = "PPTX_STREAM_SLIDE"
            slide.notes_slide.notes_text_frame.text = "PPTX_STREAM_NOTE"
            presentation.save(path)

            blocks = list(PptxStreamParser().parse(path, CancelToken()))
            text = "\n".join(block.raw_text for block in blocks)

            self.assertIn("PPTX_STREAM_SLIDE", text)
            self.assertIn("PPTX_STREAM_NOTE", text)
            self.assertTrue(all(block.slide_number == 1 for block in blocks))

    def test_pptx_resume_cursor_skips_confirmed_slides_and_reports_progress(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "resume.pptx"
            presentation = Presentation()
            for number in range(1, 4):
                slide = presentation.slides.add_slide(
                    presentation.slide_layouts[6]
                )
                box = slide.shapes.add_textbox(
                    914400,
                    914400,
                    3657600,
                    914400,
                )
                box.text = f"PPTX_RESUME_SLIDE_{number}"
                slide.notes_slide.notes_text_frame.text = (
                    f"PPTX_RESUME_NOTE_{number}"
                )
            presentation.save(path)

            progress: list[dict[str, object]] = []
            parser = PptxStreamParser()
            parser.configure_runtime(
                resume_cursor=2,
                progress_callback=progress.append,
            )
            blocks = list(parser.parse(path, CancelToken()))
            text = "\n".join(block.raw_text for block in blocks)

            self.assertTrue(parser.supports_resume)
            self.assertNotIn("PPTX_RESUME_SLIDE_1", text)
            self.assertNotIn("PPTX_RESUME_SLIDE_2", text)
            self.assertIn("PPTX_RESUME_SLIDE_3", text)
            self.assertIn("PPTX_RESUME_NOTE_3", text)
            self.assertEqual(
                [item["cursor"] for item in progress],
                [3],
            )

    def test_xlsx_streams_sparse_shared_formula_and_comment_cells(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "sample.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Data"
            sheet["D5"] = "XLSX_STREAM_VALUE"
            sheet["D5"].comment = Comment("XLSX_STREAM_COMMENT", "tester")
            sheet["E5"] = "=1+2"
            workbook.save(path)

            blocks = list(XlsxStreamParser().parse(path, CancelToken()))
            text = "\n".join(block.raw_text for block in blocks)

            self.assertIn("D5=XLSX_STREAM_VALUE", text)
            self.assertIn("XLSX_STREAM_COMMENT", text)
            self.assertIn("E5==1+2", text)
            self.assertEqual(blocks[0].sheet_name, "Data")

    def test_xlsx_stream_reports_workbook_sheet_and_row_progress(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "progress.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Large"
            for row in range(1, 301):
                sheet.cell(row=row, column=1).value = f"ROW_PROGRESS_{row:03d}"
            workbook.save(path)

            progress: list[dict[str, object]] = []
            parser = XlsxStreamParser()
            parser.configure_runtime(progress_callback=progress.append)
            blocks = list(parser.parse(path, CancelToken()))

            phases = {str(item.get("phase")) for item in progress}
            row_events = [item for item in progress if item.get("phase") == "sheet_row"]
            self.assertIn("workbook_scan", phases)
            self.assertIn("sheet_scan", phases)
            self.assertTrue(row_events)
            self.assertGreaterEqual(int(row_events[-1]["completed"]), 300)
            self.assertIn("ROW_PROGRESS_300", "\n".join(block.raw_text for block in blocks))

    def test_registry_selects_streaming_parsers_by_default(self) -> None:
        registry = ParserRegistry(AppSettings(fast_ooxml_enabled=True))

        self.assertIsInstance(registry.parser_for(Path("a.docx")), DocxStreamParser)
        self.assertIsInstance(registry.parser_for(Path("a.pptx")), PptxStreamParser)
        self.assertIsInstance(registry.parser_for(Path("a.xlsx")), XlsxStreamParser)

    def test_xlsx_parses_multiple_sheets_in_parallel_but_keeps_workbook_order(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "parallel.xlsx"
            workbook = Workbook()
            first = workbook.active
            first.title = "First"
            first["A1"] = "FIRST_SHEET"
            second = workbook.create_sheet("Second")
            second["A1"] = "SECOND_SHEET"
            third = workbook.create_sheet("Third")
            third["A1"] = "THIRD_SHEET"
            workbook.save(path)

            blocks = list(
                XlsxStreamParser(sheet_workers=2).parse(path, CancelToken())
            )

            self.assertEqual(
                [block.sheet_name for block in blocks],
                ["First", "Second", "Third"],
            )
            self.assertTrue(all(block.extra.get("sheet_parallel") for block in blocks))

    def test_xlsx_large_shared_strings_can_use_disk_backed_table(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            base = Path(tmp)
            path = base / "shared.xlsx"
            _write_minimal_shared_string_xlsx(path)

            blocks = list(
                XlsxStreamParser(
                    shared_strings_disk_threshold_bytes=1,
                    temp_dir=base / "xlsx-temp",
                ).parse(path, CancelToken())
            )

            self.assertIn("DISK_BACKED_SHARED_STRING", blocks[0].raw_text)
            self.assertEqual(blocks[0].extra["shared_strings_mode"], "disk")
            self.assertFalse(list((base / "xlsx-temp").glob("shared_strings_*")))

    def test_xlsx_resume_cursor_skips_confirmed_rows_in_large_sheet(self) -> None:
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "resume.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Rows"
            for row in range(1, 301):
                sheet.cell(row=row, column=1).value = f"RESUME_ROW_{row}"
            workbook.save(path)
            parser = XlsxStreamParser(sheet_workers=1)
            parser.configure_runtime(resume_cursor=encode_xlsx_cursor(1, 250))

            blocks = list(parser.parse(path, CancelToken()))
            text = "\n".join(block.raw_text for block in blocks)

            self.assertNotIn("RESUME_ROW_250", text)
            self.assertIn("RESUME_ROW_251", text)
            self.assertIn("RESUME_ROW_300", text)


def _write_minimal_shared_string_xlsx(path: Path) -> None:
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "xl/workbook.xml",
            """<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"
                xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
                <sheets><sheet name="Shared" sheetId="1" r:id="rId1"/></sheets>
            </workbook>""",
        )
        archive.writestr(
            "xl/_rels/workbook.xml.rels",
            """<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
                <Relationship Id="rId1"
                  Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet"
                  Target="worksheets/sheet1.xml"/>
            </Relationships>""",
        )
        archive.writestr(
            "xl/sharedStrings.xml",
            """<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"
                count="1" uniqueCount="1">
                <si><t>DISK_BACKED_SHARED_STRING</t></si>
            </sst>""",
        )
        archive.writestr(
            "xl/worksheets/sheet1.xml",
            """<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
                <dimension ref="A1:A1"/>
                <sheetData><row r="1"><c r="A1" t="s"><v>0</v></c></row></sheetData>
            </worksheet>""",
        )


if __name__ == "__main__":
    unittest.main()
