from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from local_full_text_search.core.task_manager import CancelToken
from local_full_text_search.parsers.pptx_parser import PptxParser


class PptxParserTests(unittest.TestCase):
    def test_slide_text_is_extracted_with_slide_number(self) -> None:
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx not installed")
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "sample.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(914400, 914400, 3657600, 914400)
            text_box.text = "PPTX_SLIDE_HIT"
            presentation.save(path)

            blocks = list(PptxParser().parse(path, CancelToken()))

            self.assertTrue(any(block.slide_number == 1 for block in blocks))
            self.assertTrue(any("PPTX_SLIDE_HIT" in block.raw_text for block in blocks))


if __name__ == "__main__":
    unittest.main()
