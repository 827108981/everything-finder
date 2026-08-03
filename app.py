from __future__ import annotations

import multiprocessing
import json
import os
import shutil
import sqlite3
import subprocess
import sys
import tempfile
import threading
import time
from dataclasses import asdict
from importlib import resources
from pathlib import Path

from local_full_text_search.services.startup_diagnostics import StartupDiagnostics


INDEX_STATUS_LAYOUT_CASES = (
    (1280, 800, 1.00, "ocr_running"),
    (1280, 800, 1.25, "ocr_running"),
    (1280, 800, 1.50, "ocr_running"),
    (1366, 768, 1.00, "pausing"),
    (1920, 1080, 1.50, "paused_switched"),
)


def load_application_dependencies() -> None:
    """Load optional/heavy application modules after startup diagnostics is active."""

    global APP_DISPLAY_NAME, AppSettings, CancelToken, DatabaseManager, IndexManager
    global SCHEMA_VERSION
    global SearchEngine, SearchQuery, SettingsService, TEMP_DIR, configure_logging

    from local_full_text_search.config.constants import APP_DISPLAY_NAME, TEMP_DIR
    from local_full_text_search.config.defaults import AppSettings
    from local_full_text_search.core.database import SCHEMA_VERSION, DatabaseManager
    from local_full_text_search.core.index_manager import IndexManager
    from local_full_text_search.core.search_engine import SearchEngine
    from local_full_text_search.core.task_manager import CancelToken
    from local_full_text_search.models.search_query import SearchQuery
    from local_full_text_search.services.logging_service import configure_logging
    from local_full_text_search.services.settings_service import SettingsService


def run_self_test() -> int:
    """Non-interactive smoke test used after packaging."""

    try:
        from PySide6.QtWidgets import QApplication
        from local_full_text_search.ui.main_window import MainWindow

        with tempfile.TemporaryDirectory(prefix="lfts_self_test_") as tmp:
            base = Path(tmp)
            settings_service = SettingsService(base / "settings.json")
            settings = AppSettings(enable_ocr=False, monitor_file_changes=False)
            settings_service.save(settings)
            db = DatabaseManager(base / "self_test.db")
            db.initialize()
            app = QApplication.instance() or QApplication(["LocalFullTextSearch", "--self-test"])
            apply_light_theme(app)
            window = MainWindow(db, settings, settings_service)
            window.close()
        Path("self_test_result.txt").write_text("SELF_TEST_OK\n", encoding="utf-8")
        print("SELF_TEST_OK")
        return 0
    except Exception as exc:
        Path("self_test_result.txt").write_text(f"SELF_TEST_FAILED: {exc}\n", encoding="utf-8")
        print(f"SELF_TEST_FAILED: {exc}")
        return 1


def run_ui_validation() -> int:
    """Render the desktop UI offscreen and reject blank or undersized output."""

    result_path = Path("ui_validation_result.txt")
    image_path = Path("ui_validation.png")
    index_image_path = Path("index_ui_validation.png")
    paused_index_image_path = Path("index_paused_ui_validation.png")
    try:
        from PySide6.QtGui import QImage
        from PySide6.QtWidgets import QApplication
        from local_full_text_search.ui.main_window import MainWindow

        with tempfile.TemporaryDirectory(prefix="lfts_ui_validation_") as tmp:
            base = Path(tmp)
            settings_service = SettingsService(base / "settings.json")
            settings = AppSettings(enable_ocr=False, monitor_file_changes=False)
            settings_service.save(settings)
            db = DatabaseManager(base / "ui.db")
            db.initialize()
            sample_root = base / "files"
            sample_root.mkdir()
            sample_file = sample_root / (
                "CL8000i-301140100~301140400-第一抓杯手手指组件拆卸方法测试文档.txt"
            )
            sample_file.write_text(
                "第 6 页\n操作前请拔掉3 个传感器，并标记插头。\n"
                "完整路径、日期和右侧命中内容都必须可见。",
                encoding="utf-8",
            )
            root_id = db.add_root(sample_root)
            IndexManager(db, settings).index_root(root_id)
            app = QApplication.instance() or QApplication(["LocalFullTextSearch", "--validate-ui"])
            apply_light_theme(app)
            window = MainWindow(db, settings, settings_service)
            window.resize(1280, 800)
            window.show()
            window.search_page.search_box.input.setText("拔掉 3 个传感器")
            page = SearchEngine(db).search(
                SearchQuery(text="拔掉 3 个传感器", mode="exact")
            )
            window.search_page.set_results(page)
            window.show_preview(page.results[0])
            app.processEvents()
            card = window.search_page.result_view._cards[0]
            item = window.search_page.result_view.item(0)
            if "background:#FDE68A" not in card.context.text():
                raise RuntimeError("Search result did not highlight the space-tolerant match")
            if "background-color:#fde68a" not in window.preview_panel.context.toHtml().lower():
                raise RuntimeError("Preview did not highlight the selected result")
            if item is None or card.path_label.geometry().bottom() > item.sizeHint().height():
                raise RuntimeError("Search result card clipped the path or date")
            search_box = window.search_page.search_box
            if search_box.input.geometry().right() >= search_box.action_separator.geometry().left():
                raise RuntimeError("Search actions overlap the text input")
            pixmap = window.grab()
            image = pixmap.toImage().convertToFormat(QImage.Format.Format_RGB32)
            if image.width() < 1000 or image.height() < 650:
                raise RuntimeError(f"Unexpected UI size: {image.width()}x{image.height()}")
            colors = {
                image.pixelColor(x, y).rgba()
                for x in range(0, image.width(), 32)
                for y in range(0, image.height(), 32)
            }
            if not pixmap.save(str(image_path), "PNG"):
                raise RuntimeError("Unable to save UI screenshot")
            if len(colors) < 5:
                raise RuntimeError(f"UI render appears blank: {len(colors)} sampled colors")
            window.switch_page("index")
            window.on_scan_progress(
                {
                    "stage": "indexing",
                    "phase_label": "正在解析并写入索引",
                    "indexed": 104,
                    "scanned": 337,
                    "completed_files": 104,
                    "total_files": 337,
                    "failed": 1,
                    "current_file": str(sample_root / "large-scanned-manual.pdf"),
                    "eta_seconds": 660,
                    "eta_ready": True,
                    "eta_confidence": 0.82,
                    "active_queue": "pdf",
                    "active_elapsed_seconds": 252,
                    "active_phase": "pdf_ocr_tile_recognize_microbatch",
                    "active_completed_units": 128,
                    "active_total_units": 240,
                    "no_progress_seconds": 7,
                    "retry_count": 1,
                    "excluded_video": 9,
                }
            )
            app.processEvents()
            index_pixmap = window.grab()
            if not index_pixmap.save(str(index_image_path), "PNG"):
                raise RuntimeError("Unable to save index UI screenshot")
            if "预计剩余约 11 分钟" not in window.index_page.task_eta.text():
                raise RuntimeError("Index UI did not render the single-value ETA")
            if not window.index_page.task_runtime.text().startswith("OCR 已运行"):
                raise RuntimeError("Index UI did not render the OCR runtime on row two")
            window.index_page.set_pause_state("paused", "普通模式 · 已暂停")
            app.processEvents()
            if (
                not window.index_page.start_button.isEnabled()
                or not window.index_page.performance_button.isEnabled()
            ):
                raise RuntimeError("Mode controls were not enabled after safe pause")
            paused_pixmap = window.grab()
            if not paused_pixmap.save(str(paused_index_image_path), "PNG"):
                raise RuntimeError("Unable to save paused index UI screenshot")
            window.switch_page("failed")
            window.failed_page.set_rows(
                [
                    {
                        "id": 999,
                        "path": str(sample_root / "stale-pdf-task.pdf"),
                        "filename": "stale-pdf-task.pdf",
                        "extension": ".pdf",
                        "parse_status": "pending",
                        "parse_error_message": "任务未完成",
                    }
                ]
            )
            window.failed_page.set_readiness(
                {
                    "ready": False,
                    "repairable": True,
                    "not_ready_reasons": ["blocking_files"],
                }
            )
            app.processEvents()
            force_button = window.failed_page.force_complete_button
            if force_button.text() != "强力完成本次索引":
                raise RuntimeError("Failed-files UI did not expose force completion")
            window.failed_page.set_rows([])
            window.failed_page.set_readiness(
                {
                    "ready": False,
                    "repairable": True,
                    "not_ready_reasons": [
                        "unfinished_tasks",
                        "content_fts_dirty",
                    ],
                }
            )
            app.processEvents()
            if force_button.isHidden() or not force_button.isEnabled():
                raise RuntimeError(
                    "Force completion must repair a zero-blocker residual state"
                )
            if "残留解析任务" not in force_button.toolTip():
                raise RuntimeError("Force completion did not disclose residual tasks")
            window.failed_page.set_readiness(
                {
                    "ready": True,
                    "repairable": False,
                    "not_ready_reasons": [],
                }
            )
            app.processEvents()
            if force_button.isEnabled():
                raise RuntimeError("Force completion must be disabled after readiness")
            window.failed_page.set_rows(
                [
                    {
                        "id": 999,
                        "path": str(sample_root / "stale-pdf-task.pdf"),
                        "filename": "stale-pdf-task.pdf",
                        "extension": ".pdf",
                        "parse_status": "pending",
                        "parse_error_message": "任务未完成",
                    }
                ]
            )
            window.failed_page.set_readiness(
                {
                    "ready": False,
                    "repairable": True,
                    "not_ready_reasons": ["blocking_files"],
                }
            )
            app.processEvents()
            if not force_button.isEnabled():
                raise RuntimeError(
                    "Force completion control must be enabled with blockers"
                )
            if force_button.mapTo(
                window.failed_page,
                force_button.rect().bottomLeft(),
            ).y() >= window.failed_page.table.geometry().top():
                raise RuntimeError("Force completion control is below the failed-files table")
            if force_button.geometry().intersects(
                window.failed_page.status.geometry()
            ):
                raise RuntimeError("Force completion control overlaps status text")
            window.failed_page.set_exclusion_running(True)
            window.failed_page.set_exclusion_progress(
                {
                    "stage": "rebuilding_content_fts",
                    "phase_label": "正在重建全文索引",
                    "large_fts_operation": True,
                    "can_cancel": True,
                    "elapsed_seconds": 3,
                }
            )
            app.processEvents()
            if window.failed_page.exclusion_progress.isHidden():
                raise RuntimeError("Scope exclusion progress is not visible")
            if (
                window.failed_page.exclude_button.isEnabled()
                or window.failed_page.restore_button.isEnabled()
                or force_button.isEnabled()
            ):
                raise RuntimeError("Conflicting controls remain enabled during exclusion")
            if "大型 FTS 操作" not in window.failed_page.status.text():
                raise RuntimeError("Scope exclusion FTS stage is not disclosed")
            window.failed_page.set_exclusion_running(False)
            window.close()
            app.processEvents()
            message = (
                "UI_VALIDATION_OK\n"
                f"size={image.width()}x{image.height()}; sampled_colors={len(colors)}\n"
            )
            result_path.write_text(message, encoding="utf-8")
            print(message, end="")
            return 0
    except Exception as exc:
        result_path.write_text(f"UI_VALIDATION_FAILED: {exc}\n", encoding="utf-8")
        print(f"UI_VALIDATION_FAILED: {exc}")
        return 1


def run_core_validation() -> int:
    """Build a temporary multi-format index and verify exact search hits."""

    result_path = Path("core_validation_result.txt")
    try:
        with tempfile.TemporaryDirectory(prefix="lfts_validation_") as tmp:
            base = Path(tmp)
            root = base / "files"
            root.mkdir()
            _create_validation_files(root)

            db = DatabaseManager(base / "validation.db")
            db.initialize()
            root_id = db.add_root(root)
            summary = IndexManager(db, AppSettings()).index_root(root_id)
            engine = SearchEngine(db)

            checks = {
                "TXT_VALIDATION_HIT": "txt",
                "PDF_VALIDATION_HIT": "pdf",
                "DOCX_VALIDATION_HIT": "docx",
                "XLSX_VALIDATION_HIT": "xlsx",
                "PPTX_VALIDATION_HIT": "pptx",
                "OCR TEST 123": "ocr_image",
            }
            failures: list[str] = []
            for term, label in checks.items():
                page = engine.search(
                    SearchQuery(
                        text=term,
                        mode="exact",
                        include_ocr_fuzzy=label == "ocr_image",
                    )
                )
                if page.total_confirmed < 1:
                    failures.append(f"{label}:{term}")
            if failures:
                diagnostic_text = "; ".join(
                    f"{row['path']} [{row['parse_error_code']}]: {row['parse_error_message']}"
                    for row in db.failed_files()
                )
                detail = f"; diagnostics={diagnostic_text}" if diagnostic_text else ""
                raise RuntimeError("未命中: " + ", ".join(failures) + detail)
            message = (
                "CORE_VALIDATION_OK\n"
                f"scanned={summary.scanned}; indexed={summary.indexed}; "
                f"skipped={summary.skipped}; failed={summary.failed}; unsupported={summary.unsupported}\n"
            )
            result_path.write_text(message, encoding="utf-8")
            print(message, end="")
            return 0
    except Exception as exc:
        result_path.write_text(f"CORE_VALIDATION_FAILED: {exc}\n", encoding="utf-8")
        print(f"CORE_VALIDATION_FAILED: {exc}")
        return 1


def run_process_pool_validation() -> int:
    """Force DOCX/XLSX through the spawned process parser in a frozen build."""

    result_path = Path("process_pool_validation_result.txt")
    try:
        with tempfile.TemporaryDirectory(prefix="lfts_process_validation_") as tmp:
            base = Path(tmp)
            root = base / "files"
            root.mkdir()
            _create_process_validation_files(root)

            db = DatabaseManager(base / "validation.db")
            db.initialize()
            root_id = db.add_root(root)
            queues: list[str] = []
            worker_pids: list[int] = []
            result_bytes: list[int] = []
            descriptor_bytes: list[int] = []

            def capture_progress(payload: dict[str, object]) -> None:
                queue_name = payload.get("completed_queue")
                if payload.get("stage") == "indexing" and isinstance(queue_name, str):
                    queues.append(queue_name)
                    if isinstance(payload.get("worker_pid"), int):
                        worker_pids.append(int(payload["worker_pid"]))
                    if isinstance(payload.get("process_result_bytes"), int):
                        result_bytes.append(int(payload["process_result_bytes"]))
                    if isinstance(payload.get("process_descriptor_bytes"), int):
                        descriptor_bytes.append(int(payload["process_descriptor_bytes"]))

            settings = AppSettings(
                enable_ocr=False,
                process_parser_workers=1,
                process_pending_tasks=1,
                process_max_tasks_per_child=1,
                large_office_process_min_bytes=0,
            )
            summary = IndexManager(db, settings).index_root(root_id, progress_callback=capture_progress)
            engine = SearchEngine(db)
            terms = ("PROCESS_DOCX_VALIDATION_HIT", "PROCESS_XLSX_VALIDATION_HIT")
            failures = [
                term
                for term in terms
                if engine.search(SearchQuery(text=term, mode="exact")).total_confirmed < 1
            ]
            if failures:
                raise RuntimeError("Missing process-pool hits: " + ", ".join(failures))
            if queues.count("office_process") != 2:
                raise RuntimeError(f"Unexpected parse queues: {queues}")
            if len(set(worker_pids)) != 2:
                raise RuntimeError(f"Process recycling was not observed: {worker_pids}")
            if not descriptor_bytes or max(descriptor_bytes) >= 4096:
                raise RuntimeError(f"Process descriptor transfer is unexpectedly large: {descriptor_bytes}")
            message = (
                "PROCESS_POOL_VALIDATION_OK\n"
                f"scanned={summary.scanned}; indexed={summary.indexed}; failed={summary.failed}; "
                f"office_process={queues.count('office_process')}; recycled_workers={len(set(worker_pids))}; "
                f"spooled_bytes={sum(result_bytes)}; ipc_descriptor_bytes={sum(descriptor_bytes)}\n"
            )
            result_path.write_text(message, encoding="utf-8")
            print(message, end="")
            return 0
    except Exception as exc:
        result_path.write_text(f"PROCESS_POOL_VALIDATION_FAILED: {exc}\n", encoding="utf-8")
        print(f"PROCESS_POOL_VALIDATION_FAILED: {exc}")
        return 1


def run_schema_v3_validation() -> int:
    """Verify deduplication, search expansion, metrics, and SQLite integrity."""

    result_path = Path("schema_v3_validation_result.txt")
    try:
        with tempfile.TemporaryDirectory(prefix="lfts_schema_v3_") as tmp:
            base = Path(tmp)
            root = base / "files"
            root.mkdir()
            first = root / "first.txt"
            second = root / "second.txt"
            first.write_text("SCHEMA_V3_VALIDATION_HIT", encoding="utf-8")
            shutil.copy2(first, second)
            db = DatabaseManager(base / "validation.db")
            db.initialize()
            root_id = db.add_root(root)
            summary = IndexManager(db, AppSettings(enable_ocr=False)).index_root(root_id)
            page = SearchEngine(db).search(
                SearchQuery(
                    text="SCHEMA_V3_VALIDATION_HIT",
                    search_filename=False,
                    search_path=False,
                )
            )
            stats = db.stats()
            integrity = db.integrity_report()
            if summary.indexed != 2 or stats["documents"] != 1 or stats["blocks"] != 1:
                raise RuntimeError(f"Unexpected dedup stats: summary={summary}; stats={stats}")
            if page.total_confirmed != 2:
                raise RuntimeError(f"Duplicate paths were not expanded: {page.total_confirmed}")
            if integrity["integrity"] != ["ok"] or integrity["foreign_key_errors"]:
                raise RuntimeError(f"SQLite integrity failed: {integrity}")
            with db.connect() as con:
                run = con.execute(
                    "SELECT status, discovered_files FROM index_runs ORDER BY started_at DESC LIMIT 1"
                ).fetchone()
                task_status = con.execute(
                    "SELECT status FROM parse_tasks ORDER BY id DESC LIMIT 1"
                ).fetchone()
            if run is None or run["status"] != "complete" or int(run["discovered_files"]) != 2:
                raise RuntimeError("Run metrics were not finalized")
            if task_status is None or task_status["status"] != "complete":
                raise RuntimeError("Parse task state was not finalized")
            message = (
                "SCHEMA_V3_VALIDATION_OK\n"
                f"files={stats['files']}; documents={stats['documents']}; blocks={stats['blocks']}; "
                f"expanded_hits={page.total_confirmed}\n"
            )
            result_path.write_text(message, encoding="utf-8")
            print(message, end="")
            return 0
    except Exception as exc:
        result_path.write_text(f"SCHEMA_V3_VALIDATION_FAILED: {exc}\n", encoding="utf-8")
        print(f"SCHEMA_V3_VALIDATION_FAILED: {exc}")
        return 1


def run_schema_v4_validation() -> int:
    """Migrate a token-heavy schema-v3 database and verify a real update."""

    result_path = Path("schema_v4_validation_result.txt")
    try:
        with tempfile.TemporaryDirectory(prefix="lfts_schema_v4_") as tmp:
            base = Path(tmp)
            root = base / "files"
            root.mkdir()
            source = root / "legacy-tokens.txt"
            source.write_text("SCHEMA_V4_OLD_HIT", encoding="utf-8")
            db_path = base / "validation.db"
            db = DatabaseManager(db_path)
            db.initialize()
            root_id = db.add_root(root)
            IndexManager(db, AppSettings(enable_ocr=False)).index_root(root_id)

            legacy_token_count = 250_000
            with db.connect() as con:
                block_id = int(con.execute("SELECT id FROM content_blocks LIMIT 1").fetchone()[0])
                con.execute("DROP INDEX idx_short_tokens_block")
                con.executemany(
                    "INSERT INTO short_tokens(token, block_id) VALUES (?, ?)",
                    ((f"legacy-{index}", block_id) for index in range(legacy_token_count)),
                )
                con.execute("PRAGMA user_version = 3")

            migration_started = time.perf_counter()
            migrated = DatabaseManager(db_path)
            migrated.initialize()
            migration_ms = int((time.perf_counter() - migration_started) * 1000)
            with migrated.connect() as con:
                token_rows = int(con.execute("SELECT COUNT(*) FROM short_tokens").fetchone()[0])
                version = int(con.execute("PRAGMA user_version").fetchone()[0])
                indexes = {
                    str(row["name"]) for row in con.execute("PRAGMA index_list(short_tokens)")
                }
                plan = " ".join(
                    str(row["detail"])
                    for row in con.execute(
                        "EXPLAIN QUERY PLAN DELETE FROM short_tokens WHERE block_id = ?",
                        (block_id,),
                    )
                )
            if token_rows != 0 or version != 4:
                raise RuntimeError(f"Unexpected migration state: tokens={token_rows}; version={version}")
            if "idx_short_tokens_block" not in indexes or "idx_short_tokens_block" not in plan:
                raise RuntimeError(f"Missing block lookup index: indexes={indexes}; plan={plan}")
            if migrated.requires_full_rebuild():
                raise RuntimeError("Schema-v4 migration incorrectly forced a content reparse")

            source.write_text("SCHEMA_V4_NEW_HIT", encoding="utf-8")
            update_started = time.perf_counter()
            summary = IndexManager(migrated, AppSettings(enable_ocr=False)).index_root(root_id)
            update_ms = int((time.perf_counter() - update_started) * 1000)
            hits = SearchEngine(migrated).search(
                SearchQuery(text="SCHEMA_V4_NEW_HIT", mode="exact")
            ).total_confirmed
            report = migrated.integrity_report()
            backup = base / "validation.schema-v3.backup.db"
            if not backup.is_file():
                raise RuntimeError("Schema-v3 backup was not created")
            if summary.indexed != 1 or summary.failed != 0 or hits != 1:
                raise RuntimeError(
                    f"Post-migration update failed: summary={summary}; hits={hits}"
                )
            if report["integrity"] != ["ok"] or report["foreign_key_errors"]:
                raise RuntimeError(f"SQLite integrity failed: {report}")
            message = (
                "SCHEMA_V4_VALIDATION_OK\n"
                f"legacy_tokens={legacy_token_count}; migration_ms={migration_ms}; "
                f"update_ms={update_ms}; indexed={summary.indexed}; hits={hits}\n"
            )
            result_path.write_text(message, encoding="utf-8")
            print(message, end="")
            return 0
    except Exception as exc:
        result_path.write_text(f"SCHEMA_V4_VALIDATION_FAILED: {exc}\n", encoding="utf-8")
        print(f"SCHEMA_V4_VALIDATION_FAILED: {exc}")
        return 1


def seed_failure_fallback_demo(
    db: object,
    demo_root: Path,
) -> dict[str, object]:
    """Create a repeatable mixed-success scope for hands-on fallback testing."""

    import zipfile

    from local_full_text_search.models.content_block import ContentBlock

    demo_root.mkdir(parents=True, exist_ok=True)
    sources = {
        "searchable": demo_root / "正常可搜索文件.txt",
        "broken_pdf": demo_root / "损坏的PDF.pdf",
        "password_word": demo_root / "受密码保护的Word文档.docx",
        "ocr_failed": demo_root / "OCR识别失败图片.png",
        "oversize_zip": demo_root / "成员数超过安全上限的压缩包.zip",
        "metadata_zip": demo_root / "仅元数据完成的压缩包.zip",
    }
    sources["searchable"].write_text(
        "DEMO_SEARCH_OK\n该文件用于确认排除异常项后，正常正文仍然可以搜索。\n",
        encoding="utf-8",
    )
    sources["broken_pdf"].write_bytes(b"%PDF-1.4\ninvalid demo document")
    sources["password_word"].write_bytes(b"PK\x03\x04encrypted-demo")
    sources["ocr_failed"].write_bytes(b"not-a-readable-png")
    with zipfile.ZipFile(sources["oversize_zip"], "w") as archive:
        for index in range(2001):
            archive.writestr(f"demo-{index:04d}.txt", b"")
    with zipfile.ZipFile(sources["metadata_zip"], "w") as archive:
        archive.writestr("unsupported-demo.bin", b"demo")

    root_id = db.add_root(demo_root)
    file_ids = {
        name: db.upsert_file_metadata(root_id, path)[0]
        for name, path in sources.items()
    }
    searchable = sources["searchable"]
    db.replace_file_blocks(
        file_ids["searchable"],
        searchable.name,
        str(searchable),
        [
            ContentBlock(
                file_path=str(searchable),
                block_index=0,
                block_type="paragraph",
                location_text="演示正文",
                raw_text="DEMO_SEARCH_OK 正常文件仍然可以搜索",
                normalized_text="demo_search_ok 正常文件仍然可以搜索",
            )
        ],
        parser_name="text",
    )
    db.record_failure(
        file_ids["broken_pdf"],
        "PDF_CORRUPTED",
        "演示：PDF 结构损坏，无法读取页面",
        parser_name="pdf",
    )
    db.set_file_error_status(
        file_ids["password_word"],
        "password_protected",
        "PASSWORD_PROTECTED",
        "演示：文档受密码保护",
        parser_name="docx",
    )
    db.set_file_error_status(
        file_ids["ocr_failed"],
        "ocr_failed",
        "OCR_FAILED",
        "演示：图片无法解码，OCR 识别失败",
        parser_name="image_ocr",
    )
    db.set_file_error_status(
        file_ids["oversize_zip"],
        "skipped",
        "ZIP_FILE_COUNT_LIMIT",
        "演示：压缩包成员数超过 2,000 个安全上限",
        parser_name="zip",
    )
    db.set_file_error_status(
        file_ids["metadata_zip"],
        "metadata_only",
        "ZIP_NO_SUPPORTED_MEMBER",
        "演示：压缩包内没有支持的正文文件",
        parser_name="zip",
    )
    db.update_root_scan_time(root_id, "incomplete")
    db.begin_deferred_fts()
    return {
        "root_id": root_id,
        "file_ids": file_ids,
        "demo_root": str(demo_root),
        "search_token": "DEMO_SEARCH_OK",
    }


def reset_failure_fallback_demo(app_data_dir: Path, db_path: Path) -> Path:
    app_data_dir = app_data_dir.resolve()
    db_path = db_path.resolve()
    if not db_path.is_relative_to(app_data_dir):
        raise RuntimeError("演示数据库不在隔离数据目录内，已拒绝重置")
    for suffix in ("", "-wal", "-shm"):
        Path(str(db_path) + suffix).unlink(missing_ok=True)
    demo_root = app_data_dir / "demo-files"
    if demo_root.exists():
        shutil.rmtree(demo_root)
    return demo_root


def run_failure_fallback_demo_validation() -> int:
    result_path = Path("failure_fallback_demo_validation_result.json")
    try:
        with tempfile.TemporaryDirectory(prefix="lfts_failure_demo_") as tmp:
            base = Path(tmp)
            db = DatabaseManager(base / "demo.db")
            db.initialize()
            seeded = seed_failure_fallback_demo(db, base / "demo-files")
            before = db.index_readiness()
            result = db.force_complete_current_scope(
                reason="冻结验证：排除四个演示阻断项并开放搜索",
                operation_source="failure_demo_validation",
            )
            after = db.index_readiness()
            hits = SearchEngine(db).search(
                SearchQuery(
                    text=str(seeded["search_token"]),
                    mode="exact",
                )
            ).total_confirmed
            integrity = db.integrity_report()
            payload = {
                "passed": (
                    int(before["blocking_files"]) == 4
                    and len(db.excluded_files(limit=20)) == 4
                    and int(result["excluded_files"]) == 4
                    and bool(after["ready"])
                    and hits == 1
                    and integrity["integrity"] == ["ok"]
                    and not integrity["foreign_key_errors"]
                ),
                "blocking_before": int(before["blocking_files"]),
                "excluded_after": len(db.excluded_files(limit=20)),
                "metadata_only_files": int(before["metadata_only_complete_files"]),
                "ready_after": bool(after["ready"]),
                "search_hits": hits,
                "integrity": integrity["integrity"],
                "foreign_key_errors": len(integrity["foreign_key_errors"]),
            }
        result_path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 0 if payload["passed"] else 1
    except Exception as exc:
        payload = {"passed": False, "error": f"{type(exc).__name__}: {exc}"}
        result_path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 1


def capture_index_status_layout(
    width: int,
    height: int,
    scale: float,
    state: str,
    output_path: Path,
) -> int:
    result_path = output_path.with_suffix(".json")
    try:
        from PySide6.QtGui import QFontDatabase
        from PySide6.QtWidgets import QApplication
        from local_full_text_search.ui.main_window import MainWindow

        with tempfile.TemporaryDirectory(
            prefix="lfts_index_layout_"
        ) as tmp:
            base = Path(tmp)
            settings_service = SettingsService(
                base / "settings.json"
            )
            settings = AppSettings(
                enable_ocr=False,
                monitor_file_changes=False,
            )
            settings_service.save(settings)
            database = DatabaseManager(base / "layout.db")
            database.initialize()
            application = QApplication.instance() or QApplication(
                ["LocalFullTextSearch", "--capture-index-status-layout"]
            )
            font_families = set(QFontDatabase.families())
            if os.name == "nt" and not {
                "Microsoft YaHei UI",
                "Microsoft YaHei",
            }.intersection(font_families):
                raise RuntimeError(
                    "Windows Chinese UI fonts are unavailable"
                )
            apply_light_theme(application)
            window = MainWindow(
                database,
                settings,
                settings_service,
            )
            window.resize(int(width), int(height))
            window.show()
            window.switch_page("index")
            if state == "paused_switched":
                window.index_page.set_performance_mode(True)
            window.on_scan_progress(
                {
                    "stage": "indexing",
                    "phase_label": "正在解析并写入索引",
                    "indexed": 209,
                    "scanned": 309,
                    "completed_files": 209,
                    "total_files": 309,
                    "failed": 2,
                    "current_file": (
                        "E:\\MRCODE\\化免\\免疫资料\\9000i\\"
                        "超长中文文件名_微信图片_20260421082832_32_45.jpg"
                    ),
                    "eta_seconds": 660,
                    "eta_ready": True,
                    "eta_confidence": 0.82,
                    "active_queue": "ocr",
                    "active_elapsed_seconds": 252,
                    "active_file_count": 3,
                    "active_phase": (
                        "ocr_recognize_original_regions"
                    ),
                    "active_completed_units": 128,
                    "active_total_units": 240,
                    "no_progress_seconds": 7,
                    "retry_count": 1,
                    "excluded_video": 22,
                }
            )
            if state == "pausing":
                window.index_page.set_pause_state(
                    "pausing",
                    "普通模式 · 正在暂停",
                )
            elif state == "paused_switched":
                window.index_page.set_pause_state(
                    "paused",
                    "性能模式 · 已暂停",
                )
            application.processEvents()
            page = window.index_page
            labels = {
                "overview": page.task_label,
                "eta": page.task_eta,
                "runtime": page.task_runtime,
                "phase": page.task_phase,
                "units": page.task_units,
            }
            clipped: list[str] = []
            for name, label in labels.items():
                required = label.fontMetrics().horizontalAdvance(
                    label.text()
                )
                available = label.contentsRect().width()
                if label.text() and required > available + 2:
                    clipped.append(
                        f"{name}:{required}>{available}"
                    )
            buttons = (
                page.start_button,
                page.performance_button,
                page.pause_button,
                page.cancel_button,
            )
            missing_buttons = [
                button.text()
                for button in buttons
                if not button.isVisible()
                or button.width()
                < button.sizeHint().width() - 2
            ]
            if page.task_file.toolTip() != (
                "E:\\MRCODE\\化免\\免疫资料\\9000i\\"
                "超长中文文件名_微信图片_20260421082832_32_45.jpg"
            ):
                raise RuntimeError(
                    "Full current path is not available as tooltip"
                )
            if clipped:
                raise RuntimeError(
                    "Critical index status text is clipped: "
                    + ", ".join(clipped)
                )
            if missing_buttons:
                raise RuntimeError(
                    "Index controls are clipped or hidden: "
                    + ", ".join(missing_buttons)
                )
            output_path.parent.mkdir(
                parents=True,
                exist_ok=True,
            )
            pixmap = window.grab()
            if not pixmap.save(str(output_path), "PNG"):
                raise RuntimeError("Unable to save layout screenshot")
            payload = {
                "passed": True,
                "requested_size": [width, height],
                "scale": scale,
                "state": state,
                "device_pixel_ratio": round(
                    float(window.devicePixelRatioF()),
                    3,
                ),
                "font_family_count": len(font_families),
                "image_size": [
                    pixmap.width(),
                    pixmap.height(),
                ],
                "critical_text": {
                    name: label.text()
                    for name, label in labels.items()
                },
                "file_tooltip": page.task_file.toolTip(),
                "buttons": [
                    {
                        "text": button.text(),
                        "enabled": button.isEnabled(),
                    }
                    for button in buttons
                ],
            }
            result_path.write_text(
                json.dumps(
                    payload,
                    ensure_ascii=False,
                    indent=2,
                ),
                encoding="utf-8",
            )
            window.close()
            application.processEvents()
        print(json.dumps(payload, ensure_ascii=False))
        return 0
    except Exception as exc:
        payload = {
            "passed": False,
            "error": f"{type(exc).__name__}: {exc}",
            "requested_size": [width, height],
            "scale": scale,
            "state": state,
        }
        result_path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        print(json.dumps(payload, ensure_ascii=False))
        return 1


def run_index_status_layout_validation() -> int:
    result_path = Path(
        "index_status_layout_validation_result.json"
    )
    screenshot_dir = Path("index_status_layout_screenshots")
    screenshot_dir.mkdir(parents=True, exist_ok=True)
    executable = Path(sys.executable).resolve()
    command_prefix = (
        [str(executable)]
        if getattr(sys, "frozen", False)
        else [str(executable), str(Path(__file__).resolve())]
    )
    reports: list[dict[str, object]] = []
    for width, height, scale, state in INDEX_STATUS_LAYOUT_CASES:
        label = (
            f"{width}x{height}_{round(scale * 100)}pct_{state}"
        )
        output = (
            screenshot_dir / f"{label}.png"
        ).resolve()
        environment = os.environ.copy()
        environment["QT_QPA_PLATFORM"] = (
            "windows" if os.name == "nt" else "offscreen"
        )
        environment["QT_SCALE_FACTOR"] = str(scale)
        environment["QT_ENABLE_HIGHDPI_SCALING"] = "1"
        environment["PYTHONUTF8"] = "1"
        completed = subprocess.run(
            [
                *command_prefix,
                "--capture-index-status-layout",
                str(width),
                str(height),
                str(scale),
                state,
                str(output),
            ],
            cwd=Path.cwd(),
            env=environment,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            check=False,
        )
        case_result = output.with_suffix(".json")
        if case_result.is_file():
            payload = json.loads(
                case_result.read_text(encoding="utf-8")
            )
        else:
            payload = {
                "passed": False,
                "error": (
                    "capture result is missing; "
                    f"stdout={completed.stdout}; "
                    f"stderr={completed.stderr}"
                ),
            }
        payload["exit_code"] = completed.returncode
        payload["screenshot"] = str(output)
        reports.append(payload)
    passed = bool(
        len(reports) == len(INDEX_STATUS_LAYOUT_CASES)
        and all(
            bool(report.get("passed"))
            and int(report.get("exit_code") or 0) == 0
            for report in reports
        )
    )
    result = {
        "passed": passed,
        "cases": reports,
    }
    result_path.write_text(
        json.dumps(result, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0 if passed else 1


def run_hang_recovery_validation_command() -> int:
    from local_full_text_search.core.hang_validation import (
        run_hang_recovery_validation,
    )

    report = run_hang_recovery_validation(
        Path("hang_recovery_validation_result.json"),
        timeout_seconds=1.0,
    )
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if report["passed"] else 1


def run_semantic_progress_validation_command() -> int:
    from local_full_text_search.core.hang_validation import (
        run_semantic_progress_validation,
    )

    report = run_semantic_progress_validation(
        Path("semantic_progress_validation_result.json"),
        timeout_seconds=0.5,
    )
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if report["passed"] else 1


def run_single_eta_validation_command() -> int:
    from local_full_text_search.core.eta_replay import (
        EtaHistoryContext,
        EtaReplayEvent,
        replay_eta,
    )

    result_path = Path("single_eta_validation_result.json")
    try:
        events = [
            EtaReplayEvent(
                at_seconds=0,
                event_type="progress",
                remaining_cost_by_lane={"ocr": 80},
                workers_by_lane={"ocr": 1},
            ),
            *[
                EtaReplayEvent(
                    at_seconds=float(index * 10),
                    event_type="completion",
                    lane="ocr",
                    completed_cost=10,
                    service_seconds=10,
                    remaining_cost_by_lane={"ocr": 80 - index * 10},
                )
                for index in range(1, 5)
            ],
            EtaReplayEvent(
                at_seconds=45,
                event_type="pause",
                remaining_cost_by_lane={"ocr": 40},
            ),
            EtaReplayEvent(
                at_seconds=55,
                event_type="progress",
                remaining_cost_by_lane={"ocr": 40},
            ),
            EtaReplayEvent(
                at_seconds=60,
                event_type="resume",
                remaining_cost_by_lane={"ocr": 40},
            ),
            *[
                EtaReplayEvent(
                    at_seconds=float(60 + index * 10),
                    event_type="completion",
                    lane="ocr",
                    completed_cost=10,
                    service_seconds=10,
                    remaining_cost_by_lane={
                        "ocr": max(0, 40 - index * 10)
                    },
                )
                for index in range(1, 4)
            ],
            EtaReplayEvent(
                at_seconds=100,
                event_type="finish",
                remaining_cost_by_lane={},
            ),
        ]
        report = replay_eta(events)
        baseline = EtaHistoryContext(
            parser_name="pdf",
            parser_version="7",
            ocr_enabled=True,
            ocr_strategy="adaptive",
            ocr_model_fingerprint="validation-model",
            execution_mode="normal",
            hardware_tier="validation-cpu",
            disk_class="local-ssd",
            extension=".pdf",
            size_bucket="10m-100m",
            page_bucket="21-100",
        )
        network = EtaHistoryContext(
            **{
                **asdict(baseline),
                "disk_class": "network",
            }
        )
        if report.first_ready_seconds != 30:
            raise RuntimeError(
                f"unexpected first ready time: {report.first_ready_seconds}"
            )
        if not report.pause_frozen:
            raise RuntimeError("ETA changed during the pause interval")
        if baseline.key == network.key:
            raise RuntimeError("ETA history did not isolate disk classes")
        payload = {
            "passed": True,
            "history_isolated": True,
            **report.to_dict(),
        }
        result_path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 0
    except Exception as exc:
        payload = {"passed": False, "error": str(exc)}
        result_path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 1


def run_schema_v6_validation() -> int:
    """Migrate a populated schema-v4 database through the latest schema."""

    result_path = Path("schema_v6_validation_result.txt")
    try:
        with tempfile.TemporaryDirectory(prefix="lfts_schema_v6_") as tmp:
            base = Path(tmp)
            root = base / "files"
            root.mkdir()
            source = root / "existing.txt"
            source.write_text("SCHEMA_V6_MIGRATION_PRESERVES_HIT", encoding="utf-8")
            db_path = base / "validation.db"
            db = DatabaseManager(db_path)
            db.initialize()
            root_id = db.add_root(root)
            IndexManager(db, AppSettings(enable_ocr=False)).index_root(root_id)

            con = sqlite3.connect(db_path)
            try:
                con.execute("PRAGMA foreign_keys = OFF")
                con.execute("PRAGMA legacy_alter_table = ON")
                for index_name in (
                    "idx_files_container",
                    "idx_files_source_kind",
                    "idx_files_zip_member_source",
                    "idx_files_content_hash_full",
                    "idx_files_exact_content",
                ):
                    con.execute(f"DROP INDEX IF EXISTS {index_name}")
                con.execute("ALTER TABLE files RENAME TO files_v5")
                con.execute(
                    """
                    CREATE TABLE files (
                        id INTEGER PRIMARY KEY AUTOINCREMENT,
                        root_id INTEGER NOT NULL,
                        path TEXT NOT NULL UNIQUE,
                        filename TEXT NOT NULL,
                        extension TEXT,
                        size_bytes INTEGER,
                        modified_time REAL,
                        created_time REAL,
                        quick_fingerprint TEXT,
                        content_hash TEXT,
                        content_key TEXT,
                        document_id INTEGER,
                        parse_status TEXT NOT NULL,
                        parse_error_code TEXT,
                        parse_error_message TEXT,
                        parser_name TEXT,
                        parser_version TEXT,
                        indexed_at TEXT,
                        last_seen_at TEXT,
                        is_deleted INTEGER NOT NULL DEFAULT 0,
                        FOREIGN KEY(root_id) REFERENCES roots(id),
                        FOREIGN KEY(document_id) REFERENCES documents(id)
                    )
                    """
                )
                legacy_columns = (
                    "id, root_id, path, filename, extension, size_bytes, "
                    "modified_time, created_time, quick_fingerprint, content_hash, "
                    "content_key, document_id, parse_status, parse_error_code, "
                    "parse_error_message, parser_name, parser_version, indexed_at, "
                    "last_seen_at, is_deleted"
                )
                con.execute(
                    f"INSERT INTO files({legacy_columns}) "
                    f"SELECT {legacy_columns} FROM files_v5"
                )
                con.execute("DROP TABLE files_v5")
                con.execute("PRAGMA user_version = 4")
                con.commit()
            finally:
                con.close()

            started = time.perf_counter()
            migrated = DatabaseManager(db_path)
            migrated.initialize()
            migration_ms = int((time.perf_counter() - started) * 1000)
            backup = base / "validation.schema-v4.backup.db"
            with migrated.connect() as con:
                columns = {
                    str(row["name"]) for row in con.execute("PRAGMA table_info(files)")
                }
                indexes = {
                    str(row["name"]) for row in con.execute("PRAGMA index_list(files)")
                }
                version = int(con.execute("PRAGMA user_version").fetchone()[0])
            hits = SearchEngine(migrated).search(
                SearchQuery(text="SCHEMA_V6_MIGRATION_PRESERVES_HIT", mode="exact")
            ).total_confirmed
            required_columns = {"source_kind", "container_file_id", "content_hash_full"}
            required_indexes = {
                "idx_files_container",
                "idx_files_source_kind",
                "idx_files_zip_member_source",
                "idx_files_content_hash_full",
            }
            task_columns = set()
            attempt_table = False
            with migrated.connect() as con:
                task_columns = {
                    str(row["name"])
                    for row in con.execute("PRAGMA table_info(parse_tasks)")
                }
                attempt_table = bool(
                    con.execute(
                        """
                        SELECT 1 FROM sqlite_master
                        WHERE type = 'table' AND name = 'parse_task_attempts'
                        """
                    ).fetchone()
                )
            required_task_columns = {
                "parent_task_id",
                "unit_key",
                "progress_phase",
                "progress_cursor",
                "last_semantic_progress_at",
                "worker_pid",
                "checkpoint_path",
            }
            if version != SCHEMA_VERSION or not required_columns.issubset(columns):
                raise RuntimeError(
                    f"Schema-v6 columns were not migrated: version={version}; columns={columns}"
                )
            if not required_task_columns.issubset(task_columns) or not attempt_table:
                raise RuntimeError(
                    "Schema-v6 task-progress tables are incomplete: "
                    f"columns={task_columns}; attempts={attempt_table}"
                )
            if not required_indexes.issubset(indexes):
                raise RuntimeError(f"Schema-v5 indexes are incomplete: {indexes}")
            if not backup.is_file() or hits != 1:
                raise RuntimeError(
                    f"Schema-v4 content was not preserved: backup={backup.is_file()}; hits={hits}"
                )
            message = (
                "SCHEMA_V6_VALIDATION_OK\n"
                f"migration_ms={migration_ms}; backup={backup.name}; hits={hits}\n"
            )
            result_path.write_text(message, encoding="utf-8")
            print(message, end="")
            return 0
    except Exception as exc:
        result_path.write_text(
            f"SCHEMA_V6_VALIDATION_FAILED: {exc}\n",
            encoding="utf-8",
        )
        print(f"SCHEMA_V6_VALIDATION_FAILED: {exc}")
        return 1


def run_schema_v8_validation_command() -> int:
    from local_full_text_search.core.schema_validation import (
        run_schema_v8_validation,
    )

    result_path = Path("schema_v8_validation_result.json")
    try:
        with tempfile.TemporaryDirectory(
            prefix="lfts_schema_v8_"
        ) as tmp:
            payload = run_schema_v8_validation(Path(tmp))
        result_path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 0 if payload["passed"] else 1
    except Exception as exc:
        payload = {"passed": False, "error": str(exc)}
        result_path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 1


def run_schema_v7_validation_command() -> int:
    return run_schema_v8_validation_command()


def run_manual_exclusion_validation_command() -> int:
    from local_full_text_search.core.manual_exclusion_validation import (
        run_manual_exclusion_validation,
    )

    result_path = Path("manual_exclusion_validation_result.json")
    try:
        with tempfile.TemporaryDirectory(
            prefix="lfts_manual_exclusion_"
        ) as tmp:
            payload = run_manual_exclusion_validation(Path(tmp))
        result_path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 0 if payload["passed"] else 1
    except Exception as exc:
        payload = {"passed": False, "error": str(exc)}
        result_path.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 1


def run_phase2_validation_command(name: str) -> int:
    from local_full_text_search.core.phase2_validation import (
        validate_ocr_adaptive_v2,
        validate_ocr_backend_gate,
        validate_paused_mode_switch,
        validate_pdf_page_pipeline,
        validate_safe_pause,
        write_validation_result,
    )

    validators = {
        "pdf_page_pipeline": validate_pdf_page_pipeline,
        "ocr_adaptive_v2": validate_ocr_adaptive_v2,
        "ocr_backend": (
            lambda base: validate_ocr_backend_gate(base)
        ),
        "safe_pause": validate_safe_pause,
        "paused_mode_switch": validate_paused_mode_switch,
    }
    validator = validators.get(name)
    if validator is None:
        print(f"UNKNOWN_PHASE2_VALIDATION: {name}")
        return 2
    return write_validation_result(name, validator)


def run_legacy_word_doc_generator_command() -> int:
    index = sys.argv.index("--generate-legacy-word-doc")
    if index + 1 >= len(sys.argv):
        print("--generate-legacy-word-doc requires an output path")
        return 2
    from local_full_text_search.core.phase2_validation import (
        _create_legacy_word_doc,
    )

    try:
        _create_legacy_word_doc(Path(sys.argv[index + 1]))
    except Exception as exc:
        print(f"LEGACY_WORD_DOC_GENERATION_FAILED: {exc}")
        return 1
    print("LEGACY_WORD_DOC_GENERATION_OK")
    return 0


def _safe_command_print(payload: object) -> None:
    try:
        print(payload)
    except OSError:
        pass


def run_compare_index_command(
    baseline_path: Path,
    candidate_path: Path,
    queries_path: Path,
    output_path: Path | None = None,
) -> int:
    from local_full_text_search.core.corpus_validation import (
        compare_index_databases,
    )

    target = (
        Path(output_path).resolve()
        if output_path is not None
        else Path("compare_index_result.json").resolve()
    )
    try:
        payload = compare_index_databases(
            baseline_path,
            candidate_path,
            queries_path,
        )
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        _safe_command_print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 0 if payload["passed"] else 1
    except Exception as exc:
        payload = {"passed": False, "error": str(exc)}
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        _safe_command_print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 1


def run_cold_index_benchmark_parent(
    root: Path,
    output_path: Path,
    *,
    performance_mode: bool = False,
) -> int:
    output_path = Path(output_path).resolve()
    state_dir = output_path.parent / (
        output_path.stem + ".cold-state"
    )
    if output_path.exists():
        print(f"冷索引输出已存在，拒绝覆盖: {output_path}")
        return 2
    if state_dir.exists():
        print(f"冷索引状态目录已存在，拒绝复用缓存: {state_dir}")
        return 2
    output_path.parent.mkdir(parents=True, exist_ok=True)
    executable = Path(sys.executable).resolve()
    command = (
        [str(executable)]
        if getattr(sys, "frozen", False)
        else [str(executable), str(Path(__file__).resolve())]
    )
    environment = os.environ.copy()
    environment["LFTS_APP_DATA_DIR"] = str(state_dir)
    worker_command = [
            *command,
            "--benchmark-cold-index-worker",
            str(Path(root).resolve()),
            str(output_path),
            str(state_dir),
        ]
    if performance_mode:
        worker_command.append("--performance")
    result = subprocess.run(
        worker_command,
        cwd=output_path.parent,
        env=environment,
        check=False,
    )
    return int(result.returncode)


def run_cold_index_benchmark_worker(
    root: Path,
    output_path: Path,
    state_dir: Path,
    *,
    performance_mode: bool = False,
) -> int:
    from local_full_text_search.core.corpus_validation import (
        run_cold_index_benchmark,
    )

    try:
        payload = run_cold_index_benchmark(
            root,
            output_path,
            state_dir=state_dir,
            performance_mode=performance_mode,
        )
        _safe_command_print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 0 if payload["passed"] else 1
    except Exception as exc:
        payload = {"passed": False, "error": str(exc)}
        target = Path(output_path).resolve()
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(
            json.dumps(payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        _safe_command_print(json.dumps(payload, ensure_ascii=False, indent=2))
        return 1


def _option_after(flag: str) -> str | None:
    if flag not in sys.argv:
        return None
    index = sys.argv.index(flag)
    if index + 1 >= len(sys.argv):
        return None
    return str(sys.argv[index + 1])


def run_field_reindex_validation(db_path: Path, settings_path: Path) -> int:
    """Run a real configured update against an explicitly supplied database copy."""

    result_path = Path("field_reindex_validation_result.txt")
    try:
        settings = SettingsService(settings_path).load()
        db = DatabaseManager(db_path)
        migration_started = time.perf_counter()
        db.initialize()
        migration_ms = int((time.perf_counter() - migration_started) * 1000)
        latest_progress: dict[str, object] = {}
        last_reported = 0.0

        def capture_progress(payload: dict[str, object]) -> None:
            nonlocal latest_progress, last_reported
            latest_progress = dict(payload)
            now = time.monotonic()
            if now - last_reported >= 10.0 or payload.get("stage") in {
                "finished",
                "cancelled",
            }:
                print(
                    "FIELD_PROGRESS "
                    + json.dumps(
                        {
                            "stage": payload.get("stage"),
                            "completed": payload.get("completed_files"),
                            "total": payload.get("total_files"),
                            "failed": payload.get("failed"),
                            "queue": payload.get("queue"),
                            "current_file": payload.get("current_file"),
                        },
                        ensure_ascii=False,
                    ),
                    flush=True,
                )
                last_reported = now

        update_started = time.perf_counter()
        summary = IndexManager(db, settings).index_enabled_roots(
            progress_callback=capture_progress
        )
        update_ms = int((time.perf_counter() - update_started) * 1000)
        report = db.integrity_report()
        with db.connect() as con:
            version = int(con.execute("PRAGMA user_version").fetchone()[0])
            short_tokens = int(con.execute("SELECT COUNT(*) FROM short_tokens").fetchone()[0])
            blocks = int(con.execute("SELECT COUNT(*) FROM content_blocks").fetchone()[0])
            fts_rows = int(con.execute("SELECT COUNT(*) FROM content_fts").fetchone()[0])
            latest_run = con.execute(
                "SELECT id, status, discovered_files, write_ms, fts_ms, total_ms "
                "FROM index_runs ORDER BY started_at DESC LIMIT 1"
            ).fetchone()
            incomplete_states = {
                str(row["key"]): str(row["value"])
                for row in con.execute(
                    "SELECT key, value FROM index_state WHERE key IN ("
                    "'content_fts_dirty', 'full_batch_incomplete', 'schema_v3_rebuild_required')"
                )
            }
            task_counts = {
                str(row["status"]): int(row["n"])
                for row in con.execute(
                    "SELECT status, COUNT(*) AS n FROM parse_tasks "
                    "WHERE run_id = ? GROUP BY status",
                    (str(latest_run["id"]),),
                )
            } if latest_run is not None else {}
        if summary.cancelled or summary.failed:
            raise RuntimeError(f"Real update reported failures: {summary}")
        if version != 4 or short_tokens != 0:
            raise RuntimeError(f"Migration incomplete: version={version}; short_tokens={short_tokens}")
        if latest_run is None or str(latest_run["status"]) != "complete":
            raise RuntimeError(f"Latest run did not complete: {dict(latest_run or {})}")
        if fts_rows != blocks:
            raise RuntimeError(f"FTS mismatch: fts={fts_rows}; blocks={blocks}")
        if any(value != "0" for value in incomplete_states.values()):
            raise RuntimeError(f"Incomplete index state: {incomplete_states}")
        if report["integrity"] != ["ok"] or report["foreign_key_errors"]:
            raise RuntimeError(f"SQLite integrity failed: {report}")
        message = (
            "FIELD_REINDEX_VALIDATION_OK\n"
            f"migration_ms={migration_ms}; update_ms={update_ms}; scanned={summary.scanned}; "
            f"indexed={summary.indexed}; skipped={summary.skipped}; failed={summary.failed}; "
            f"partial_success={summary.partial_success}; blocks={blocks}; fts_rows={fts_rows}; "
            f"write_ms={int(latest_run['write_ms'])}; fts_ms={int(latest_run['fts_ms'])}; "
            f"tasks={json.dumps(task_counts, ensure_ascii=False, sort_keys=True)}; "
            f"last_stage={latest_progress.get('stage', '')}\n"
        )
        result_path.write_text(message, encoding="utf-8")
        print(message, end="")
        return 0
    except Exception as exc:
        result_path.write_text(f"FIELD_REINDEX_VALIDATION_FAILED: {exc}\n", encoding="utf-8")
        print(f"FIELD_REINDEX_VALIDATION_FAILED: {exc}")
        return 1


def run_database_lock_validation() -> int:
    """Verify busy-writer tolerance and migrated FTS rebuild performance."""

    result_path = Path("database_lock_validation_result.txt")
    try:
        with tempfile.TemporaryDirectory(prefix="lfts_database_lock_") as tmp:
            base = Path(tmp)
            root = base / "files"
            root.mkdir()
            file_count = 64
            for index in range(file_count):
                (root / f"document-{index:03d}.txt").write_text(
                    f"DATABASE_LOCK_VALIDATION_HIT_{index:03d}",
                    encoding="utf-8",
                )

            db = DatabaseManager(base / "validation.db")
            db.initialize()
            root_id = db.add_root(root)
            settings = AppSettings(
                enable_ocr=False,
                parser_workers=4,
                normal_pending_tasks=4,
                index_write_batch_size=4,
                db_write_batch_blocks=4,
                db_write_max_delay_ms=10,
                defer_fts_during_full_scan=True,
                enable_parse_cache=False,
            )
            first = IndexManager(db, settings).index_root(root_id)
            if first.indexed != file_count:
                raise RuntimeError(f"Unable to prepare baseline index: {first}")

            with db.connect() as con:
                probe_file_id = int(
                    con.execute("SELECT id FROM files ORDER BY id LIMIT 1").fetchone()[0]
                )
            probe_task_id = db.create_parse_task(
                probe_file_id,
                "database-lock-probe",
                "normal",
            )
            locker = sqlite3.connect(db.db_path, timeout=0.1)
            try:
                locker.execute("BEGIN IMMEDIATE")
                lock_skip_started = time.perf_counter()
                running_updated = db.try_mark_tasks_running(
                    [probe_task_id],
                    timeout_seconds=0.0,
                )
                spool_updated = db.try_mark_task_spooled(
                    probe_task_id,
                    base / "lock-probe.spool",
                    "lock-probe-checksum",
                    timeout_seconds=0.0,
                )
                lock_skip_elapsed = time.perf_counter() - lock_skip_started
            finally:
                locker.rollback()
                locker.close()
            if running_updated or spool_updated or lock_skip_elapsed >= 0.2:
                raise RuntimeError(
                    "Busy diagnostic writes did not skip promptly: "
                    f"running={running_updated}; spool={spool_updated}; "
                    f"elapsed={lock_skip_elapsed:.3f}s"
                )
            db.mark_task_failed(probe_task_id, "LOCK_PROBE_COMPLETE", "validation probe")

            stale_rows = 40_000
            with db.connect() as con:
                con.executemany(
                    """
                    INSERT INTO content_fts(
                        rowid, block_id, file_id, filename, path, location_text, normalized_text
                    ) VALUES (?, ?, 0, 'stale', 'stale', '', 'MIGRATED_STALE_FTS')
                    """,
                    ((1_000_000 + index, 1_000_000 + index) for index in range(stale_rows)),
                )
                con.execute(
                    "INSERT INTO index_state(key, value) VALUES ('content_fts_dirty', '1') "
                    "ON CONFLICT(key) DO UPDATE SET value='1'"
                )
                con.execute(
                    "INSERT INTO index_state(key, value) VALUES ('full_batch_incomplete', '1') "
                    "ON CONFLICT(key) DO UPDATE SET value='1'"
                )
                con.execute(
                    "INSERT INTO index_state(key, value) VALUES ('schema_v3_rebuild_required', '1') "
                    "ON CONFLICT(key) DO UPDATE SET value='1'"
                )
                con.execute(
                    "UPDATE files SET parser_version = NULL, parse_status = 'pending' "
                    "WHERE is_deleted = 0"
                )

            started = time.perf_counter()
            summary = IndexManager(db, settings).index_root(root_id)
            elapsed = time.perf_counter() - started
            report = db.integrity_report()
            with db.connect() as con:
                fts_rows = int(con.execute("SELECT COUNT(*) FROM content_fts").fetchone()[0])
                block_rows = int(con.execute("SELECT COUNT(*) FROM content_blocks").fetchone()[0])
                latest_run = con.execute(
                    "SELECT status, write_ms FROM index_runs ORDER BY started_at DESC LIMIT 1"
                ).fetchone()
                states = {
                    str(row["key"]): str(row["value"])
                    for row in con.execute(
                        "SELECT key, value FROM index_state WHERE key IN ("
                        "'content_fts_dirty', 'full_batch_incomplete', 'schema_v3_rebuild_required')"
                    )
                }
            hits = SearchEngine(db).search(
                SearchQuery(text="DATABASE_LOCK_VALIDATION_HIT_063", mode="exact")
            ).total_confirmed
            if summary.indexed != file_count or summary.failed != 0:
                raise RuntimeError(f"Unexpected rebuild summary: {summary}")
            if elapsed >= 20:
                raise RuntimeError(f"Migrated FTS rebuild remained too slow: {elapsed:.3f}s")
            if latest_run is None or latest_run["status"] != "complete":
                raise RuntimeError(f"Index run did not complete: {dict(latest_run or {})}")
            if fts_rows != block_rows or hits != 1:
                raise RuntimeError(
                    f"FTS rebuild mismatch: fts={fts_rows}; blocks={block_rows}; hits={hits}"
                )
            state_keys = (
                "content_fts_dirty",
                "full_batch_incomplete",
                "schema_v3_rebuild_required",
            )
            if any(states.get(key) != "0" for key in state_keys):
                raise RuntimeError(f"Rebuild flags were not cleared: {states}")
            if report["integrity"] != ["ok"] or report["foreign_key_errors"]:
                raise RuntimeError(f"SQLite integrity failed: {report}")
            message = (
                "DATABASE_LOCK_VALIDATION_OK\n"
                f"files={file_count}; stale_fts_rows={stale_rows}; elapsed_seconds={elapsed:.3f}; "
                f"lock_skip_ms={int(lock_skip_elapsed * 1000)}; "
                f"write_ms={int(latest_run['write_ms'])}; fts_rows={fts_rows}\n"
            )
            result_path.write_text(message, encoding="utf-8")
            print(message, end="")
            return 0
    except Exception as exc:
        result_path.write_text(f"DATABASE_LOCK_VALIDATION_FAILED: {exc}\n", encoding="utf-8")
        print(f"DATABASE_LOCK_VALIDATION_FAILED: {exc}")
        return 1


def run_shutdown_validation() -> int:
    """Force-cancel an active spawned process index and enforce a bounded exit."""

    result_path = Path("shutdown_validation_result.txt")
    try:
        from docx import Document

        with tempfile.TemporaryDirectory(prefix="lfts_shutdown_") as tmp:
            base = Path(tmp)
            root = base / "files"
            root.mkdir()
            for file_index in range(8):
                document = Document()
                for paragraph in range(800):
                    document.add_paragraph(
                        f"SHUTDOWN_VALIDATION_{file_index}_{paragraph} " + "payload " * 12
                    )
                document.save(root / f"document-{file_index}.docx")
            db = DatabaseManager(base / "validation.db")
            db.initialize()
            root_id = db.add_root(root)
            settings = AppSettings(
                enable_ocr=False,
                large_office_process_min_bytes=0,
                process_parser_workers=2,
                process_pending_tasks=2,
            )
            manager = IndexManager(db, settings)
            token = CancelToken()
            planned = threading.Event()
            failure: list[BaseException] = []

            def progress(payload: dict[str, object]) -> None:
                if payload.get("stage") in {"planning", "indexing"}:
                    planned.set()

            def run_index() -> None:
                try:
                    manager.index_root(root_id, token, progress)
                except BaseException as exc:
                    failure.append(exc)

            thread = threading.Thread(target=run_index, name="shutdown-validation")
            thread.start()
            planned.wait(timeout=10)
            time.sleep(0.25)
            started = time.perf_counter()
            token.cancel(force=True)
            manager.force_terminate_processes()
            thread.join(timeout=4)
            elapsed = time.perf_counter() - started
            if thread.is_alive():
                raise RuntimeError("Index thread did not stop within 4 seconds")
            if failure:
                raise RuntimeError(f"Index cancellation raised: {failure[0]}")
            message = f"SHUTDOWN_VALIDATION_OK\nstop_seconds={elapsed:.3f}\n"
            result_path.write_text(message, encoding="utf-8")
            print(message, end="")
            return 0
    except Exception as exc:
        result_path.write_text(f"SHUTDOWN_VALIDATION_FAILED: {exc}\n", encoding="utf-8")
        print(f"SHUTDOWN_VALIDATION_FAILED: {exc}")
        return 1


def run_checkpoint_timeout_validation() -> int:
    """Prove a long OCR task survives while it keeps making semantic progress."""

    result_path = Path("checkpoint_timeout_validation_result.txt")
    try:
        import fitz
        from PIL import Image, ImageDraw, ImageFont

        with tempfile.TemporaryDirectory(prefix="lfts_checkpoint_timeout_") as tmp:
            base = Path(tmp)
            root = base / "files"
            root.mkdir()
            image_path = base / "scanned-page.png"
            image = Image.new("RGB", (1600, 1200), "white")
            draw = ImageDraw.Draw(image)
            font_path = Path("C:/Windows/Fonts/arial.ttf")
            font = (
                ImageFont.truetype(str(font_path), 72)
                if font_path.exists()
                else ImageFont.load_default()
            )
            validation_nonce = time.time_ns()
            for line in range(12):
                draw.text(
                    (60, 50 + line * 90),
                    f"OCR TIMEOUT {validation_nonce} PAGE {line:02d}",
                    fill="black",
                    font=font,
                )
            image.save(image_path)

            pdf_path = root / "checkpoint.pdf"
            document = fitz.open()
            native_page = document.new_page()
            native_page.insert_text((72, 72), "NATIVE_CHECKPOINT_SURVIVES_TIMEOUT")
            for _ in range(3):
                scanned_page = document.new_page(width=800, height=600)
                scanned_page.insert_image(scanned_page.rect, filename=str(image_path))
            document.save(pdf_path)
            document.close()

            db = DatabaseManager(base / "validation.db")
            db.initialize()
            root_id = db.add_root(root)
            settings = AppSettings(
                enable_ocr=True,
                ocr_images=True,
                ocr_scanned_pdf=True,
                ocr_cpu_threads=1,
                ocr_no_progress_timeout_seconds=15,
                process_no_progress_timeout_seconds=15,
                no_progress_max_retries=0,
                process_parser_workers=1,
                process_pending_tasks=1,
                pdf_parser_workers=1,
                pdf_pending_tasks=1,
                ocr_workers=1,
                ocr_pending_tasks=1,
                ocr_microbatch_parent_jobs=1,
            )
            started = time.perf_counter()
            summary = IndexManager(db, settings).index_root(root_id)
            elapsed = time.perf_counter() - started
            try:
                hit_count = SearchEngine(db).search(
                    SearchQuery(text="NATIVE_CHECKPOINT_SURVIVES_TIMEOUT", mode="exact")
                ).total_confirmed
            except Exception as exc:
                raise RuntimeError(
                    "Completed semantic-progress index was not searchable: "
                    f"{type(exc).__name__}: {exc}"
                ) from exc
            with db.connect() as con:
                row = con.execute(
                    "SELECT parse_status, parse_error_code FROM files WHERE path = ?",
                    (str(pdf_path),),
                ).fetchone()
                block_count = int(con.execute("SELECT COUNT(*) FROM content_blocks").fetchone()[0])
            if row is None:
                raise RuntimeError("Checkpoint PDF was not recorded")
            if row["parse_status"] != "success" or row["parse_error_code"] is not None:
                raise RuntimeError(f"Semantic progress was incorrectly timed out: {dict(row)}")
            if summary.failed != 0 or hit_count != 1 or block_count < 4:
                raise RuntimeError(
                    f"Semantic progress result was incomplete: summary={summary}; "
                    f"hits={hit_count}; blocks={block_count}"
                )
            if elapsed <= settings.process_no_progress_timeout_seconds:
                raise RuntimeError(
                    "Validation did not run longer than the no-progress threshold: "
                    f"elapsed={elapsed:.3f}s"
                )
            message = (
                "SEMANTIC_PROGRESS_VALIDATION_OK\n"
                f"elapsed_seconds={elapsed:.3f}; failed={summary.failed}; "
                f"hits={hit_count}; published_blocks={block_count}\n"
            )
            result_path.write_text(message, encoding="utf-8")
            print(message, end="")
            return 0
    except Exception as exc:
        result_path.write_text(
            f"CHECKPOINT_TIMEOUT_VALIDATION_FAILED: {exc}\n",
            encoding="utf-8",
        )
        print(f"CHECKPOINT_TIMEOUT_VALIDATION_FAILED: {exc}")
        return 1


def run_legacy_shutdown_validation(source: Path) -> int:
    """Start a real Office conversion, force-cancel it, and verify PID cleanup."""

    result_path = Path("legacy_shutdown_validation_result.txt")
    try:
        source = source.resolve()
        if not source.is_file() or source.suffix.lower() not in {".doc", ".xls", ".ppt"}:
            raise ValueError("需要现有的 .doc/.xls/.ppt 文件")
        runtime_parent = Path.cwd() if Path.cwd().anchor == source.anchor else None
        with tempfile.TemporaryDirectory(
            prefix="lfts_legacy_shutdown_",
            dir=runtime_parent,
        ) as tmp:
            base = Path(tmp)
            root = base / "files"
            root.mkdir()
            linked = root / source.name
            try:
                os.link(source, linked)
            except OSError:
                shutil.copy2(source, linked)
            db = DatabaseManager(base / "validation.db")
            db.initialize()
            root_id = db.add_root(root)
            settings = AppSettings(
                enable_ocr=False,
                process_parser_workers=1,
                process_pending_tasks=1,
                legacy_no_progress_timeout_seconds=600,
                legacy_conversion_cache=False,
            )
            manager = IndexManager(db, settings)
            token = CancelToken()
            failure: list[BaseException] = []

            def run_index() -> None:
                try:
                    manager.index_root(root_id, token)
                except BaseException as exc:
                    failure.append(exc)

            thread = threading.Thread(target=run_index, name="legacy-shutdown-validation")
            started = time.perf_counter()
            started_wall = time.time()
            thread.start()
            record: dict[str, object] | None = None
            deadline = time.monotonic() + 45
            while time.monotonic() < deadline:
                records = [
                    path
                    for path in (TEMP_DIR / "process_results").glob("*/office_processes/*.json")
                    if path.stat().st_mtime >= started_wall
                ]
                if records:
                    record = json.loads(records[-1].read_text(encoding="ascii"))
                    break
                if not thread.is_alive():
                    break
                time.sleep(0.1)
            if record is None:
                raise RuntimeError("未观察到应用登记的 Office 自动化进程")
            office_pid = int(record["pid"])
            cancel_started = time.perf_counter()
            token.cancel(force=True)
            manager.force_terminate_processes()
            thread.join(timeout=4)
            stop_seconds = time.perf_counter() - cancel_started
            if thread.is_alive():
                raise RuntimeError("老版 Office 索引线程未在 4 秒内停止")
            if failure:
                raise RuntimeError(f"取消过程抛出异常: {failure[0]}")
            try:
                import psutil

                if psutil.pid_exists(office_pid):
                    process = psutil.Process(office_pid)
                    expected = float(record.get("create_time") or 0.0)
                    if not expected or abs(process.create_time() - expected) <= 1.0:
                        raise RuntimeError(f"Office 自动化进程仍然存在: PID {office_pid}")
            except ImportError:
                pass
            message = (
                "LEGACY_SHUTDOWN_VALIDATION_OK\n"
                f"office_pid={office_pid}; stop_seconds={stop_seconds:.3f}; "
                f"elapsed_seconds={time.perf_counter() - started:.3f}\n"
            )
            result_path.write_text(message, encoding="utf-8")
            print(message, end="")
            return 0
    except Exception as exc:
        result_path.write_text(
            f"LEGACY_SHUTDOWN_VALIDATION_FAILED: {exc}\n",
            encoding="utf-8",
        )
        print(f"LEGACY_SHUTDOWN_VALIDATION_FAILED: {exc}")
        return 1


def _create_process_validation_files(root: Path) -> None:
    from docx import Document
    from openpyxl import Workbook

    document = Document()
    for index in range(100):
        document.add_paragraph(f"PROCESS_DOCX_VALIDATION_HIT paragraph {index}")
    document.save(root / "process_sample.docx")

    workbook = Workbook(write_only=True)
    sheet = workbook.create_sheet("ProcessValidation")
    for index in range(1000):
        sheet.append([index, "PROCESS_XLSX_VALIDATION_HIT", f"row-{index}"])
    workbook.save(root / "process_sample.xlsx")


def _create_validation_files(root: Path) -> None:
    (root / "sample.txt").write_text("TXT_VALIDATION_HIT\nBS-2800M2", encoding="utf-8")

    from docx import Document
    doc = Document()
    doc.add_paragraph("DOCX_VALIDATION_HIT")
    doc.save(root / "sample.docx")

    from openpyxl import Workbook
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    sheet["A1"] = "XLSX_VALIDATION_HIT"
    workbook.save(root / "sample.xlsx")

    from pptx import Presentation
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(914400, 914400, 3657600, 914400)
    text_box.text = "PPTX_VALIDATION_HIT"
    presentation.save(root / "sample.pptx")

    import fitz
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "PDF_VALIDATION_HIT")
    pdf.save(root / "sample.pdf")
    pdf.close()

    from PIL import Image, ImageDraw, ImageFont
    image = Image.new("RGB", (900, 260), "white")
    draw = ImageDraw.Draw(image)
    font_path = Path("C:/Windows/Fonts/arial.ttf")
    font = ImageFont.truetype(str(font_path), 86) if font_path.exists() else ImageFont.load_default()
    draw.text((48, 80), "OCR TEST 123", fill="black", font=font)
    image.save(root / "sample_ocr.png")


def apply_light_theme(app: object) -> None:
    try:
        qss = resources.files("local_full_text_search.ui.styles").joinpath("light.qss").read_text(encoding="utf-8")
        app.setStyleSheet(qss)
    except Exception:
        # A missing stylesheet should never prevent the search tool from opening.
        pass


def main() -> int:
    diagnostics = StartupDiagnostics()
    diagnostics.begin()

    def finish_command(result: int) -> int:
        diagnostics.mark_completed()
        return result

    try:
        diagnostics.stage_started("加载应用模块")
        load_application_dependencies()
        diagnostics.stage_completed()
        diagnostics.stage_started("配置日志")
        configure_logging()
        diagnostics.stage_completed()
        if "--self-test" in sys.argv:
            return finish_command(run_self_test())
        if "--validate-core" in sys.argv:
            return finish_command(run_core_validation())
        if "--validate-process-pool" in sys.argv:
            return finish_command(run_process_pool_validation())
        if "--validate-schema-v3" in sys.argv:
            return finish_command(run_schema_v3_validation())
        if "--validate-schema-v4" in sys.argv:
            return finish_command(run_schema_v4_validation())
        if "--validate-schema-v6" in sys.argv:
            return finish_command(run_schema_v6_validation())
        if "--validate-schema-v7" in sys.argv:
            return finish_command(run_schema_v7_validation_command())
        if "--validate-schema-v8" in sys.argv:
            return finish_command(run_schema_v8_validation_command())
        if "--validate-manual-exclusion" in sys.argv:
            return finish_command(
                run_manual_exclusion_validation_command()
            )
        if "--validate-failure-demo" in sys.argv:
            return finish_command(run_failure_fallback_demo_validation())
        if "--validate-pdf-page-pipeline" in sys.argv:
            return finish_command(
                run_phase2_validation_command(
                    "pdf_page_pipeline"
                )
            )
        if "--validate-ocr-adaptive-v2" in sys.argv:
            return finish_command(
                run_phase2_validation_command("ocr_adaptive_v2")
            )
        if "--validate-ocr-backend" in sys.argv:
            return finish_command(
                run_phase2_validation_command("ocr_backend")
            )
        if "--capture-index-status-layout" in sys.argv:
            index = sys.argv.index(
                "--capture-index-status-layout"
            )
            if index + 5 >= len(sys.argv):
                print(
                    "--capture-index-status-layout "
                    "需要 width height scale state output"
                )
                return finish_command(2)
            return finish_command(
                capture_index_status_layout(
                    int(sys.argv[index + 1]),
                    int(sys.argv[index + 2]),
                    float(sys.argv[index + 3]),
                    str(sys.argv[index + 4]),
                    Path(sys.argv[index + 5]),
                )
            )
        if "--validate-index-status-layout" in sys.argv:
            return finish_command(
                run_index_status_layout_validation()
            )
        if "--validate-safe-pause" in sys.argv:
            return finish_command(
                run_phase2_validation_command("safe_pause")
            )
        if "--validate-paused-mode-switch" in sys.argv:
            return finish_command(
                run_phase2_validation_command(
                    "paused_mode_switch"
                )
            )
        if "--generate-legacy-word-doc" in sys.argv:
            return finish_command(run_legacy_word_doc_generator_command())
        if "--benchmark-cold-index-worker" in sys.argv:
            index = sys.argv.index("--benchmark-cold-index-worker")
            if index + 3 >= len(sys.argv):
                print(
                    "--benchmark-cold-index-worker "
                    "需要 folder output state-dir"
                )
                return finish_command(2)
            return finish_command(
                run_cold_index_benchmark_worker(
                    Path(sys.argv[index + 1]),
                    Path(sys.argv[index + 2]),
                    Path(sys.argv[index + 3]),
                    performance_mode="--performance" in sys.argv,
                )
            )
        if "--benchmark-cold-index" in sys.argv:
            index = sys.argv.index("--benchmark-cold-index")
            output = _option_after("--output")
            if index + 1 >= len(sys.argv) or output is None:
                print(
                    "--benchmark-cold-index "
                    "需要 <folder> --output <json>"
                )
                return finish_command(2)
            return finish_command(
                run_cold_index_benchmark_parent(
                    Path(sys.argv[index + 1]),
                    Path(output),
                    performance_mode="--performance" in sys.argv,
                )
            )
        if "--compare-index" in sys.argv:
            index = sys.argv.index("--compare-index")
            queries = _option_after("--queries")
            output = _option_after("--output")
            if index + 2 >= len(sys.argv) or queries is None:
                print(
                    "--compare-index 需要 <baseline-db> "
                    "<candidate-db> --queries <json>"
                )
                return finish_command(2)
            return finish_command(
                run_compare_index_command(
                    Path(sys.argv[index + 1]),
                    Path(sys.argv[index + 2]),
                    Path(queries),
                    Path(output) if output else None,
                )
            )
        if "--validate-field-reindex" in sys.argv:
            index = sys.argv.index("--validate-field-reindex")
            if index + 2 >= len(sys.argv):
                print("--validate-field-reindex 需要数据库副本和设置文件路径")
                return finish_command(2)
            return finish_command(
                run_field_reindex_validation(
                    Path(sys.argv[index + 1]),
                    Path(sys.argv[index + 2]),
                )
            )
        if "--validate-database-lock" in sys.argv:
            return finish_command(run_database_lock_validation())
        if "--validate-shutdown" in sys.argv:
            return finish_command(run_shutdown_validation())
        if "--validate-checkpoint-timeout" in sys.argv:
            return finish_command(run_checkpoint_timeout_validation())
        if "--validate-hang-recovery" in sys.argv:
            return finish_command(run_hang_recovery_validation_command())
        if "--validate-semantic-progress" in sys.argv:
            return finish_command(run_semantic_progress_validation_command())
        if "--validate-single-eta" in sys.argv:
            return finish_command(run_single_eta_validation_command())
        if "--validate-legacy-shutdown" in sys.argv:
            index = sys.argv.index("--validate-legacy-shutdown")
            if index + 1 >= len(sys.argv):
                print("--validate-legacy-shutdown 需要一个 .doc/.xls/.ppt 路径")
                return finish_command(2)
            return finish_command(run_legacy_shutdown_validation(Path(sys.argv[index + 1])))
        if "--validate-ui" in sys.argv:
            return finish_command(run_ui_validation())
        diagnostics.stage_started("加载图形界面")
        from PySide6.QtWidgets import QApplication, QMessageBox
        from local_full_text_search.ui.main_window import MainWindow
        diagnostics.stage_completed()
    except Exception as exc:
        diagnostics.stage_failed(diagnostics.current_stage, exc)
        diagnostics.show_native_error("本地全文搜索启动失败", diagnostics.format_error_message(exc))
        return 2
    try:
        diagnostics.stage_started("读取设置")
        settings_service = SettingsService()
        settings = settings_service.load()
        diagnostics.stage_completed()
        diagnostics.stage_started("初始化索引数据库")
        db = DatabaseManager()
        demo_mode = "--failure-fallback-demo" in sys.argv
        demo_root: Path | None = None
        if demo_mode:
            app_data = Path(os.environ.get("LFTS_APP_DATA_DIR", "")).resolve()
            if not os.environ.get("LFTS_APP_DATA_DIR"):
                raise RuntimeError("异常保底演示必须使用独立的 LFTS_APP_DATA_DIR")
            demo_root = reset_failure_fallback_demo(app_data, db.db_path)
        db.initialize()
        if demo_root is not None:
            seed_failure_fallback_demo(db, demo_root)
        diagnostics.stage_completed()
        diagnostics.stage_started("创建图形界面")
        app = QApplication(sys.argv)
        app.setApplicationName(APP_DISPLAY_NAME)
        apply_light_theme(app)
        window = MainWindow(db, settings, settings_service)
        if demo_mode:
            window.setWindowTitle(f"{APP_DISPLAY_NAME} - 异常保底演示（隔离数据）")
            window.switch_page("failed")
        window.show()
        diagnostics.stage_completed()
        diagnostics.mark_window_visible()
        if diagnostics.previous_failure:
            QMessageBox.warning(
                window,
                "检测到上次启动未完成",
                "上次启动没有正常显示主窗口。\n\n"
                f"上次阶段：{diagnostics.previous_failure.get('stage', '未知')}\n"
                f"诊断日志：{diagnostics.log_path}",
            )
        return app.exec()
    except Exception as exc:
        diagnostics.stage_failed(diagnostics.current_stage, exc)
        diagnostics.show_native_error("本地全文搜索启动失败", diagnostics.format_error_message(exc))
        return 1


if __name__ == "__main__":
    multiprocessing.freeze_support()
    raise SystemExit(main())
